#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR

# Load the template
template_path = '/Users/karsten/Nextcloud/Documents/200_Presentations/101_ITU/100_Cloud Fundamentals/Cloud Fundamentals.pptx'
prs = Presentation(template_path)

# Clear all existing slides (keep only layouts)
while len(prs.slides) > 0:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])

# Define font settings - NO BOLD!
TITLE_FONT = "Open Sans"
CONTENT_FONT = "Open Sans Light"
TITLE_SIZE = Pt(44)
SLIDE_TITLE_SIZE = Pt(36)
CONTENT_SIZE = Pt(18)

# Helper function to find placeholder by index
def find_placeholder(slide, idx):
    for placeholder in slide.placeholders:
        try:
            if placeholder.placeholder_format.idx == idx:
                return placeholder
        except:
            continue
    return None

def add_content_slide(prs, slide_title, content_items, notes_text=None):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Use the title placeholder
    title = slide.shapes.title
    title.text = slide_title
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = TITLE_FONT
            run.font.size = SLIDE_TITLE_SIZE
            run.font.bold = False
            run.font.color.rgb = RGBColor(68, 68, 68)
    
    # Replace subtitle with content textbox
    subtitle = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 4:  # SUBTITLE
            subtitle = shape
            break
    
    if subtitle:
        # Clear subtitle and replace with content
        tf = subtitle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Pt(12)
        tf.margin_bottom = Pt(12)
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        
        for i, item in enumerate(content_items):
            p = tf.add_paragraph()
            p.text = item
            p.font.name = CONTENT_FONT
            p.font.size = CONTENT_SIZE
            p.level = 0
            p.space_after = Pt(10)
            p.space_before = Pt(4) if i > 0 else Pt(0)
    
    # Add trainer notes if provided
    if notes_text:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes_text
    
    return slide

# Slide 1: Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = find_placeholder(slide, 1)

title.text = "Software-Architektur Training"
for paragraph in title.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = TITLE_FONT
        run.font.size = TITLE_SIZE
        run.font.bold = False
        run.font.color.rgb = RGBColor(68, 68, 68)

if subtitle:
    subtitle.text = "Telekom Architecture Training 2025\nSoftware-Design und Architektur-Prinzipien\nBest Practices & Patterns"
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = CONTENT_FONT
            run.font.size = Pt(22)
            run.font.bold = False
            run.font.color.rgb = RGBColor(100, 100, 100)

# Slide 2: Trainer Introduction
trainer_info = [
    "Software-Architekt und Senior Developer",
    "15+ Jahre Erfahrung in Enterprise-Architektur", 
    "Spezialist für Design Patterns und Clean Code",
    "Telekom-Projekte: Skalierbare Systeme & Microservices",
    "Autor verschiedener Architektur-Guidelines",
    "Zertifiziert in Software-Architecture (CPSA-F/A)"
]
notes_text = "Trainer-Notizen:\n- Persönliche Vorstellung anpassen\n- Eigene Erfahrungen mit Telekom-Projekten erwähnen\n- Interaktion fördern: 'Was ist eure Erfahrung mit Architektur-Patterns?'\n- Zeitrahmen: 3-5 Minuten"
add_content_slide(prs, "Ihr Trainer", trainer_info, notes_text)

# Slide 3: Workshop Overview
overview_items = [
    "Intensive Schulung zu Software-Architektur",
    "Fokus auf praktische Anwendung von Design Prinzipien",
    "Telekom-spezifische Beispiele und Use Cases",
    "Hands-on Sessions mit realen Problemen",
    "Clean Code und SOLID Prinzipien im Detail", 
    "Architektur-Entscheidungen fundiert treffen lernen"
]
add_content_slide(prs, "Workshop-Übersicht", overview_items)

# Slide 4: Workshop-Agenda
agenda_items = [
    "Einführung: Software-Architektur Grundlagen & Prinzipien",
    "Clean Architecture: Strukturierung für Wartbarkeit", 
    "SOLID Prinzipien: Fundament für guten Code",
    "Design Patterns: Bewährte Lösungsmuster verstehen",
    "Code-Smells: Anti-Patterns erkennen und vermeiden",
    "Praktische Anwendung: Telekom-spezifische Beispiele"
]
add_content_slide(prs, "Workshop-Agenda", agenda_items)

# Slide 5: Lernziele
objectives = [
    "Software-Architektur Grundprinzipien verstehen und anwenden",
    "Clean Architecture Konzepte für wartbare Systeme",
    "SOLID Prinzipien als Fundament für guten Code",
    "Design Patterns als bewährte Lösungsmuster",
    "Code-Smells erkennen und vermeiden lernen", 
    "Telekom-spezifische Architektur-Herausforderungen lösen"
]
notes_text = "Trainer-Notizen:\n- Betonen Sie die praktische Anwendung der Prinzipien\n- Verweisen Sie auf die Telekom-spezifischen Beispiele\n- Erwähnen Sie, dass alle Konzepte mit Code-Beispielen unterlegt sind\n- Zeitrahmen: ca. 5 Minuten"
add_content_slide(prs, "Lernziele des Workshops", objectives, notes_text)

# Slide 6: Methodik
methodology_items = [
    "Theoretische Einführung mit praxisnahen Beispielen",
    "Live-Coding Sessions für direkte Anwendung",
    "Telekom-spezifische Use Cases und Herausforderungen",
    "Hands-on Übungen in kleinen Gruppen",
    "Code Reviews und Best-Practice Diskussionen",
    "Q&A Sessions für individuelle Fragen",
    "Praktische Refactoring-Workshops"
]
notes_text = "Trainer-Notizen:\n- Interaktiven Charakter des Trainings betonen\n- Teilnehmer ermutigen, eigene Beispiele einzubringen\n- Laptop/IDE-Setup für Hands-on Sessions erwähnen\n- Zeitrahmen: ca. 3 Minuten"
add_content_slide(prs, "Methodik", methodology_items, notes_text)

# Slide 7: Software-Architektur Definition
fundamentals = [
    "Strukturierung von Software-Systemen",
    "Entscheidungen über Komponenten und deren Beziehungen",
    "Balance zwischen Flexibilität und Komplexität",
    "Qualitätsattribute: Performance, Wartbarkeit, Skalierbarkeit",
    "Architektur-Patterns als bewährte Lösungen",
    "Dokumentation und Kommunikation der Architektur"
]
add_content_slide(prs, "Software-Architektur Definition", fundamentals)

# Slide 8: Architektur vs Design
architecture_vs_design = [
    "Architektur: Strategische, strukturelle Entscheidungen",
    "Design: Taktische, detaillierte Implementierung",
    "Architektur: Schwer änderbare Grundentscheidungen",
    "Design: Flexibel anpassbare Implementierungsdetails",
    "Architektur: System-weite Auswirkungen und Constraints",
    "Design: Lokale Optimierungen und Verfeinerungen"
]
notes_text = "Trainer-Notizen:\n- Unterschied anhand Telekom-Beispielen erklären\n- Microservices (Architektur) vs. Service-Implementation (Design)\n- Kosten von Architektur-Änderungen vs. Design-Änderungen\n- Zeitrahmen: ca. 7-8 Minuten"
add_content_slide(prs, "Architektur vs Design", architecture_vs_design, notes_text)

# Slide 9: Clean Architecture
clean_principles = [
    "Unabhängigkeit von Frameworks",
    "Testbarkeit als Kernprinzip",
    "UI-Unabhängigkeit der Geschäftslogik",
    "Datenbank-Unabhängigkeit",
    "Schichtenarchitektur mit klaren Abhängigkeiten",
    "Dependency Rule: Abhängigkeiten zeigen nach innen"
]
add_content_slide(prs, "Clean Architecture Prinzipien", clean_principles)

# Slide 10: SOLID Prinzipien
solid_items = [
    "Single Responsibility - Eine Klasse, eine Verantwortung",
    "Open/Closed - Offen für Erweiterung, geschlossen für Änderung",
    "Liskov Substitution - Austauschbarkeit von Subtypen",
    "Interface Segregation - Kleine, spezifische Interfaces",
    "Dependency Inversion - Abhängigkeit von Abstraktionen"
]
add_content_slide(prs, "SOLID Prinzipien", solid_items)

# Slide 11: DRY, KISS, YAGNI (Special formatting)
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "DRY, KISS, YAGNI"
for paragraph in title.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = TITLE_FONT
        run.font.size = SLIDE_TITLE_SIZE
        run.font.bold = False
        run.font.color.rgb = RGBColor(68, 68, 68)

# Replace subtitle with special content
subtitle = None
for shape in slide.shapes:
    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 4:  # SUBTITLE
        subtitle = shape
        break

if subtitle:
    tf = subtitle.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_top = Pt(12)
    
    dry_kiss_yagni = [
        ("DRY - Don't Repeat Yourself", 0, 20),
        ("→ Wiederholungen eliminieren, aber nicht übertreiben", 1, 16),
        ("→ Abstraktion vs. Premature Abstraction abwägen", 1, 16),
        ("", 0, 16),
        ("KISS - Keep It Simple, Stupid", 0, 20),
        ("→ Einfachheit vor Cleverness bevorzugen", 1, 16),
        ("→ Verständlichkeit und Wartbarkeit priorisieren", 1, 16),
        ("", 0, 16),
        ("YAGNI - You Aren't Gonna Need It", 0, 20),
        ("→ Keine Funktionen 'auf Vorrat' entwickeln", 1, 16),
        ("→ Aktuelle Anforderungen lösen, nicht hypothetische", 1, 16)
    ]
    
    for i, (text, level, size) in enumerate(dry_kiss_yagni):
        if text == "":
            p = tf.add_paragraph()
            p.text = ""
            p.space_after = Pt(6)
        else:
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.name = TITLE_FONT if level == 0 else CONTENT_FONT
            p.font.size = Pt(size)
            p.font.bold = False
            p.space_after = Pt(8) if level == 0 else Pt(4)
            p.space_before = Pt(4) if i > 0 and text != "" else Pt(0)

# Slide 12: Code-Smells
code_smells = [
    "Long Method - Methoden mit zu vielen Zeilen",
    "Large Class - Klassen mit zu vielen Verantwortungen",
    "Duplicate Code - Identische oder ähnliche Code-Blöcke",
    "Feature Envy - Klasse nutzt andere Klasse übermäßig",
    "Data Clumps - Gruppierte Daten ohne Struktur",
    "Primitive Obsession - Übernutzung primitiver Datentypen",
    "Long Parameter List - Methoden mit zu vielen Parametern",
    "Shotgun Surgery - Änderung erfordert viele Stellen"
]
notes_text = "Trainer-Notizen:\n- Code-Smells anhand von Telekom-Beispielen erklären\n- Interaktion: 'Welche Code-Smells habt ihr schon erlebt?'\n- Refactoring-Techniken als Lösung erwähnen\n- Überleitung zu Design Patterns als Lösungsansatz\n- Zeitrahmen: ca. 10 Minuten"
add_content_slide(prs, "Code-Smells", code_smells, notes_text)

# ============================================================================
# CREATIONAL PATTERNS SECTION (25-30 slides)
# ============================================================================

def add_code_slide(prs, slide_title, code_content, language="java", notes_text=None):
    """Add a slide with formatted code using Source Code Pro font"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Title
    title = slide.shapes.title
    title.text = slide_title
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = TITLE_FONT
            run.font.size = SLIDE_TITLE_SIZE
            run.font.bold = False
            run.font.color.rgb = RGBColor(68, 68, 68)
    
    # Code content
    subtitle = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 4:  # SUBTITLE
            subtitle = shape
            break
    
    if subtitle:
        tf = subtitle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Pt(12)
        tf.margin_bottom = Pt(12)
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        
        p = tf.add_paragraph()
        p.text = code_content
        p.font.name = "Source Code Pro"  # Code font
        p.font.size = Pt(14)
        p.font.bold = False
        p.space_after = Pt(0)
    
    # Add trainer notes if provided
    if notes_text:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes_text
    
    return slide

# Slide 13: Creational Patterns Section Divider
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Design Patterns"
for paragraph in title.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = TITLE_FONT
        run.font.size = Pt(48)
        run.font.bold = False
        run.font.color.rgb = RGBColor(68, 68, 68)

subtitle = None
for shape in slide.shapes:
    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 4:  # SUBTITLE
        subtitle = shape
        break

if subtitle:
    subtitle.text = "Creational Patterns\n\nBewährte Lösungen für Objekterzeugung"
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = CONTENT_FONT
            run.font.size = Pt(28)
            run.font.bold = False
            run.font.color.rgb = RGBColor(100, 100, 100)
        paragraph.alignment = PP_ALIGN.CENTER

# Slide 14: Creational Patterns Overview
creational_overview = [
    "Factory Method - Objekterzeugung über Subklassen",
    "Abstract Factory - Familien verwandter Objekte erstellen",
    "Builder - Komplexe Objekte Schritt für Schritt konstruieren", 
    "Prototype - Objekte durch Klonen erzeugen",
    "Singleton - Eindeutige Instanz einer Klasse gewährleisten",
    "",
    "Gemeinsames Ziel: Kontrolle über Objekterzeugung",
    "Flexibilität und Entkopplung bei der Instanzierung"
]
notes_text = "Trainer-Notizen:\n- Überblick über alle 5 Creational Patterns\n- Gemeinsame Ziele betonen: Entkopplung, Flexibilität\n- Unterschiedliche Einsatzgebiete kurz erwähnen\n- Telekom-Kontext: Verschiedene Service-Typen, Kunden-Kategorien\n- Zeitrahmen: ca. 5 Minuten"
add_content_slide(prs, "Creational Patterns - Überblick", creational_overview, notes_text)

# Slide 15: Factory Method - Problem
factory_method_problem = [
    "Problematische Objekterzeugung in Telekom-Systemen:",
    "",
    "• Switch/If-Statements für Kunden-Typen",
    "• Harte Abhängigkeiten zu konkreten Klassen",
    "• Open/Closed Principle wird verletzt",
    "• Code-Duplikation bei ähnlichen Objekten",
    "• Schwierige Erweiterung um neue Typen",
    "",
    "Beispiel: CustomerManager mit Switch-Statement",
    "→ Jeder neue Kunden-Typ erfordert Code-Änderung"
]
add_content_slide(prs, "Factory Method - Problem", factory_method_problem)

# Slide 16: Factory Method Code - Problematisch
problematic_code = '''public class CustomerManager {
    public Customer createCustomer(String type, String name, String contractId) {
        switch (type) {
            case "PRIVATE":
                Customer privateCustomer = new Customer();
                privateCustomer.setName(name);
                privateCustomer.setContractId(contractId);
                privateCustomer.setTariffOptions(Arrays.asList("Basic", "Comfort"));
                privateCustomer.setPaymentMethod("SEPA");
                return privateCustomer;
                
            case "BUSINESS":
                Customer businessCustomer = new Customer();
                businessCustomer.setName(name);
                businessCustomer.setContractId(contractId);
                businessCustomer.setTariffOptions(Arrays.asList("Professional"));
                businessCustomer.setPaymentMethod("Invoice");
                return businessCustomer;
                
            default:
                throw new IllegalArgumentException("Unknown type: " + type);
        }
    }
}'''
notes_text = "Trainer-Notizen:\n- Code-Smells identifizieren lassen\n- Long Method, Switch Statement, Duplicate Code\n- Open/Closed Principle Verletzung zeigen\n- Zeitrahmen: ca. 7 Minuten"
add_code_slide(prs, "Factory Method - Problematischer Code", problematic_code, "java", notes_text)

# Slide 17: Factory Method - Lösung Struktur
factory_method_solution = [
    "Factory Method Pattern - Struktur:",
    "",
    "• Creator (abstract): Definiert Factory Method",
    "• ConcreteCreator: Implementiert spezifische Erzeugung",
    "• Product (interface): Definiert Produkt-Schnittstelle",
    "• ConcreteProduct: Spezifische Implementierung",
    "",
    "Vorteile:",
    "• Polymorphismus statt Switch-Statements",
    "• Open/Closed Principle erfüllt",
    "• Single Responsibility für jede Factory"
]
add_content_slide(prs, "Factory Method - Lösung", factory_method_solution)

# Slide 18: Factory Method Code - Customer Interface
customer_interface = '''// Product Interface
public interface Customer {
    String getName();
    String getContractId();
    List<String> getTariffOptions();
    String getPaymentMethod();
    String getSupportLevel();
    void processContract();
}

// Konkrete Implementierung
public class PrivateCustomer implements Customer {
    private String name;
    private String contractId;
    private List<String> tariffOptions;
    
    public PrivateCustomer(String name, String contractId) {
        this.name = name;
        this.contractId = contractId;
        this.tariffOptions = Arrays.asList("Basic", "Comfort");
    }
    
    @Override
    public String getPaymentMethod() {
        return "SEPA";
    }
    
    @Override
    public void processContract() {
        validatePersonalData();
        setupBasicServices();
    }
}'''
add_code_slide(prs, "Factory Method - Customer Interface", customer_interface)

# Slide 19: Factory Method Code - Factory Implementation
factory_implementation = '''// Abstract Creator
public abstract class CustomerFactory {
    // Factory Method - zu implementieren von Subklassen
    protected abstract Customer createCustomer(String name, String contractId);
    
    // Template Method - verwendet Factory Method
    public Customer processNewCustomer(String name, String contractId) {
        Customer customer = createCustomer(name, contractId);
        
        // Gemeinsame Geschäftslogik
        validateContract(customer);
        persistCustomer(customer);
        sendWelcomeMessage(customer);
        
        return customer;
    }
    
    private void validateContract(Customer customer) { /* ... */ }
    private void persistCustomer(Customer customer) { /* ... */ }
    private void sendWelcomeMessage(Customer customer) { /* ... */ }
}

// Concrete Creator
public class PrivateCustomerFactory extends CustomerFactory {
    @Override
    protected Customer createCustomer(String name, String contractId) {
        return new PrivateCustomer(name, contractId);
    }
}'''
add_code_slide(prs, "Factory Method - Factory Implementation", factory_implementation)

# Slide 20: Factory Method - Telekom Beispiel
telekom_factory_example = [
    "Telekom-spezifisches Factory Method Beispiel:",
    "",
    "ServiceProvisioningFactory für verschiedene Dienste:",
    "• DSLServiceFactory → Erstellt DSL-Services",
    "• FiberServiceFactory → Erstellt Glasfaser-Services", 
    "• MobileServiceFactory → Erstellt Mobil-Services",
    "",
    "Vorteile im Telekom-Kontext:",
    "• Neue Services ohne Code-Änderung hinzufügbar",
    "• Service-spezifische Konfigurationen gekapselt",
    "• Einheitliche Service-Aktivierung über Template Method"
]
notes_text = "Trainer-Notizen:\n- Bezug zu realen Telekom-Services herstellen\n- Template Method Pattern Integration zeigen\n- Skalierbarkeit für neue Services betonen\n- Zeitrahmen: ca. 6 Minuten"
add_content_slide(prs, "Factory Method - Telekom Beispiel", telekom_factory_example, notes_text)

# Slide 21: Abstract Factory - Problem
abstract_factory_problem = [
    "Problem: Familien von verwandten Objekten",
    "",
    "Telekom-Szenario: UI-Komponenten für verschiedene Plattformen",
    "• Desktop-Anwendung (Swing)",
    "• Web-Anwendung (HTML)",
    "• Mobile-App (Native)",
    "",
    "Jede Plattform benötigt:",
    "• Button, TextField, Dialog, Menu",
    "• Konsistentes Look & Feel pro Plattform",
    "• Austauschbarkeit zwischen Plattformen"
]
add_content_slide(prs, "Abstract Factory - Problem", abstract_factory_problem)

# Slide 22: Abstract Factory Code - Structure
abstract_factory_structure = '''// Abstract Factory
public interface UIComponentFactory {
    Button createButton();
    TextField createTextField();
    Dialog createDialog();
}

// Product Interfaces
public interface Button {
    void click();
    void setLabel(String text);
}

public interface TextField {
    void setText(String text);
    String getText();
}

// Concrete Factory
public class WebUIFactory implements UIComponentFactory {
    @Override
    public Button createButton() {
        return new HTMLButton();
    }
    
    @Override
    public TextField createTextField() {
        return new HTMLTextField();
    }
    
    @Override
    public Dialog createDialog() {
        return new HTMLDialog();
    }
}'''
add_code_slide(prs, "Abstract Factory - Struktur", abstract_factory_structure)

# Slide 23: Abstract Factory - Telekom Use Case
abstract_factory_usecase = [
    "Telekom Abstract Factory - Service-Bundles:",
    "",
    "TelecomServiceBundleFactory:",
    "• createInternetService() → Internet-Tarif",
    "• createPhoneService() → Telefon-Service", 
    "• createTVService() → Entertainment-Paket",
    "",
    "Konkrete Bundle-Factories:",
    "• BasicBundleFactory → Basis-Services",
    "• PremiumBundleFactory → Premium-Services",
    "• BusinessBundleFactory → Business-Services",
    "",
    "Vorteil: Konsistente Service-Kombinationen"
]
notes_text = "Trainer-Notizen:\n- Abstract Factory für Service-Bundles erklären\n- Konsistenz zwischen verwandten Services betonen\n- Unterschied zu Factory Method verdeutlichen\n- Zeitrahmen: ca. 8 Minuten"
add_content_slide(prs, "Abstract Factory - Telekom Use Case", abstract_factory_usecase, notes_text)

# Slide 24: Builder Pattern - Problem
builder_problem = [
    "Problem: Komplexe Objektkonstruktion",
    "",
    "Telekom Service-Konfiguration mit vielen Parametern:",
    "• 15+ optionale Konfigurationsparameter",
    "• Verschiedene Konfigurationskombinationen",
    "• Validierungen zwischen Parametern",
    "• Schrittweise Konfiguration erforderlich",
    "",
    "Probleme mit Konstruktor-Ansatz:",
    "• Konstruktor mit vielen Parametern (Telescoping)",
    "• Unübersichtliche Parameter-Reihenfolge",
    "• Schwierige Validierung"
]
add_content_slide(prs, "Builder Pattern - Problem", builder_problem)

# Slide 25: Builder Pattern Code
builder_code = '''// Product
public class TelecomServiceConfig {
    private String serviceName;
    private int bandwidth;
    private boolean hasIPv6;
    private String qosLevel;
    private List<String> addOns;
    
    // Private constructor - nur Builder kann erstellen
    private TelecomServiceConfig(Builder builder) {
        this.serviceName = builder.serviceName;
        this.bandwidth = builder.bandwidth;
        this.hasIPv6 = builder.hasIPv6;
        this.qosLevel = builder.qosLevel;
        this.addOns = builder.addOns;
    }
    
    // Builder Class
    public static class Builder {
        private String serviceName;
        private int bandwidth = 100; // Default
        private boolean hasIPv6 = false;
        private String qosLevel = "STANDARD";
        private List<String> addOns = new ArrayList<>();
        
        public Builder(String serviceName) {
            this.serviceName = serviceName;
        }
        
        public Builder bandwidth(int bandwidth) {
            this.bandwidth = bandwidth;
            return this;
        }
        
        public Builder enableIPv6() {
            this.hasIPv6 = true;
            return this;
        }
        
        public Builder qosLevel(String level) {
            this.qosLevel = level;
            return this;
        }
        
        public TelecomServiceConfig build() {
            // Validierungen hier
            return new TelecomServiceConfig(this);
        }
    }
}'''
add_code_slide(prs, "Builder Pattern - Implementation", builder_code)

# Slide 26: Builder Pattern - Usage
builder_usage = '''// Fluent Interface Usage
TelecomServiceConfig config = new TelecomServiceConfig.Builder("DSL-Premium")
    .bandwidth(1000)
    .enableIPv6()
    .qosLevel("PREMIUM")
    .build();

// Director Pattern für vordefinierte Konfigurationen
public class ServiceConfigDirector {
    private TelecomServiceConfig.Builder builder;
    
    public ServiceConfigDirector(TelecomServiceConfig.Builder builder) {
        this.builder = builder;
    }
    
    public TelecomServiceConfig createBasicConfig() {
        return builder
            .bandwidth(100)
            .qosLevel("STANDARD")
            .build();
    }
    
    public TelecomServiceConfig createPremiumConfig() {
        return builder
            .bandwidth(1000)
            .enableIPv6()
            .qosLevel("PREMIUM")
            .build();
    }
}

// Usage with Director
ServiceConfigDirector director = new ServiceConfigDirector(
    new TelecomServiceConfig.Builder("Fiber")
);
TelecomServiceConfig premiumConfig = director.createPremiumConfig();'''
notes_text = "Trainer-Notizen:\n- Fluent Interface Vorteile zeigen\n- Director Pattern als Ergänzung erklären\n- Validierung im build() betonen\n- Zeitrahmen: ca. 10 Minuten"
add_code_slide(prs, "Builder Pattern - Verwendung", builder_usage, "java", notes_text)

# Slide 27: Prototype Pattern - Problem
prototype_problem = [
    "Problem: Kostspielige Objekterstellung",
    "",
    "Telekom-Szenario: Service-Templates laden",
    "• Konfigurationen aus externen Systemen laden",
    "• Komplexe Berechnungen für Service-Parameter",
    "• Datenbankabfragen für Referenzdaten",
    "• Netzwerk-Latenz bei API-Aufrufen",
    "",
    "Traditionelle Lösung: Jedes Mal neu erstellen",
    "• Performance-Problem bei häufiger Nutzung",
    "• Redundante externe Aufrufe",
    "• Ressourcenverschwendung"
]
add_content_slide(prs, "Prototype Pattern - Problem", prototype_problem)

# Slide 28: Prototype Pattern Code
prototype_code = '''// Prototype Interface
public interface ServiceTemplate extends Cloneable {
    ServiceTemplate clone();
    void customize(Map<String, Object> parameters);
}

// Concrete Prototype
public class DSLServiceTemplate implements ServiceTemplate {
    private String templateName;
    private Map<String, Object> baseConfiguration;
    private List<String> requiredComponents;
    
    public DSLServiceTemplate() {
        // Kostspielige Initialisierung
        loadFromExternalSystem();
        calculateBaseParameters();
        loadRequiredComponents();
    }
    
    @Override
    public ServiceTemplate clone() {
        try {
            DSLServiceTemplate cloned = (DSLServiceTemplate) super.clone();
            // Deep Copy für mutable Objekte
            cloned.baseConfiguration = new HashMap<>(this.baseConfiguration);
            cloned.requiredComponents = new ArrayList<>(this.requiredComponents);
            return cloned;
        } catch (CloneNotSupportedException e) {
            throw new RuntimeException("Clone not supported", e);
        }
    }
    
    @Override
    public void customize(Map<String, Object> parameters) {
        // Template an spezifische Anforderungen anpassen
        baseConfiguration.putAll(parameters);
    }
    
    private void loadFromExternalSystem() {
        // Simuliert kostspielige Operation
        try { Thread.sleep(1000); } catch (InterruptedException e) {}
    }
}'''
add_code_slide(prs, "Prototype Pattern - Implementation", prototype_code)

# Slide 29: Prototype Pattern - Deep vs Shallow Copy
copy_types = [
    "Shallow Copy vs Deep Copy - Wichtiger Unterschied:",
    "",
    "Shallow Copy:",
    "• Kopiert nur Referenzen auf Objekte",
    "• Mutable Objekte werden geteilt",
    "• Änderungen betreffen alle Klone",
    "• Performance: Sehr schnell",
    "",
    "Deep Copy:",
    "• Kopiert auch referenzierte Objekte",
    "• Jeder Klon ist vollständig unabhängig",
    "• Keine unerwünschten Seiteneffekte",
    "• Performance: Langsamer, aber sicherer",
    "",
    "Telekom-Regel: Deep Copy für Konfigurationsobjekte"
]
notes_text = "Trainer-Notizen:\n- Unterschied anhand von Beispielen erklären\n- Risiken von Shallow Copy bei mutable Objects\n- Performance-Trade-offs diskutieren\n- Zeitrahmen: ca. 6 Minuten"
add_content_slide(prs, "Prototype - Deep vs Shallow Copy", copy_types, notes_text)

# Slide 30: Singleton Pattern - Problem
singleton_problem = [
    "Problem: Eindeutige Instanz erforderlich",
    "",
    "Telekom-Szenarien für Singleton:",
    "• Configuration Manager - Eine zentrale Konfiguration",
    "• Logger - Einheitliches Logging-System", 
    "• Database Connection Pool - Ressourcen-Management",
    "• License Manager - Lizenz-Validierung",
    "",
    "Herausforderungen:",
    "• Thread-Safety in Multithreading-Umgebung",
    "• Lazy vs Eager Initialization",
    "• Performance vs Sicherheit",
    "• Testbarkeit (schwierig zu mocken)"
]
add_content_slide(prs, "Singleton Pattern - Problem", singleton_problem)

# Slide 31: Singleton Pattern - Thread-Safe Implementation
singleton_code = '''// Thread-Safe Singleton mit Lazy Initialization
public class TelecomConfigManager {
    private static volatile TelecomConfigManager instance;
    private Map<String, String> configurations;
    
    private TelecomConfigManager() {
        // Private constructor verhindert direkte Instantiierung
        configurations = loadConfigurations();
    }
    
    // Double-Checked Locking Pattern
    public static TelecomConfigManager getInstance() {
        if (instance == null) {
            synchronized (TelecomConfigManager.class) {
                if (instance == null) {
                    instance = new TelecomConfigManager();
                }
            }
        }
        return instance;
    }
    
    public String getConfig(String key) {
        return configurations.get(key);
    }
    
    public void setConfig(String key, String value) {
        synchronized (this) {
            configurations.put(key, value);
        }
    }
    
    private Map<String, String> loadConfigurations() {
        // Load from external source
        return new HashMap<>();
    }
}

// Alternative: Enum-based Singleton (recommended)
public enum ConfigManager {
    INSTANCE;
    
    private Map<String, String> configurations;
    
    ConfigManager() {
        configurations = new HashMap<>();
    }
    
    public String getConfig(String key) {
        return configurations.get(key);
    }
}'''
notes_text = "Trainer-Notizen:\n- Double-Checked Locking Pattern erklären\n- Enum-Singleton als beste Praxis\n- Thread-Safety Probleme diskutieren\n- Anti-Pattern: Public static Felder\n- Zeitrahmen: ca. 12 Minuten"
add_code_slide(prs, "Singleton - Thread-Safe Implementation", singleton_code, "java", notes_text)

# Slide 32: Singleton Pattern - Probleme und Alternativen
singleton_problems = [
    "Singleton Pattern - Kritische Betrachtung:",
    "",
    "Probleme:",
    "• Schwierig zu testen (Global State)",
    "• Versteckte Abhängigkeiten",
    "• Verletzung des Single Responsibility Principle",
    "• Threading-Probleme",
    "• Schwierige Vererbung",
    "",
    "Moderne Alternativen:",
    "• Dependency Injection Container",
    "• Monostate Pattern (Shared State)",
    "• Factory mit Instance-Verwaltung",
    "• Spring @Singleton Annotation"
]
notes_text = "Trainer-Notizen:\n- Singleton als Anti-Pattern diskutieren\n- Wann ist Singleton wirklich nötig?\n- Dependency Injection als moderne Alternative\n- Testbarkeit stark betonen\n- Zeitrahmen: ca. 8 Minuten"
add_content_slide(prs, "Singleton - Probleme und Alternativen", singleton_problems, notes_text)

# Slide 33: Creational Patterns Comparison
patterns_comparison = [
    "Creational Patterns - Wann welches Pattern?",
    "",
    "Factory Method:",
    "• Verschiedene, aber ähnliche Objekte",
    "• Subklassen entscheiden über konkrete Typen",
    "",
    "Abstract Factory:",
    "• Familien verwandter Objekte",
    "• Konsistenz zwischen Produkt-Familien",
    "",
    "Builder:",
    "• Komplexe Objekte mit vielen Parametern", 
    "• Schrittweise Konstruktion erforderlich",
    "",
    "Prototype:",
    "• Kostspielige Objekterstellung",
    "• Ähnliche Objekte mit kleinen Variationen",
    "",
    "Singleton:",
    "• Eindeutige Instanz zwingend erforderlich",
    "• Mit Vorsicht verwenden!"
]
add_content_slide(prs, "Creational Patterns - Vergleich", patterns_comparison)

# Slide 34: Pattern Selection Guide
selection_guide = [
    "Entscheidungshilfe - Pattern Selection:",
    "",
    "Fragen zur Pattern-Auswahl:",
    "",
    "1. Brauche ich verschiedene Varianten ähnlicher Objekte?",
    "   → Factory Method",
    "",
    "2. Brauche ich Familien zusammengehöriger Objekte?", 
    "   → Abstract Factory",
    "",
    "3. Hat mein Objekt sehr viele Konfigurationsoptionen?",
    "   → Builder",
    "",
    "4. Ist die Objekterstellung sehr kostspielig?",
    "   → Prototype",
    "",
    "5. Brauche ich wirklich nur eine einzige Instanz?",
    "   → Singleton (mit Vorsicht!)"
]
notes_text = "Trainer-Notizen:\n- Entscheidungsbaum interaktiv durchgehen\n- Teilnehmer eigene Beispiele finden lassen\n- Pattern-Kombinationen erwähnen\n- Telekom-Kontext: Welche Patterns wo einsetzen?\n- Zeitrahmen: ca. 8 Minuten"
add_content_slide(prs, "Pattern Selection Guide", selection_guide, notes_text)

# ============================================================================
# END OF CREATIONAL PATTERNS SECTION
# ============================================================================

# ============================================================================
# STRUCTURAL PATTERNS SECTION (25-30 slides)
# ============================================================================

# Slide 35: Structural Patterns Section Divider
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Design Patterns"
for paragraph in title.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = TITLE_FONT
        run.font.size = Pt(48)
        run.font.bold = False
        run.font.color.rgb = RGBColor(68, 68, 68)

subtitle = None
for shape in slide.shapes:
    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 4:  # SUBTITLE
        subtitle = shape
        break

if subtitle:
    subtitle.text = "Structural Patterns\n\nObjektzusammensetzung und -beziehungen"
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = CONTENT_FONT
            run.font.size = Pt(28)
            run.font.bold = False
            run.font.color.rgb = RGBColor(100, 100, 100)
        paragraph.alignment = PP_ALIGN.CENTER

# Slide 36: Structural Patterns Overview
structural_overview = [
    "Adapter - Legacy-Systeme in neue Architekturen integrieren",
    "Bridge - Abstraktion von Implementierung trennen",
    "Composite - Hierarchische Baumstrukturen verwalten",
    "Decorator - Verhalten dynamisch zur Laufzeit erweitern", 
    "Facade - Komplexe Subsysteme vereinfacht zugänglich machen",
    "Flyweight - Speicheroptimierung durch geteilte Objekte",
    "Proxy - Stellvertreter für aufwändige oder entfernte Objekte",
    "",
    "Gemeinsames Ziel: Flexible Objektzusammensetzung",
    "Entkopplung zwischen Komponenten und deren Verwendung"
]
notes_text = "Trainer-Notizen:\n- Überblick über alle 7 Structural Patterns\n- Fokus auf Objektbeziehungen und -zusammensetzung\n- Unterschied zu Creational Patterns erklären\n- Telekom-Kontext: Legacy-Integration, Service-Orchestrierung\n- Zeitrahmen: ca. 6 Minuten"
add_content_slide(prs, "Structural Patterns - Überblick", structural_overview, notes_text)

# Slide 37: Adapter Pattern - Problem
adapter_problem = [
    "Problem: Integration inkompatibler Interfaces",
    "",
    "Telekom-Szenario: Legacy-System Integration",
    "• Altes Billing-System mit proprietärer Schnittstelle",
    "• Neue Microservice-Architektur mit REST-APIs",
    "• Verschiedene Datenformate (XML vs JSON)",
    "• Inkompatible Methodensignaturen",
    "",
    "Herausforderungen:",
    "• Bestehender Code soll nicht verändert werden",
    "• Neue Services erwarten moderne Interfaces",
    "• Mehrere Legacy-Systeme müssen integriert werden",
    "• Schrittweise Migration erforderlich"
]
add_content_slide(prs, "Adapter Pattern - Problem", adapter_problem)

# Slide 38: Adapter Pattern - Structure
adapter_structure = [
    "Adapter Pattern - Struktur:",
    "",
    "• Target - Interface, das der Client erwartet",
    "• Client - Verwendet das Target Interface",
    "• Adaptee - Legacy-Klasse mit inkompatiblem Interface",
    "• Adapter - Vermittelt zwischen Target und Adaptee",
    "",
    "Zwei Varianten:",
    "• Object Adapter - Komposition (empfohlen)",
    "• Class Adapter - Mehrfachvererbung (Java: nicht möglich)",
    "",
    "Vorteil: Legacy-Code bleibt unverändert"
]
add_content_slide(prs, "Adapter Pattern - Struktur", adapter_structure)

# Slide 39: Adapter Pattern Code - Legacy System
adapter_legacy_code = '''// Legacy Billing System (Adaptee)
public class LegacyBillingSystem {
    public void processLegacyPayment(String customerId, double amount, String currency) {
        System.out.println("Processing legacy payment:");
        System.out.println("Customer: " + customerId);
        System.out.println("Amount: " + amount + " " + currency);
        
        // Complex legacy processing logic
        validateCustomerInLegacyDB(customerId);
        processLegacyTransaction(amount, currency);
        updateLegacyBillingRecord(customerId, amount);
    }
    
    public String getLegacyCustomerInfo(String customerId) {
        // Returns XML format
        return "<?xml version='1.0'?><customer><id>" + customerId + 
               "</id><status>ACTIVE</status></customer>";
    }
    
    private void validateCustomerInLegacyDB(String customerId) {
        // Legacy database operations
    }
    
    private void processLegacyTransaction(double amount, String currency) {
        // Complex legacy business logic
    }
}

// Modern Interface (Target)
public interface ModernPaymentProcessor {
    PaymentResult processPayment(PaymentRequest request);
    CustomerInfo getCustomer(String customerId);
}

public class PaymentRequest {
    private String customerId;
    private BigDecimal amount;
    private String currency;
    // constructors, getters, setters
}'''
add_code_slide(prs, "Adapter Pattern - Legacy System", adapter_legacy_code)

# Slide 40: Adapter Pattern Code - Implementation
adapter_implementation = '''// Adapter Implementation
public class BillingSystemAdapter implements ModernPaymentProcessor {
    private LegacyBillingSystem legacySystem;
    
    public BillingSystemAdapter(LegacyBillingSystem legacySystem) {
        this.legacySystem = legacySystem;
    }
    
    @Override
    public PaymentResult processPayment(PaymentRequest request) {
        try {
            // Adapt modern interface to legacy interface
            legacySystem.processLegacyPayment(
                request.getCustomerId(),
                request.getAmount().doubleValue(),
                request.getCurrency()
            );
            
            return new PaymentResult(true, "Payment processed successfully");
        } catch (Exception e) {
            return new PaymentResult(false, "Payment failed: " + e.getMessage());
        }
    }
    
    @Override
    public CustomerInfo getCustomer(String customerId) {
        // Get legacy XML format
        String xmlData = legacySystem.getLegacyCustomerInfo(customerId);
        
        // Convert XML to modern object format
        return parseXMLToCustomerInfo(xmlData);
    }
    
    private CustomerInfo parseXMLToCustomerInfo(String xmlData) {
        // Parse XML and create CustomerInfo object
        // Using XML parser to extract customer data
        return new CustomerInfo(extractId(xmlData), extractStatus(xmlData));
    }
}

// Usage
ModernPaymentProcessor processor = new BillingSystemAdapter(
    new LegacyBillingSystem()
);

PaymentRequest request = new PaymentRequest("CUST123", 
    new BigDecimal("99.99"), "EUR");
PaymentResult result = processor.processPayment(request);'''
notes_text = "Trainer-Notizen:\n- Object Adapter Pattern zeigen\n- Datenkonvertierung zwischen Formaten\n- Legacy-System bleibt unverändert betonen\n- Telekom-Beispiel: Billing-System Integration\n- Zeitrahmen: ca. 10 Minuten"
add_code_slide(prs, "Adapter Pattern - Implementation", adapter_implementation, "java", notes_text)

# Slide 41: Bridge Pattern - Problem
bridge_problem = [
    "Problem: Abstraktion und Implementierung gekoppelt",
    "",
    "Telekom-Szenario: Notification-System",
    "• Verschiedene Nachrichtentypen (Email, SMS, Push)",
    "• Verschiedene Implementierungen (Inhouse, Amazon SES, Twilio)",
    "• Jede Kombination würde eigene Klasse benötigen",
    "• EmailAmazonSES, EmailTwilio, SMSAmazonSNS, etc.",
    "",
    "Probleme ohne Bridge:",
    "• Klassen-Explosion bei n×m Kombinationen",
    "• Änderungen in Abstraktion betreffen alle Implementierungen",
    "• Schwierige Erweiterung um neue Typen/Anbieter",
    "• Code-Duplikation in ähnlichen Klassen"
]
add_content_slide(prs, "Bridge Pattern - Problem", bridge_problem)

# Slide 42: Bridge Pattern - Structure
bridge_structure = [
    "Bridge Pattern - Struktur:",
    "",
    "• Abstraction - Definiert Interface für Client",
    "• Refined Abstraction - Erweiterte Abstraktion",
    "• Implementor - Interface für Implementierung",
    "• Concrete Implementor - Spezifische Implementierung",
    "",
    "Kernidee:",
    "• Trennung von 'Was' (Abstraktion) und 'Wie' (Implementierung)",
    "• Komposition statt Vererbung",
    "• Unabhängige Erweiterung beider Dimensionen",
    "",
    "Vorteil: n + m Klassen statt n × m Kombinationen"
]
add_content_slide(prs, "Bridge Pattern - Struktur", bridge_structure)

# Slide 43: Bridge Pattern Code - Implementation
bridge_implementation = '''// Implementor Interface
public interface NotificationSender {
    void sendMessage(String recipient, String subject, String content);
    boolean isAvailable();
}

// Concrete Implementors
public class EmailSESImplementation implements NotificationSender {
    @Override
    public void sendMessage(String recipient, String subject, String content) {
        System.out.println("Sending via Amazon SES:");
        System.out.println("To: " + recipient);
        System.out.println("Subject: " + subject);
        // Amazon SES specific implementation
        sendViaSES(recipient, subject, content);
    }
    
    @Override
    public boolean isAvailable() {
        return checkSESConnection();
    }
}

public class SMSTwilioImplementation implements NotificationSender {
    @Override
    public void sendMessage(String recipient, String subject, String content) {
        System.out.println("Sending SMS via Twilio:");
        System.out.println("To: " + recipient);
        System.out.println("Message: " + content);
        // Twilio SMS specific implementation
        sendViaTwilio(recipient, content);
    }
    
    @Override
    public boolean isAvailable() {
        return checkTwilioAPI();
    }
}

// Abstraction
public abstract class Notification {
    protected NotificationSender sender;
    
    public Notification(NotificationSender sender) {
        this.sender = sender;
    }
    
    public abstract void send(String recipient, String content);
}'''
add_code_slide(prs, "Bridge Pattern - Implementation", bridge_implementation)

# Slide 44: Bridge Pattern Code - Refined Abstractions
bridge_refined_abstractions = '''// Refined Abstractions
public class UrgentNotification extends Notification {
    public UrgentNotification(NotificationSender sender) {
        super(sender);
    }
    
    @Override
    public void send(String recipient, String content) {
        if (!sender.isAvailable()) {
            throw new RuntimeException("Urgent notification sender not available!");
        }
        
        String urgentContent = "🚨 URGENT: " + content;
        sender.sendMessage(recipient, "URGENT NOTIFICATION", urgentContent);
        
        // Retry logic for urgent notifications
        if (!wasDelivered()) {
            retryWithBackup(recipient, urgentContent);
        }
    }
}

public class MarketingNotification extends Notification {
    private boolean respectOptOut = true;
    
    public MarketingNotification(NotificationSender sender) {
        super(sender);
    }
    
    @Override
    public void send(String recipient, String content) {
        if (respectOptOut && hasOptedOut(recipient)) {
            System.out.println("User has opted out, skipping marketing notification");
            return;
        }
        
        String marketingContent = content + "\\n\\n📧 Unsubscribe: link";
        sender.sendMessage(recipient, "Telekom Marketing", marketingContent);
    }
}

// Usage - Flexible Combinations
Notification emailUrgent = new UrgentNotification(new EmailSESImplementation());
Notification smsUrgent = new UrgentNotification(new SMSTwilioImplementation());
Notification emailMarketing = new MarketingNotification(new EmailSESImplementation());

emailUrgent.send("customer@example.com", "Service disruption detected");
smsUrgent.send("+4917012345678", "Critical system alert");'''
notes_text = "Trainer-Notizen:\n- Flexible Kombinationen demonstrieren\n- Unabhängige Erweiterung beider Dimensionen\n- Telekom-Kontext: Verschiedene Notification-Anbieter\n- Unterschied zu Strategy Pattern diskutieren\n- Zeitrahmen: ca. 10 Minuten"
add_code_slide(prs, "Bridge Pattern - Refined Abstractions", bridge_refined_abstractions, "java", notes_text)

# Slide 45: Composite Pattern - Problem
composite_problem = [
    "Problem: Hierarchische Strukturen einheitlich behandeln",
    "",
    "Telekom-Szenario: Organisationsstruktur & Services",
    "• Kunden können Individual-Personen sein",
    "• Kunden können Firmen mit Abteilungen sein",
    "• Firmen können Tochter-Unternehmen haben",
    "• Services können einzeln oder als Pakete angeboten werden",
    "",
    "Herausforderungen:",
    "• Unterschiedliche Behandlung von Einzelobjekten vs Gruppen",
    "• Rekursive Operationen (Gesamtkosten berechnen)",
    "• Code-Duplikation bei ähnlichen Operationen",
    "• Schwierige Navigation in der Hierarchie"
]
add_content_slide(prs, "Composite Pattern - Problem", composite_problem)

# Slide 46: Composite Pattern - Structure
composite_structure = [
    "Composite Pattern - Struktur:",
    "",
    "• Component - Gemeinsames Interface für alle Objekte",
    "• Leaf - Einzelne Objekte ohne Kinder",
    "• Composite - Container-Objekte mit Kindern",
    "• Client - Behandelt alle Objekte einheitlich",
    "",
    "Kernprinzip:",
    "• Einheitliches Interface für Einzelobjekte und Gruppen",
    "• Rekursive Komposition möglich",
    "• Teil-Ganzes-Hierarchien einfach darstellbar",
    "",
    "Vorteil: Client muss nicht zwischen Typen unterscheiden"
]
add_content_slide(prs, "Composite Pattern - Struktur", composite_structure)

# Slide 47: Composite Pattern Code - Component Interface
composite_component_code = '''// Component Interface
public abstract class TelecomEntity {
    protected String name;
    protected String entityId;
    
    public TelecomEntity(String name, String entityId) {
        this.name = name;
        this.entityId = entityId;
    }
    
    // Operationen für alle Entities
    public abstract BigDecimal calculateTotalCost();
    public abstract int getServiceCount();
    public abstract List<String> getAllServices();
    public abstract void printStructure(int indent);
    
    // Composite-spezifische Operationen (default implementation)
    public void addEntity(TelecomEntity entity) {
        throw new UnsupportedOperationException("Cannot add to leaf entity");
    }
    
    public void removeEntity(TelecomEntity entity) {
        throw new UnsupportedOperationException("Cannot remove from leaf entity");
    }
    
    public List<TelecomEntity> getChildren() {
        return Collections.emptyList();
    }
    
    // Getter
    public String getName() { return name; }
    public String getEntityId() { return entityId; }
}

// Leaf - Individual Customer/Service
public class IndividualCustomer extends TelecomEntity {
    private BigDecimal baseFee;
    private List<String> subscribedServices;
    
    public IndividualCustomer(String name, String customerId, BigDecimal baseFee) {
        super(name, customerId);
        this.baseFee = baseFee;
        this.subscribedServices = new ArrayList<>();
    }
    
    @Override
    public BigDecimal calculateTotalCost() {
        return baseFee;
    }
    
    @Override
    public int getServiceCount() {
        return subscribedServices.size();
    }
    
    @Override
    public List<String> getAllServices() {
        return new ArrayList<>(subscribedServices);
    }
    
    @Override
    public void printStructure(int indent) {
        String indentation = "  ".repeat(indent);
        System.out.println(indentation + "👤 " + name + " (€" + baseFee + ")");
    }
}'''
add_code_slide(prs, "Composite Pattern - Component Interface", composite_component_code)

# Slide 48: Composite Pattern Code - Composite Implementation
composite_implementation = '''// Composite - Organization/Service Package
public class CorporateCustomer extends TelecomEntity {
    private List<TelecomEntity> children;
    private BigDecimal managementFee;
    
    public CorporateCustomer(String name, String customerId, BigDecimal managementFee) {
        super(name, customerId);
        this.children = new ArrayList<>();
        this.managementFee = managementFee;
    }
    
    @Override
    public void addEntity(TelecomEntity entity) {
        children.add(entity);
    }
    
    @Override
    public void removeEntity(TelecomEntity entity) {
        children.remove(entity);
    }
    
    @Override
    public List<TelecomEntity> getChildren() {
        return new ArrayList<>(children);
    }
    
    @Override
    public BigDecimal calculateTotalCost() {
        BigDecimal total = managementFee;
        for (TelecomEntity child : children) {
            total = total.add(child.calculateTotalCost());
        }
        return total;
    }
    
    @Override
    public int getServiceCount() {
        int total = 0;
        for (TelecomEntity child : children) {
            total += child.getServiceCount();
        }
        return total;
    }
    
    @Override
    public List<String> getAllServices() {
        List<String> allServices = new ArrayList<>();
        for (TelecomEntity child : children) {
            allServices.addAll(child.getAllServices());
        }
        return allServices;
    }
    
    @Override
    public void printStructure(int indent) {
        String indentation = "  ".repeat(indent);
        System.out.println(indentation + "🏢 " + name + " (Management: €" + managementFee + ")");
        for (TelecomEntity child : children) {
            child.printStructure(indent + 1);
        }
    }
}'''
add_code_slide(prs, "Composite Pattern - Composite Implementation", composite_implementation)

# Slide 49: Composite Pattern - Usage Example
composite_usage = '''// Usage Example - Building Hierarchy
public class CompositePatternDemo {
    public static void main(String[] args) {
        // Individual customers
        TelecomEntity individual1 = new IndividualCustomer("Max Mustermann", "IND001", 
            new BigDecimal("29.99"));
        TelecomEntity individual2 = new IndividualCustomer("Anna Schmidt", "IND002", 
            new BigDecimal("49.99"));
        
        // Small department
        CorporateCustomer itDepartment = new CorporateCustomer("IT Department", "DEPT001", 
            new BigDecimal("100.00"));
        itDepartment.addEntity(individual1);
        itDepartment.addEntity(individual2);
        
        // Large corporation
        CorporateCustomer telekomCorp = new CorporateCustomer("Telekom AG", "CORP001", 
            new BigDecimal("500.00"));
        telekomCorp.addEntity(itDepartment);
        
        // Add more departments
        CorporateCustomer hrDepartment = new CorporateCustomer("HR Department", "DEPT002", 
            new BigDecimal("80.00"));
        hrDepartment.addEntity(new IndividualCustomer("HR Manager", "IND003", 
            new BigDecimal("39.99")));
        telekomCorp.addEntity(hrDepartment);
        
        // Client treats all entities uniformly
        System.out.println("Structure:");
        telekomCorp.printStructure(0);
        
        System.out.println("\\nTotal Cost: €" + telekomCorp.calculateTotalCost());
        System.out.println("Total Services: " + telekomCorp.getServiceCount());
        System.out.println("All Services: " + telekomCorp.getAllServices());
    }
}

// Output:
// 🏢 Telekom AG (Management: €500)
//   🏢 IT Department (Management: €100)
//     👤 Max Mustermann (€29.99)
//     👤 Anna Schmidt (€49.99)
//   🏢 HR Department (Management: €80)
//     👤 HR Manager (€39.99)
//
// Total Cost: €799.97'''
notes_text = "Trainer-Notizen:\n- Einheitliche Behandlung aller Objekte demonstrieren\n- Rekursive Operationen zeigen\n- Hierarchie-Aufbau interaktiv entwickeln\n- Telekom-Kontext: Kunden-Organisationen\n- Zeitrahmen: ca. 10 Minuten"
add_code_slide(prs, "Composite Pattern - Usage Example", composite_usage, "java", notes_text)

# Slide 50: Decorator Pattern - Problem
decorator_problem = [
    "Problem: Verhalten dynamisch zur Laufzeit erweitern",
    "",
    "Telekom-Szenario: Service-Erweiterungen",
    "• Basis-DSL-Service (100 Mbit/s)",
    "• Optionale Erweiterungen: IPv6, QoS, Monitoring, Backup",
    "• Kunden können beliebige Kombinationen wählen",
    "• Verhalten soll zur Laufzeit hinzufügbar sein",
    "",
    "Probleme mit Vererbung:",
    "• 2^n Klassen für n Erweiterungen",
    "• DSLServiceWithIPv6WithQoS, DSLServiceWithMonitoring...",
    "• Statische Kombination zur Compile-Zeit",
    "• Schwierige Wartung bei vielen Kombinationen"
]
add_content_slide(prs, "Decorator Pattern - Problem", decorator_problem)

# Slide 51: Decorator Pattern - Structure
decorator_structure = [
    "Decorator Pattern - Struktur:",
    "",
    "• Component - Interface für Basis-Objekt und Dekorierer",
    "• Concrete Component - Basis-Implementierung",
    "• Decorator - Abstrakte Basis für alle Dekorierer",
    "• Concrete Decorator - Spezifische Erweiterungen",
    "",
    "Kernprinzip:",
    "• Wrapper-Objekte erweitern andere Objekte",
    "• Gleiche Schnittstelle wie das ursprüngliche Objekt",
    "• Ketten von Dekorierern möglich",
    "• Zur Laufzeit zusammenstellbar",
    "",
    "Vorteil: Flexible Kombination von Verhalten"
]
add_content_slide(prs, "Decorator Pattern - Struktur", decorator_structure)

# Slide 52: Decorator Pattern Code - Component and Base
decorator_component_code = '''// Component Interface
public interface TelecomService {
    String getDescription();
    BigDecimal getCost();
    Map<String, Object> getConfiguration();
    void activate();
    void deactivate();
}

// Concrete Component - Base Service
public class DSLBasicService implements TelecomService {
    private String customerId;
    private int bandwidth;
    
    public DSLBasicService(String customerId, int bandwidth) {
        this.customerId = customerId;
        this.bandwidth = bandwidth;
    }
    
    @Override
    public String getDescription() {
        return "DSL Basic Service (" + bandwidth + " Mbit/s)";
    }
    
    @Override
    public BigDecimal getCost() {
        return new BigDecimal("39.99");
    }
    
    @Override
    public Map<String, Object> getConfiguration() {
        Map<String, Object> config = new HashMap<>();
        config.put("customerId", customerId);
        config.put("bandwidth", bandwidth);
        config.put("type", "DSL_BASIC");
        return config;
    }
    
    @Override
    public void activate() {
        System.out.println("Activating DSL Basic Service for " + customerId);
        // Service activation logic
    }
    
    @Override
    public void deactivate() {
        System.out.println("Deactivating DSL Basic Service for " + customerId);
        // Service deactivation logic
    }
}'''
add_code_slide(prs, "Decorator Pattern - Component & Base", decorator_component_code)

# Slide 53: Decorator Pattern Code - Base Decorator
decorator_base_decorator = '''// Base Decorator
public abstract class ServiceDecorator implements TelecomService {
    protected TelecomService wrappedService;
    
    public ServiceDecorator(TelecomService service) {
        this.wrappedService = service;
    }
    
    @Override
    public String getDescription() {
        return wrappedService.getDescription();
    }
    
    @Override
    public BigDecimal getCost() {
        return wrappedService.getCost();
    }
    
    @Override
    public Map<String, Object> getConfiguration() {
        return new HashMap<>(wrappedService.getConfiguration());
    }
    
    @Override
    public void activate() {
        wrappedService.activate();
    }
    
    @Override
    public void deactivate() {
        wrappedService.deactivate();
    }
}

// Concrete Decorator - IPv6 Support
public class IPv6ServiceDecorator extends ServiceDecorator {
    public IPv6ServiceDecorator(TelecomService service) {
        super(service);
    }
    
    @Override
    public String getDescription() {
        return wrappedService.getDescription() + " + IPv6 Support";
    }
    
    @Override
    public BigDecimal getCost() {
        return wrappedService.getCost().add(new BigDecimal("9.99"));
    }
    
    @Override
    public Map<String, Object> getConfiguration() {
        Map<String, Object> config = super.getConfiguration();
        config.put("ipv6Enabled", true);
        config.put("ipv6Prefix", "2001:db8::/64");
        return config;
    }
    
    @Override
    public void activate() {
        super.activate();
        System.out.println("Configuring IPv6 support...");
        // IPv6 specific activation
    }
}'''
add_code_slide(prs, "Decorator Pattern - Base Decorator", decorator_base_decorator)

# Slide 54: Decorator Pattern Code - Multiple Decorators
decorator_multiple = '''// QoS Decorator
public class QoSServiceDecorator extends ServiceDecorator {
    private String qosLevel;
    
    public QoSServiceDecorator(TelecomService service, String qosLevel) {
        super(service);
        this.qosLevel = qosLevel;
    }
    
    @Override
    public String getDescription() {
        return wrappedService.getDescription() + " + QoS (" + qosLevel + ")";
    }
    
    @Override
    public BigDecimal getCost() {
        BigDecimal qosCost = "PREMIUM".equals(qosLevel) ? 
            new BigDecimal("19.99") : new BigDecimal("9.99");
        return wrappedService.getCost().add(qosCost);
    }
    
    @Override
    public void activate() {
        super.activate();
        System.out.println("Setting up QoS level: " + qosLevel);
    }
}

// Monitoring Decorator
public class MonitoringServiceDecorator extends ServiceDecorator {
    public MonitoringServiceDecorator(TelecomService service) {
        super(service);
    }
    
    @Override
    public String getDescription() {
        return wrappedService.getDescription() + " + 24/7 Monitoring";
    }
    
    @Override
    public BigDecimal getCost() {
        return wrappedService.getCost().add(new BigDecimal("14.99"));
    }
    
    @Override
    public void activate() {
        super.activate();
        System.out.println("Starting monitoring services...");
        setupMonitoringAgent();
        configureAlerts();
    }
    
    private void setupMonitoringAgent() { /* monitoring setup */ }
    private void configureAlerts() { /* alert configuration */ }
}

// Usage - Flexible Decoration
TelecomService service = new DSLBasicService("CUST123", 100);

// Add decorations dynamically
service = new IPv6ServiceDecorator(service);
service = new QoSServiceDecorator(service, "PREMIUM");
service = new MonitoringServiceDecorator(service);

System.out.println(service.getDescription());
System.out.println("Total cost: €" + service.getCost());
service.activate();'''
notes_text = "Trainer-Notizen:\n- Flexible Decorator-Kombination zeigen\n- Zur Laufzeit zusammenstellbar betonen\n- Unterschied zu Vererbung verdeutlichen\n- Telekom-Kontext: Service-Add-ons\n- Zeitrahmen: ca. 12 Minuten"
add_code_slide(prs, "Decorator Pattern - Multiple Decorators", decorator_multiple, "java", notes_text)

# Slide 55: Facade Pattern - Problem
facade_problem = [
    "Problem: Komplexe Subsysteme schwer zu verwenden",
    "",
    "Telekom-Szenario: Kunden-Onboarding",
    "• Kundendaten-Validierung (CRM-System)",
    "• Kreditprüfung (Scoring-Service)",
    "• Vertragsabwicklung (Contract-Management)",
    "• Service-Aktivierung (Provisioning-System)",
    "• Billing-Setup (Billing-Engine)",
    "",
    "Probleme ohne Facade:",
    "• Client muss alle Subsysteme kennen",
    "• Komplexe Abhängigkeiten zwischen Systemen",
    "• Fehleranfällige Orchestrierung",
    "• Code-Duplikation bei ähnlichen Workflows"
]
add_content_slide(prs, "Facade Pattern - Problem", facade_problem)

# Slide 56: Facade Pattern Code - Subsystems
facade_subsystems_code = '''// Complex Subsystems
public class CustomerValidationService {
    public boolean validatePersonalData(String name, String address, String id) {
        System.out.println("Validating personal data for: " + name);
        // Complex validation logic
        return performDataValidation(name, address, id);
    }
    
    public boolean checkDuplicateCustomer(String email, String phone) {
        System.out.println("Checking for duplicate customer...");
        return !customerExists(email, phone);
    }
}

public class CreditScoringService {
    public CreditScore performCreditCheck(String customerId, String socialSecurityNumber) {
        System.out.println("Performing credit check for: " + customerId);
        // External credit scoring API call
        return callCreditScoringAPI(socialSecurityNumber);
    }
    
    public boolean isEligibleForPremiumServices(CreditScore score) {
        return score.getScore() >= 700;
    }
}

public class ContractManagementService {
    public Contract createContract(String customerId, String serviceType, BigDecimal amount) {
        System.out.println("Creating contract for: " + customerId);
        // Complex contract generation
        return generateContractDocument(customerId, serviceType, amount);
    }
    
    public void sendContractForSignature(Contract contract) {
        System.out.println("Sending contract for digital signature...");
        // Electronic signature workflow
    }
}

public class ServiceProvisioningSystem {
    public void activateService(String customerId, String serviceType, Map<String, Object> config) {
        System.out.println("Activating " + serviceType + " for: " + customerId);
        // Service activation logic
        configureNetworkElements(customerId, config);
        updateInventorySystem(customerId, serviceType);
    }
}

public class BillingEngine {
    public void setupBillingAccount(String customerId, Contract contract) {
        System.out.println("Setting up billing for: " + customerId);
        // Billing account creation
        createBillingProfile(customerId, contract);
        scheduleBillingCycle(customerId);
    }
}'''
add_code_slide(prs, "Facade Pattern - Complex Subsystems", facade_subsystems_code)

# Slide 57: Facade Pattern Code - Facade Implementation
facade_implementation_code = '''// Facade - Simplified Interface
public class CustomerOnboardingFacade {
    // Subsystem references
    private CustomerValidationService validationService;
    private CreditScoringService scoringService;
    private ContractManagementService contractService;
    private ServiceProvisioningSystem provisioningSystem;
    private BillingEngine billingEngine;
    
    public CustomerOnboardingFacade() {
        this.validationService = new CustomerValidationService();
        this.scoringService = new CreditScoringService();
        this.contractService = new ContractManagementService();
        this.provisioningSystem = new ServiceProvisioningSystem();
        this.billingEngine = new BillingEngine();
    }
    
    // Simplified high-level operation
    public OnboardingResult onboardNewCustomer(CustomerData customerData, 
                                               String serviceType, 
                                               BigDecimal monthlyFee) {
        try {
            // Step 1: Validate customer data
            boolean dataValid = validationService.validatePersonalData(
                customerData.getName(), 
                customerData.getAddress(), 
                customerData.getId()
            );
            
            if (!dataValid) {
                return new OnboardingResult(false, "Invalid customer data");
            }
            
            // Step 2: Check for duplicates
            boolean noDuplicates = validationService.checkDuplicateCustomer(
                customerData.getEmail(), 
                customerData.getPhone()
            );
            
            if (!noDuplicates) {
                return new OnboardingResult(false, "Customer already exists");
            }
            
            // Step 3: Credit check
            CreditScore score = scoringService.performCreditCheck(
                customerData.getId(), 
                customerData.getSocialSecurityNumber()
            );
            
            if (score.getScore() < 500) {
                return new OnboardingResult(false, "Credit check failed");
            }
            
            // Step 4: Create contract
            Contract contract = contractService.createContract(
                customerData.getId(), serviceType, monthlyFee
            );
            contractService.sendContractForSignature(contract);
            
            // Step 5: Setup billing
            billingEngine.setupBillingAccount(customerData.getId(), contract);
            
            // Step 6: Activate service (if premium eligible)
            if (scoringService.isEligibleForPremiumServices(score)) {
                Map<String, Object> config = createPremiumConfig();
                provisioningSystem.activateService(customerData.getId(), 
                    serviceType + "_PREMIUM", config);
            } else {
                provisioningSystem.activateService(customerData.getId(), 
                    serviceType, Collections.emptyMap());
            }
            
            return new OnboardingResult(true, "Customer onboarded successfully");
            
        } catch (Exception e) {
            return new OnboardingResult(false, "Onboarding failed: " + e.getMessage());
        }
    }
}'''
add_code_slide(prs, "Facade Pattern - Facade Implementation", facade_implementation_code)

# Slide 58: Facade Pattern - Usage and Benefits
facade_benefits = [
    "Facade Pattern - Verwendung und Vorteile:",
    "",
    "Einfache Verwendung:",
    "```java",
    "CustomerOnboardingFacade facade = new CustomerOnboardingFacade();",
    "CustomerData data = new CustomerData(/*...*/);",
    "OnboardingResult result = facade.onboardNewCustomer(data, \"DSL\", €39.99);",
    "```",
    "",
    "Vorteile:",
    "• Vereinfachte Schnittstelle für komplexe Operationen",
    "• Entkopplung zwischen Client und Subsystemen",
    "• Zentrale Orchestrierung der Geschäftslogik",
    "• Fehlerbehandlung an einer Stelle",
    "• Einfachere Wartung und Tests",
    "",
    "Telekom-Nutzen: One-Click Kunden-Onboarding"
]
notes_text = "Trainer-Notizen:\n- Vereinfachung für Clients betonen\n- Orchestrierung komplexer Workflows zeigen\n- Unterschied zu Adapter Pattern erklären\n- Telekom-Kontext: Kundenprozesse vereinfachen\n- Zeitrahmen: ca. 10 Minuten"
add_content_slide(prs, "Facade Pattern - Usage & Benefits", facade_benefits, notes_text)

# Slide 59: Flyweight Pattern - Problem
flyweight_problem = [
    "Problem: Speicherverbrauch bei vielen ähnlichen Objekten",
    "",
    "Telekom-Szenario: Netzwerk-Topologie Darstellung",
    "• 50.000+ Router in der Netzwerk-Karte",
    "• Jeder Router hat Position, Status, Verbindungen",
    "• Viele Router des gleichen Typs (Cisco 2900, Huawei AR)",
    "• Gemeinsame Eigenschaften: Icon, Spezifikationen, Rendering",
    "",
    "Problem ohne Flyweight:",
    "• Jeder Router-Objekt speichert redundante Daten",
    "• Icon-Daten millionenfach im Speicher",
    "• Hoher Memory-Footprint",
    "• Schlechte Performance bei Darstellung"
]
add_content_slide(prs, "Flyweight Pattern - Problem", flyweight_problem)

# Slide 60: Flyweight Pattern - Structure
flyweight_structure = [
    "Flyweight Pattern - Struktur:",
    "",
    "• Flyweight - Interface für intrinsische Daten",
    "• Concrete Flyweight - Implementierung mit geteilten Daten",
    "• Context - Externe (extrinsische) Daten",
    "• Flyweight Factory - Verwaltet und teilt Flyweights",
    "",
    "Kernkonzept:",
    "• Intrinsische Daten: Unveränderlich, geteilt",
    "• Extrinsische Daten: Veränderlich, als Parameter übergeben",
    "• Factory Pattern für Flyweight-Verwaltung",
    "",
    "Vorteil: Drastische Speicher-Reduzierung bei vielen Objekten"
]
add_content_slide(prs, "Flyweight Pattern - Struktur", flyweight_structure)

# Slide 61: Flyweight Pattern Code - Implementation
flyweight_implementation_code = '''// Flyweight Interface
public interface NetworkDeviceFlyweight {
    void render(NetworkDeviceContext context, Graphics2D g2d);
    String getDeviceType();
    Map<String, Object> getSpecifications();
}

// Concrete Flyweight - Shared intrinsic data
public class RouterFlyweight implements NetworkDeviceFlyweight {
    private final String deviceType;
    private final String manufacturer;
    private final BufferedImage icon;
    private final Map<String, Object> specifications;
    
    public RouterFlyweight(String deviceType, String manufacturer) {
        this.deviceType = deviceType;
        this.manufacturer = manufacturer;
        this.icon = loadDeviceIcon(deviceType); // Expensive operation
        this.specifications = loadSpecifications(deviceType);
        
        System.out.println("Created flyweight for: " + deviceType);
    }
    
    @Override
    public void render(NetworkDeviceContext context, Graphics2D g2d) {
        // Use extrinsic data from context
        int x = context.getX();
        int y = context.getY();
        String status = context.getStatus();
        
        // Render using intrinsic data
        g2d.drawImage(icon, x, y, null);
        
        // Status-dependent rendering
        if ("OFFLINE".equals(status)) {
            g2d.setColor(Color.RED);
            g2d.fillOval(x + 20, y + 20, 10, 10);
        } else if ("WARNING".equals(status)) {
            g2d.setColor(Color.YELLOW);
            g2d.fillOval(x + 20, y + 20, 10, 10);
        }
        
        // Draw device label
        g2d.setColor(Color.BLACK);
        g2d.drawString(context.getDeviceName(), x, y - 5);
    }
    
    @Override
    public String getDeviceType() { return deviceType; }
    
    @Override
    public Map<String, Object> getSpecifications() { return specifications; }
}

// Context - Extrinsic data
public class NetworkDeviceContext {
    private String deviceName;
    private int x, y;
    private String status;
    private String ipAddress;
    private List<String> connections;
    
    public NetworkDeviceContext(String deviceName, int x, int y, String status) {
        this.deviceName = deviceName;
        this.x = x;
        this.y = y;
        this.status = status;
        this.connections = new ArrayList<>();
    }
    
    // Getters and setters for extrinsic data
    public String getDeviceName() { return deviceName; }
    public int getX() { return x; }
    public int getY() { return y; }
    public String getStatus() { return status; }
}'''
add_code_slide(prs, "Flyweight Pattern - Implementation", flyweight_implementation_code)

# Slide 62: Flyweight Pattern Code - Factory and Usage
flyweight_factory_code = '''// Flyweight Factory
public class NetworkDeviceFlyweightFactory {
    private static final Map<String, NetworkDeviceFlyweight> flyweights = new HashMap<>();
    
    public static NetworkDeviceFlyweight getFlyweight(String deviceType) {
        NetworkDeviceFlyweight flyweight = flyweights.get(deviceType);
        
        if (flyweight == null) {
            // Create new flyweight only if doesn't exist
            switch (deviceType) {
                case "CISCO_2900":
                    flyweight = new RouterFlyweight("CISCO_2900", "Cisco");
                    break;
                case "HUAWEI_AR":
                    flyweight = new RouterFlyweight("HUAWEI_AR", "Huawei");
                    break;
                case "JUNIPER_SRX":
                    flyweight = new RouterFlyweight("JUNIPER_SRX", "Juniper");
                    break;
                default:
                    flyweight = new RouterFlyweight("GENERIC", "Unknown");
            }
            flyweights.put(deviceType, flyweight);
        }
        
        return flyweight;
    }
    
    public static int getFlyweightCount() {
        return flyweights.size();
    }
}

// Usage - Network Visualization
public class NetworkVisualization {
    private List<NetworkDeviceContext> devices;
    
    public void addDevice(String name, String type, int x, int y, String status) {
        NetworkDeviceContext context = new NetworkDeviceContext(name, x, y, status);
        devices.add(context);
        // No need to store flyweight - factory manages sharing
    }
    
    public void renderNetwork(Graphics2D g2d) {
        for (NetworkDeviceContext context : devices) {
            // Get appropriate flyweight for device type
            String deviceType = determineDeviceType(context);
            NetworkDeviceFlyweight flyweight = NetworkDeviceFlyweightFactory.getFlyweight(deviceType);
            
            // Render using shared flyweight with device-specific context
            flyweight.render(context, g2d);
        }
        
        System.out.println("Rendered " + devices.size() + " devices using " 
            + NetworkDeviceFlyweightFactory.getFlyweightCount() + " flyweights");
    }
}

// Memory Comparison:
// Without Flyweight: 50,000 devices × 500KB (icon + specs) = 25GB
// With Flyweight: 10 types × 500KB + 50,000 × 100B (context) = 10MB
// Memory savings: 99.96%!'''
notes_text = "Trainer-Notizen:\n- Drastische Speicher-Einsparungen demonstrieren\n- Intrinsic vs Extrinsic Data klar trennen\n- Factory Pattern Integration zeigen\n- Performance-Verbesserungen betonen\n- Zeitrahmen: ca. 12 Minuten"
add_code_slide(prs, "Flyweight Pattern - Factory & Usage", flyweight_factory_code, "java", notes_text)

# Slide 63: Proxy Pattern - Problem
proxy_problem = [
    "Problem: Kontrollierter oder aufgeschobener Objektzugriff",
    "",
    "Telekom-Szenarien für Proxy:",
    "• Schwere Service-Objekte erst bei Bedarf laden",
    "• Remote-Services über Netzwerk transparent nutzen",
    "• Zugriffskontrolle auf sensitive Daten",
    "• Caching für teure Operationen",
    "",
    "Herausforderungen ohne Proxy:",
    "• Direkte Erstellung aller Objekte (auch ungenutzte)",
    "• Keine transparente Remote-Kommunikation", 
    "• Fehlende Sicherheitskontrolle",
    "• Keine Optimierungen (Caching, Lazy Loading)"
]
add_content_slide(prs, "Proxy Pattern - Problem", proxy_problem)

# Slide 64: Proxy Pattern - Types
proxy_types = [
    "Proxy Pattern - Verschiedene Typen:",
    "",
    "Virtual Proxy (Smart Reference):",
    "• Lazy Loading schwerer Objekte",
    "• Beispiel: Große Konfigurationsdateien laden",
    "",
    "Remote Proxy:",
    "• Repräsentiert Objekte in anderem Adressraum",
    "• Beispiel: REST-API Client für externe Services",
    "",
    "Protection Proxy:",
    "• Zugriffskontrolle und Berechtigungsprüfung", 
    "• Beispiel: Admin-Services mit Rollenkontrolle",
    "",
    "Cache Proxy:",
    "• Caching teurer Operationen",
    "• Beispiel: Datenbankabfragen zwischenspeichern"
]
add_content_slide(prs, "Proxy Pattern - Typen", proxy_types)

# Slide 65: Proxy Pattern Code - Virtual Proxy
proxy_virtual_code = '''// Subject Interface
public interface TelecomServiceConfiguration {
    String getConfigValue(String key);
    Map<String, String> getAllConfigurations();
    boolean updateConfiguration(String key, String value);
    void reloadConfiguration();
}

// Real Subject - Heavy to load
public class ServiceConfigurationManager implements TelecomServiceConfiguration {
    private Map<String, String> configurations;
    private boolean loaded = false;
    
    public ServiceConfigurationManager() {
        System.out.println("ServiceConfigurationManager created");
        // Note: Heavy loading happens in loadConfigurations()
    }
    
    private void loadConfigurations() {
        if (!loaded) {
            System.out.println("Loading service configurations... (heavy operation)");
            configurations = new HashMap<>();
            
            // Simulate expensive loading
            try {
                Thread.sleep(2000); // Database queries, file parsing
                
                // Load from multiple sources
                loadFromDatabase();
                loadFromConfigFiles();
                loadFromExternalAPIs();
                
                loaded = true;
                System.out.println("Configuration loading complete");
            } catch (InterruptedException e) {
                Thread.currentThread().interrupt();
            }
        }
    }
    
    @Override
    public String getConfigValue(String key) {
        loadConfigurations(); // Lazy loading
        return configurations.get(key);
    }
    
    @Override
    public Map<String, String> getAllConfigurations() {
        loadConfigurations();
        return new HashMap<>(configurations);
    }
    
    @Override
    public boolean updateConfiguration(String key, String value) {
        loadConfigurations();
        configurations.put(key, value);
        return persistConfiguration(key, value);
    }
    
    private void loadFromDatabase() { /* ... */ }
    private void loadFromConfigFiles() { /* ... */ }
    private void loadFromExternalAPIs() { /* ... */ }
}'''
add_code_slide(prs, "Proxy Pattern - Virtual Proxy", proxy_virtual_code)

# Slide 66: Proxy Pattern Code - Virtual Proxy Implementation
proxy_virtual_implementation = '''// Virtual Proxy - Controls access to real subject
public class ServiceConfigurationProxy implements TelecomServiceConfiguration {
    private ServiceConfigurationManager realSubject;
    private final String serviceId;
    
    public ServiceConfigurationProxy(String serviceId) {
        this.serviceId = serviceId;
        System.out.println("ServiceConfigurationProxy created for: " + serviceId);
        // Real subject is NOT created here - lazy initialization
    }
    
    private ServiceConfigurationManager getRealSubject() {
        if (realSubject == null) {
            System.out.println("First access - creating real configuration manager");
            realSubject = new ServiceConfigurationManager();
        }
        return realSubject;
    }
    
    @Override
    public String getConfigValue(String key) {
        // Additional proxy logic
        if (key == null || key.trim().isEmpty()) {
            return null;
        }
        
        System.out.println("Proxy: Getting config value for key: " + key);
        return getRealSubject().getConfigValue(key);
    }
    
    @Override
    public Map<String, String> getAllConfigurations() {
        System.out.println("Proxy: Getting all configurations");
        return getRealSubject().getAllConfigurations();
    }
    
    @Override
    public boolean updateConfiguration(String key, String value) {
        // Proxy validation
        if (!isValidConfiguration(key, value)) {
            System.out.println("Proxy: Invalid configuration rejected");
            return false;
        }
        
        System.out.println("Proxy: Updating configuration");
        return getRealSubject().updateConfiguration(key, value);
    }
    
    @Override
    public void reloadConfiguration() {
        if (realSubject != null) {
            realSubject.reloadConfiguration();
        }
    }
    
    private boolean isValidConfiguration(String key, String value) {
        // Validation logic
        return key != null && !key.trim().isEmpty() && value != null;
    }
}

// Usage demonstrates lazy loading
public class ProxyDemo {
    public static void main(String[] args) {
        System.out.println("Creating proxy...");
        TelecomServiceConfiguration config = new ServiceConfigurationProxy("DSL_SERVICE");
        
        System.out.println("\\nProxy created, no heavy loading yet!");
        System.out.println("\\nFirst access triggers loading:");
        
        String value = config.getConfigValue("max_bandwidth");
        System.out.println("Retrieved value: " + value);
    }
}'''
notes_text = "Trainer-Notizen:\n- Lazy Loading Verhalten demonstrieren\n- Zusätzliche Proxy-Funktionalität zeigen\n- Performance-Vorteile diskutieren\n- Telekom-Kontext: Configuration Management\n- Zeitrahmen: ca. 8 Minuten"
add_code_slide(prs, "Proxy Pattern - Virtual Proxy Implementation", proxy_virtual_implementation, "java", notes_text)

# Slide 67: Proxy Pattern Code - Protection Proxy
proxy_protection_code = '''// Protection Proxy - Access Control
public class SecureServiceConfigurationProxy implements TelecomServiceConfiguration {
    private ServiceConfigurationManager realSubject;
    private final UserRole userRole;
    private final String userId;
    
    public SecureServiceConfigurationProxy(UserRole userRole, String userId) {
        this.userRole = userRole;
        this.userId = userId;
    }
    
    private boolean hasReadPermission() {
        return userRole == UserRole.ADMIN || 
               userRole == UserRole.OPERATOR || 
               userRole == UserRole.VIEWER;
    }
    
    private boolean hasWritePermission() {
        return userRole == UserRole.ADMIN || 
               userRole == UserRole.OPERATOR;
    }
    
    private ServiceConfigurationManager getRealSubject() {
        if (realSubject == null) {
            realSubject = new ServiceConfigurationManager();
        }
        return realSubject;
    }
    
    @Override
    public String getConfigValue(String key) {
        if (!hasReadPermission()) {
            logSecurityViolation("READ", key);
            throw new SecurityException("Access denied: Insufficient permissions");
        }
        
        logAccess("READ", key);
        return getRealSubject().getConfigValue(key);
    }
    
    @Override
    public Map<String, String> getAllConfigurations() {
        if (!hasReadPermission()) {
            logSecurityViolation("READ_ALL", "all_configs");
            throw new SecurityException("Access denied: Cannot read configurations");
        }
        
        Map<String, String> configs = getRealSubject().getAllConfigurations();
        
        // Filter sensitive configurations based on role
        if (userRole != UserRole.ADMIN) {
            return filterSensitiveConfigs(configs);
        }
        
        return configs;
    }
    
    @Override
    public boolean updateConfiguration(String key, String value) {
        if (!hasWritePermission()) {
            logSecurityViolation("WRITE", key);
            throw new SecurityException("Access denied: Cannot update configurations");
        }
        
        // Additional validation for critical settings
        if (isCriticalSetting(key) && userRole != UserRole.ADMIN) {
            logSecurityViolation("WRITE_CRITICAL", key);
            throw new SecurityException("Access denied: Admin required for critical settings");
        }
        
        logAccess("WRITE", key);
        return getRealSubject().updateConfiguration(key, value);
    }
    
    private void logAccess(String operation, String key) {
        System.out.println("ACCESS LOG: User " + userId + " (" + userRole + 
                          ") performed " + operation + " on " + key);
    }
    
    private void logSecurityViolation(String operation, String key) {
        System.out.println("SECURITY VIOLATION: User " + userId + " (" + userRole + 
                          ") attempted unauthorized " + operation + " on " + key);
    }
}'''
add_code_slide(prs, "Proxy Pattern - Protection Proxy", proxy_protection_code)

# Slide 68: Structural Patterns Comparison Matrix
patterns_comparison_matrix = [
    "Structural Patterns - Vergleichsmatrix:",
    "",
    "Adapter:",
    "• Zweck: Interface-Kompatibilität | Wann: Legacy-Integration",
    "",
    "Bridge:",
    "• Zweck: Abstraktion/Implementation trennen | Wann: Mehrere Dimensionen",
    "",
    "Composite:",
    "• Zweck: Hierarchische Strukturen | Wann: Teil-Ganzes-Beziehungen",
    "",
    "Decorator:",
    "• Zweck: Verhalten dynamisch erweitern | Wann: Flexible Erweiterungen",
    "",
    "Facade:",
    "• Zweck: Komplexität verstecken | Wann: Vereinfachte Schnittstelle",
    "",
    "Flyweight:",
    "• Zweck: Memory-Optimierung | Wann: Viele ähnliche Objekte",
    "",
    "Proxy:",
    "• Zweck: Kontrollierter Zugriff | Wann: Lazy Loading, Security, Caching"
]
notes_text = "Trainer-Notizen:\n- Entscheidungshilfe für Pattern-Auswahl\n- Überschneidungen und Unterschiede diskutieren\n- Kombinationsmöglichkeiten erwähnen\n- Telekom-Kontext: Welche Patterns wo einsetzen?\n- Zeitrahmen: ca. 8 Minuten"
add_content_slide(prs, "Structural Patterns - Comparison Matrix", patterns_comparison_matrix, notes_text)

# Slide 69: Structural Patterns - Real-world Telekom Use Cases
structural_telekom_cases = [
    "Structural Patterns - Telekom Use Cases:",
    "",
    "Adapter Pattern:",
    "• Legacy Billing-System in Cloud-Architektur integrieren",
    "• Verschiedene Protokolle (SOAP ↔ REST) adaptieren",
    "",
    "Bridge Pattern:",
    "• Notification-Services (Email/SMS) mit verschiedenen Providern",
    "• UI-Framework unabhängig von Rendering-Engine",
    "",
    "Composite Pattern:",
    "• Kunden-Organisationen (Einzelkunden, Firmen, Konzerne)",
    "• Service-Pakete und Bundle-Strukturen",
    "",
    "Decorator Pattern:",
    "• Service-Add-ons (IPv6, QoS, Monitoring) flexibel kombinieren",
    "• Security-Layer für verschiedene Authentifizierungsverfahren"
]
add_content_slide(prs, "Structural Patterns - Telekom Use Cases (1/2)", structural_telekom_cases)

# Slide 70: Structural Patterns - More Telekom Use Cases
more_telekom_cases = [
    "Structural Patterns - Telekom Use Cases (2/2):",
    "",
    "Facade Pattern:",
    "• Kunden-Onboarding Orchestrierung",
    "• Service-Activation Workflow vereinfachen",
    "• API Gateway für Microservice-Architektur",
    "",
    "Flyweight Pattern:",
    "• Netzwerk-Topologie Visualisierung",
    "• Configuration Templates für ähnliche Services",
    "• Routing-Table Optimierung",
    "",
    "Proxy Pattern:",
    "• Service Configuration Lazy Loading",
    "• Remote API-Clients mit Caching",
    "• Admin-Panel mit Rollen-basierter Zugriffskontrolle",
    "• Performance-Monitoring mit intelligenter Pufferung"
]
notes_text = "Trainer-Notizen:\n- Konkrete Anwendungsfälle in Telekom-Umgebung\n- Pattern-Kombinationen in realen Systemen\n- Diskussion: Welche Erfahrungen haben Teilnehmer?\n- Überleitung zu Implementierungs-Workshop\n- Zeitrahmen: ca. 6 Minuten"
add_content_slide(prs, "Structural Patterns - Telekom Use Cases (2/2)", more_telekom_cases, notes_text)

# ============================================================================
# END OF STRUCTURAL PATTERNS SECTION
# ============================================================================

# ============================================================================
# BEHAVIORAL PATTERNS SECTION (SLIDES 71-105)
# ============================================================================

# Slide 71: Behavioral Patterns Section Divider
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = find_placeholder(slide, 1)

title.text = "3. Behavioral Patterns"
for paragraph in title.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.name = TITLE_FONT
        run.font.size = TITLE_SIZE
        run.font.bold = False
        run.font.color.rgb = RGBColor(68, 68, 68)

if subtitle:
    subtitle.text = "Verhalten und Interaktion zwischen Objekten\nKommunikations-Patterns\nAlgorithmus-Flexibilität"
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = CONTENT_FONT
            run.font.size = Pt(22)
            run.font.bold = False
            run.font.color.rgb = RGBColor(68, 68, 68)

# Slide 72: Behavioral Patterns Overview
behavioral_overview = [
    "• Chain of Responsibility - Verkettung von Verarbeitungsschritten",
    "• Command - Kapselung von Operationen als Objekte",
    "• Iterator - Sequenzielle Navigation durch Sammlungen",
    "• Mediator - Zentrale Vermittlung zwischen Komponenten",
    "• Memento - Zustandsspeicherung und -wiederherstellung",
    "• Observer - Event-basierte Benachrichtigungen",
    "• State - Zustandsabhängiges Verhalten",
    "• Strategy - Austauschbare Algorithmen",
    "• Template Method - Gemeinsame Algorithmus-Struktur",
    "• Visitor - Operation auf verschiedenen Objekttypen"
]
notes_text = "Trainer-Notizen:\n- Behavioral Patterns fokussieren auf Kommunikation\n- Entkopplung von Sendern und Empfängern\n- Flexibilität in Algorithmen und Verhalten\n- Besonders wichtig für Event-Driven Architecture\n- Telekom: Workflow-Management, Benachrichtigungen\n- Zeitrahmen: 3 Minuten Überblick"
add_content_slide(prs, "Behavioral Patterns - Übersicht", behavioral_overview, notes_text)

# Slide 73: Chain of Responsibility - Problem
chain_problem = [
    "• Problem: Verschiedene Handler für eine Anfrage",
    "• Beispiel: Support-Ticket-System",
    "  - Level 1: Basic Support (einfache Fragen)",
    "  - Level 2: Technical Support (technische Probleme)", 
    "  - Level 3: Expert Support (komplexe Architektur)",
    "• Herausforderungen:",
    "  - Sender kennt nicht den richtigen Handler",
    "  - Statische Zuweisung ist unflexibel",
    "  - Code ist gekoppelt an Handler-Hierarchie"
]
notes_text = "Trainer-Notizen:\n- Reales Telekom-Beispiel: Support-Eskalation\n- Auch: Firewalls, Genehmigungsworkflows\n- Problem der starren Handler-Zuordnung\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Chain of Responsibility - Problem", chain_problem, notes_text)

# Slide 74: Chain of Responsibility - Structure
chain_structure = [
    "• Handler Interface:",
    "  - handleRequest(request): boolean",
    "  - setNext(handler): Handler",
    "• Concrete Handlers:",
    "  - Implementieren spezifische Logik",
    "  - Entscheiden: bearbeiten oder weiterleiten",
    "• Client:",
    "  - Sendet Request an ersten Handler",
    "  - Kennt nicht die Handler-Kette",
    "• Chain Setup:",
    "  - Handler werden verkettet",
    "  - Dynamische Konfiguration möglich"
]
notes_text = "Trainer-Notizen:\n- UML-Diagramm zeigen\n- Handler-Interface ist zentral\n- setNext() für Verkettung\n- Jeder Handler entscheidet selbst\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Chain of Responsibility - Struktur", chain_structure, notes_text)

# Slide 75: Chain of Responsibility - Implementation
chain_implementation = [
    "```java",
    "// Handler Interface",
    "interface SupportHandler {",
    "    void setNext(SupportHandler handler);",
    "    void handleTicket(SupportTicket ticket);",
    "}",
    "",
    "// Base Handler",
    "abstract class AbstractSupportHandler implements SupportHandler {",
    "    private SupportHandler nextHandler;",
    "    ",
    "    public void setNext(SupportHandler handler) {",
    "        this.nextHandler = handler;",
    "    }",
    "    ",
    "    protected void forwardToNext(SupportTicket ticket) {",
    "        if (nextHandler != null) {",
    "            nextHandler.handleTicket(ticket);",
    "        }",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Abstract Base Handler Pattern\n- Gemeinsame Verkettungslogik\n- Template für Concrete Handler\n- forwardToNext() Utility-Methode\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Chain of Responsibility - Implementation", chain_implementation, notes_text)

# Slide 76: Chain of Responsibility - Concrete Handlers
chain_concrete = [
    "```java",
    "class BasicSupportHandler extends AbstractSupportHandler {",
    "    public void handleTicket(SupportTicket ticket) {",
    "        if (ticket.getComplexity() == Complexity.BASIC) {",
    "            System.out.println(\"Basic Support: \" + ticket.getIssue());",
    "            ticket.resolve(\"Standard solution applied\");",
    "        } else {",
    "            forwardToNext(ticket);",
    "        }",
    "    }",
    "}",
    "",
    "class TechnicalSupportHandler extends AbstractSupportHandler {",
    "    public void handleTicket(SupportTicket ticket) {",
    "        if (ticket.getComplexity() == Complexity.TECHNICAL) {",
    "            System.out.println(\"Technical Support: \" + ticket.getIssue());",
    "            ticket.resolve(\"Technical analysis completed\");",
    "        } else {",
    "            forwardToNext(ticket);",
    "        }",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Konkrete Handler-Implementierung\n- Responsibility-Check am Anfang\n- Bei Nicht-Zuständigkeit: weiterleiten\n- Telekom: Network, Service, Expert Level\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Chain of Responsibility - Concrete Handlers", chain_concrete, notes_text)

# Slide 77: Command Pattern - Problem
command_problem = [
    "• Problem: Operationen als Objekte behandeln",
    "• Use Cases:",
    "  - Undo/Redo Funktionalität",
    "  - Makro-Aufzeichnung",
    "  - Queuing und Scheduling",
    "  - Transaction-ähnliches Verhalten",
    "• Herausforderungen:",
    "  - Direkte Methodenaufrufe sind nicht umkehrbar",
    "  - Operationen können nicht gespeichert werden",
    "  - Keine parametrisierbare Ausführung"
]
notes_text = "Trainer-Notizen:\n- Editor-Beispiel mit Undo/Redo\n- Telekom: Network Configuration Commands\n- Job-Scheduling in Batch-Systemen\n- Wichtig für Event Sourcing\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Command Pattern - Problem", command_problem, notes_text)

# Slide 78: Command Pattern - Structure
command_structure = [
    "• Command Interface:",
    "  - execute(): void",
    "  - undo(): void (optional)",
    "• Concrete Commands:",
    "  - Implementieren spezifische Operationen",
    "  - Speichern Receiver und Parameter",
    "• Receiver:",
    "  - Führt die eigentliche Arbeit aus",
    "• Invoker:",
    "  - Ruft execute() auf Commands auf",
    "  - Kann Command-History verwalten",
    "• Client:",
    "  - Erstellt Commands und setzt Receiver"
]
notes_text = "Trainer-Notizen:\n- Klare Trennung von Aufruf und Ausführung\n- Command speichert alle benötigten Daten\n- Invoker kennt nur das Command Interface\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Command Pattern - Struktur", command_structure, notes_text)

# Slide 79: Command Pattern - Implementation
command_implementation = [
    "```java",
    "// Command Interface",
    "interface NetworkCommand {",
    "    void execute();",
    "    void undo();",
    "}",
    "",
    "// Concrete Command",
    "class ConfigurePortCommand implements NetworkCommand {",
    "    private NetworkSwitch receiver;",
    "    private int port;",
    "    private PortConfig newConfig;",
    "    private PortConfig oldConfig;",
    "    ",
    "    public ConfigurePortCommand(NetworkSwitch receiver, int port, PortConfig config) {",
    "        this.receiver = receiver;",
    "        this.port = port;",
    "        this.newConfig = config;",
    "    }",
    "    ",
    "    public void execute() {",
    "        oldConfig = receiver.getPortConfig(port);",
    "        receiver.configurePort(port, newConfig);",
    "    }",
    "    ",
    "    public void undo() {",
    "        receiver.configurePort(port, oldConfig);",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Network Configuration Beispiel\n- Command speichert alten Zustand für Undo\n- Receiver macht die eigentliche Arbeit\n- Parameter werden im Command gespeichert\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Command Pattern - Implementation", command_implementation, notes_text)

# Slide 80: Command Pattern - Invoker & Macro
command_macro = [
    "```java",
    "class NetworkConfigurationManager {",
    "    private Stack<NetworkCommand> history = new Stack<>();",
    "    ",
    "    public void execute(NetworkCommand command) {",
    "        command.execute();",
    "        history.push(command);",
    "    }",
    "    ",
    "    public void undo() {",
    "        if (!history.isEmpty()) {",
    "            NetworkCommand command = history.pop();",
    "            command.undo();",
    "        }",
    "    }",
    "}",
    "",
    "// Macro Command",
    "class NetworkSetupMacro implements NetworkCommand {",
    "    private List<NetworkCommand> commands;",
    "    ",
    "    public void execute() {",
    "        commands.forEach(NetworkCommand::execute);",
    "    }",
    "    ",
    "    public void undo() {",
    "        Collections.reverse(commands);",
    "        commands.forEach(NetworkCommand::undo);",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Invoker mit Command History\n- Undo-Stack Implementation\n- Macro Command für Batch-Operationen\n- Reihenfolge beim Undo beachten\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Command Pattern - Invoker & Macro", command_macro, notes_text)

# Slide 81: Iterator Pattern - Problem
iterator_problem = [
    "• Problem: Navigation durch komplexe Datenstrukturen",
    "• Verschiedene Traversal-Arten:",
    "  - Sequential, Tree (Pre/In/Post-Order)",
    "  - Filtered, Sorted Traversal",
    "• Herausforderungen:",
    "  - Collection-interne Struktur exposiert",
    "  - Client-Code abhängig von Datenstruktur",
    "  - Mehrere gleichzeitige Iterationen schwierig",
    "• Beispiel: Telekom Service-Hierarchie",
    "  - Services → SubServices → Components"
]
notes_text = "Trainer-Notizen:\n- Organisationsstruktur als Beispiel\n- Problem der Exposition interner Struktur\n- Java Collections Framework als Referenz\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Iterator Pattern - Problem", iterator_problem, notes_text)

# Slide 82: Iterator Pattern - Structure
iterator_structure = [
    "• Iterator Interface:",
    "  - hasNext(): boolean",
    "  - next(): T",
    "  - remove(): void (optional)",
    "• Concrete Iterator:",
    "  - Implementiert Traversal-Logik",
    "  - Hält aktuellen Zustand",
    "• Iterable Interface:",
    "  - iterator(): Iterator<T>",
    "• Concrete Collection:",
    "  - Implementiert Iterable",
    "  - Erstellt passende Iterator-Instanzen"
]
notes_text = "Trainer-Notizen:\n- Standard Iterator Interface Pattern\n- Iterator hält eigenen Zustand\n- Collection bleibt unverändert\n- Mehrere Iterator gleichzeitig möglich\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Iterator Pattern - Struktur", iterator_structure, notes_text)

# Slide 83: Iterator Pattern - Implementation
iterator_implementation = [
    "```java",
    "// Service Hierarchy",
    "class ServiceComponent {",
    "    private String name;",
    "    private List<ServiceComponent> children = new ArrayList<>();",
    "    ",
    "    // getters, setters, add methods...",
    "}",
    "",
    "// Iterator Interface",
    "interface ServiceIterator {",
    "    boolean hasNext();",
    "    ServiceComponent next();",
    "}",
    "",
    "// Depth-First Iterator",
    "class DepthFirstServiceIterator implements ServiceIterator {",
    "    private Stack<ServiceComponent> stack = new Stack<>();",
    "    ",
    "    public DepthFirstServiceIterator(ServiceComponent root) {",
    "        if (root != null) stack.push(root);",
    "    }",
    "    ",
    "    public boolean hasNext() {",
    "        return !stack.isEmpty();",
    "    }",
    "    ",
    "    public ServiceComponent next() {",
    "        ServiceComponent current = stack.pop();",
    "        // Add children in reverse order for correct traversal",
    "        for (int i = current.getChildren().size() - 1; i >= 0; i--) {",
    "            stack.push(current.getChildren().get(i));",
    "        }",
    "        return current;",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Tree-Traversal Beispiel\n- Stack-basierte Depth-First Implementation\n- Reihenfolge der Children wichtig\n- Iterator hält eigenen Zustand\n- Zeitrahmen: 7 Minuten"
add_content_slide(prs, "Iterator Pattern - Implementation", iterator_implementation, notes_text)

# Slide 84: Iterator Pattern - Usage & Benefits
iterator_usage = [
    "```java",
    "// Usage Example",
    "ServiceComponent telekomServices = buildServiceHierarchy();",
    "",
    "// Depth-First Traversal",
    "ServiceIterator iterator = new DepthFirstServiceIterator(telekomServices);",
    "while (iterator.hasNext()) {",
    "    ServiceComponent service = iterator.next();",
    "    System.out.println(\"Processing: \" + service.getName());",
    "}",
    "",
    "// Level-Order Iterator (alternative implementation)",
    "ServiceIterator levelIterator = new LevelOrderServiceIterator(telekomServices);",
    "",
    "// Benefits:",
    "• Collection-interne Struktur bleibt verborgen",
    "• Verschiedene Traversal-Algorithmen möglich",
    "• Mehrere gleichzeitige Iterationen",
    "• Einheitliche Navigation-Schnittstelle"
]
notes_text = "Trainer-Notizen:\n- Verschiedene Iterator-Implementierungen\n- Collection bleibt unverändert\n- Client-Code unabhängig von Struktur\n- Java für-each Loop nutzt Iterator\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Iterator Pattern - Usage & Benefits", iterator_usage, notes_text)

# Slide 85: Mediator Pattern - Problem
mediator_problem = [
    "• Problem: Komplexe Kommunikation zwischen Komponenten",
    "• Beispiel: Telekom Network Management System",
    "  - Router, Switches, Monitoring, Logging",
    "  - Jede Komponente kommuniziert mit mehreren anderen",
    "• Herausforderungen:",
    "  - Tight Coupling zwischen Komponenten",
    "  - Schwer testbar und wartbar",
    "  - Änderungen propagieren durch das System",
    "  - Komplexe Abhängigkeitsgraphen"
]
notes_text = "Trainer-Notizen:\n- Network Operations Center als Beispiel\n- Problem der Punkt-zu-Punkt Kommunikation\n- Spaghetti-Code durch direkte Kopplungen\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Mediator Pattern - Problem", mediator_problem, notes_text)

# Slide 86: Mediator Pattern - Structure
mediator_structure = [
    "• Mediator Interface:",
    "  - notify(sender, event): void",
    "• Concrete Mediator:",
    "  - Orchestriert Kommunikation",
    "  - Kennt alle beteiligten Komponenten",
    "• Colleague Interface:",
    "  - setMediator(mediator): void",
    "• Concrete Colleagues:",
    "  - Kommunizieren nur über Mediator",
    "  - Senden Events an Mediator",
    "• Vorteile:",
    "  - Loose Coupling zwischen Colleagues",
    "  - Zentrale Koordination"
]
notes_text = "Trainer-Notizen:\n- Hub and Spoke Architecture\n- Mediator ist der zentrale Koordinator\n- Colleagues kennen sich nicht direkt\n- Kommunikation nur über Events\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Mediator Pattern - Struktur", mediator_structure, notes_text)

# Slide 87: Mediator Pattern - Implementation
mediator_implementation = [
    "```java",
    "// Mediator Interface",
    "interface NetworkMediator {",
    "    void notify(NetworkComponent sender, String event, Object data);",
    "}",
    "",
    "// Network Component Base",
    "abstract class NetworkComponent {",
    "    protected NetworkMediator mediator;",
    "    ",
    "    public NetworkComponent(NetworkMediator mediator) {",
    "        this.mediator = mediator;",
    "    }",
    "}",
    "",
    "// Concrete Component",
    "class Router extends NetworkComponent {",
    "    public Router(NetworkMediator mediator) {",
    "        super(mediator);",
    "    }",
    "    ",
    "    public void routePacket(Packet packet) {",
    "        // Route the packet",
    "        if (packet.isHighPriority()) {",
    "            mediator.notify(this, \"HIGH_PRIORITY_TRAFFIC\", packet);",
    "        }",
    "    }",
    "    ",
    "    public void handleNetworkCongestion() {",
    "        System.out.println(\"Router: Adjusting routing table\");",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Network Management Beispiel\n- Components kennen nur den Mediator\n- Events werden an Mediator weitergeleitet\n- Mediator orchestriert Reaktionen\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Mediator Pattern - Implementation", mediator_implementation, notes_text)

# Slide 88: Mediator Pattern - Concrete Mediator
mediator_concrete = [
    "```java",
    "class NetworkOperationsCenter implements NetworkMediator {",
    "    private Router router;",
    "    private Switch networkSwitch;",
    "    private MonitoringSystem monitor;",
    "    private AlertSystem alerts;",
    "    ",
    "    public void notify(NetworkComponent sender, String event, Object data) {",
    "        switch (event) {",
    "            case \"HIGH_PRIORITY_TRAFFIC\":",
    "                handleHighPriorityTraffic((Packet) data);",
    "                break;",
    "            case \"NETWORK_CONGESTION\":",
    "                handleNetworkCongestion(sender);",
    "                break;",
    "            case \"DEVICE_FAILURE\":",
    "                handleDeviceFailure(sender);",
    "                break;",
    "        }",
    "    }",
    "    ",
    "    private void handleHighPriorityTraffic(Packet packet) {",
    "        networkSwitch.prioritizeTraffic(packet);",
    "        monitor.logPriorityTraffic(packet);",
    "        if (packet.isCritical()) {",
    "            alerts.sendAlert(\"Critical packet routed\");",
    "        }",
    "    }",
    "}",
    "// Registration methods for components..."
]
notes_text = "Trainer-Notizen:\n- Zentrale Orchestrierung aller Events\n- Event-basierte Kommunikation\n- NOC als reales Mediator-Beispiel\n- Koordination verschiedener Systeme\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Mediator Pattern - Concrete Mediator", mediator_concrete, notes_text)

# Slide 89: Memento Pattern - Problem
memento_problem = [
    "• Problem: Zustandsspeicherung ohne Kapselung zu verletzen",
    "• Use Cases:",
    "  - Undo/Redo Funktionalität",
    "  - Checkpoints und Snapshots",
    "  - Transaction Rollback",
    "  - Configuration Backup",
    "• Herausforderungen:",
    "  - Private Daten müssen zugänglich sein",
    "  - Performance bei großen Zuständen",
    "  - Memory Management",
    "• Beispiel: Network Device Configuration",
    "  - Backup vor Änderungen",
    "  - Rollback bei Fehlern"
]
notes_text = "Trainer-Notizen:\n- Configuration Management als Beispiel\n- Problem der Kapselungsverletzung\n- Snapshot-Funktionalität in Netzwerken\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Memento Pattern - Problem", memento_problem, notes_text)

# Slide 90: Memento Pattern - Structure
memento_structure = [
    "• Originator:",
    "  - Erstellt Memento von aktuellem Zustand",
    "  - Stellt Zustand aus Memento wieder her",
    "• Memento:",
    "  - Speichert Snapshot des Originator-Zustands",
    "  - Immutable und opaque für anderen Code",
    "• Caretaker:",
    "  - Verwaltet Mementos",
    "  - Löst Backup/Restore aus",
    "  - Kennt Memento-Inhalt nicht",
    "• Kapselierung:",
    "  - Nur Originator kann Memento-Inhalt lesen",
    "  - Caretaker behandelt Memento als Black Box"
]
notes_text = "Trainer-Notizen:\n- Drei-Rollen-Pattern\n- Kapselung bleibt erhalten\n- Originator hat exklusiven Zugriff\n- Caretaker macht Lifecycle-Management\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Memento Pattern - Struktur", memento_structure, notes_text)

# Slide 91: Memento Pattern - Implementation
memento_implementation = [
    "```java",
    "// Network Configuration (Originator)",
    "class NetworkConfiguration {",
    "    private String routingTable;",
    "    private Map<String, String> interfaceConfigs;",
    "    private List<SecurityRule> securityRules;",
    "    ",
    "    // Inner Memento Class - nur Originator hat Zugriff",
    "    public class ConfigurationMemento {",
    "        private final String routingTable;",
    "        private final Map<String, String> interfaceConfigs;",
    "        private final List<SecurityRule> securityRules;",
    "        private final long timestamp;",
    "        ",
    "        private ConfigurationMemento() {",
    "            this.routingTable = NetworkConfiguration.this.routingTable;",
    "            this.interfaceConfigs = new HashMap<>(NetworkConfiguration.this.interfaceConfigs);",
    "            this.securityRules = new ArrayList<>(NetworkConfiguration.this.securityRules);",
    "            this.timestamp = System.currentTimeMillis();",
    "        }",
    "    }",
    "    ",
    "    public ConfigurationMemento createSnapshot() {",
    "        return new ConfigurationMemento();",
    "    }",
    "    ",
    "    public void restore(ConfigurationMemento memento) {",
    "        this.routingTable = memento.routingTable;",
    "        this.interfaceConfigs = new HashMap<>(memento.interfaceConfigs);",
    "        this.securityRules = new ArrayList<>(memento.securityRules);",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Inner Class Pattern für Kapselung\n- Deep Copy der Datenstrukturen\n- Timestamp für Versioning\n- Nur Originator kann Memento erstellen/lesen\n- Zeitrahmen: 7 Minuten"
add_content_slide(prs, "Memento Pattern - Implementation", memento_implementation, notes_text)

# Slide 92: Memento Pattern - Caretaker
memento_caretaker = [
    "```java",
    "// Configuration Manager (Caretaker)",
    "class NetworkConfigurationManager {",
    "    private NetworkConfiguration config;",
    "    private Stack<NetworkConfiguration.ConfigurationMemento> history;",
    "    private final int maxHistorySize = 10;",
    "    ",
    "    public NetworkConfigurationManager(NetworkConfiguration config) {",
    "        this.config = config;",
    "        this.history = new Stack<>();",
    "    }",
    "    ",
    "    public void backup() {",
    "        if (history.size() >= maxHistorySize) {",
    "            // Remove oldest snapshot",
    "            history.remove(0);",
    "        }",
    "        history.push(config.createSnapshot());",
    "        System.out.println(\"Configuration backed up\");",
    "    }",
    "    ",
    "    public void rollback() {",
    "        if (!history.isEmpty()) {",
    "            NetworkConfiguration.ConfigurationMemento memento = history.pop();",
    "            config.restore(memento);",
    "            System.out.println(\"Configuration rolled back\");",
    "        }",
    "    }",
    "    ",
    "    public void performSafeConfigChange(Runnable configChange) {",
    "        backup();  // Create backup before change",
    "        try {",
    "            configChange.run();",
    "        } catch (Exception e) {",
    "            rollback();  // Automatic rollback on error",
    "            throw e;",
    "        }",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Caretaker verwaltet Memento-Lifecycle\n- History-Größe begrenzen für Memory Management\n- Safe Configuration Change Pattern\n- Automatic Rollback bei Fehlern\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Memento Pattern - Caretaker", memento_caretaker, notes_text)

# Slide 93: Observer Pattern - Problem
observer_problem = [
    "• Problem: Benachrichtigung über Zustandsänderungen",
    "• Beispiele:",
    "  - Network Monitoring Dashboard",
    "  - Service Status Updates",
    "  - Configuration Change Notifications",
    "• Herausforderungen:",
    "  - Subject kennt alle Observer",
    "  - Tight Coupling Subject ↔ Observer",
    "  - Synchrone Notification kann langsam sein",
    "  - Memory Leaks durch Observer-Referenzen",
    "• Telekom Use Case:",
    "  - Service Health Monitoring",
    "  - Multiple Dashboards, Alerts, Logs"
]
notes_text = "Trainer-Notizen:\n- Publish-Subscribe Pattern\n- GUI Model-View-Controller verwandt\n- Event-Driven Architecture Foundation\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Observer Pattern - Problem", observer_problem, notes_text)

# Slide 94: Observer Pattern - Structure
observer_structure = [
    "• Subject (Observable):",
    "  - attach(observer): void",
    "  - detach(observer): void", 
    "  - notify(): void",
    "• Observer Interface:",
    "  - update(subject): void",
    "• Concrete Subject:",
    "  - Speichert Zustand",
    "  - Benachrichtigt Observer bei Änderungen",
    "• Concrete Observer:",
    "  - Implementiert Update-Logik",
    "  - Kann Zustand vom Subject abfragen",
    "• Notification Patterns:",
    "  - Push Model: Subject sendet Daten",
    "  - Pull Model: Observer fragt Daten ab"
]
notes_text = "Trainer-Notizen:\n- Classic GoF Observer Pattern\n- Push vs Pull Model erklären\n- Observer List Management\n- Notification Loop Prevention\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Observer Pattern - Struktur", observer_structure, notes_text)

# Slide 95: Observer Pattern - Implementation
observer_implementation = [
    "```java",
    "// Observer Interface",
    "interface ServiceObserver {",
    "    void onServiceStatusChanged(ServiceStatus status, String serviceName);",
    "}",
    "",
    "// Subject (Observable)",
    "class TelekomService {",
    "    private String serviceName;",
    "    private ServiceStatus status;",
    "    private List<ServiceObserver> observers = new ArrayList<>();",
    "    ",
    "    public void addObserver(ServiceObserver observer) {",
    "        observers.add(observer);",
    "    }",
    "    ",
    "    public void removeObserver(ServiceObserver observer) {",
    "        observers.remove(observer);",
    "    }",
    "    ",
    "    public void setStatus(ServiceStatus newStatus) {",
    "        if (this.status != newStatus) {",
    "            this.status = newStatus;",
    "            notifyObservers();",
    "        }",
    "    }",
    "    ",
    "    private void notifyObservers() {",
    "        // Copy list to prevent ConcurrentModificationException",
    "        List<ServiceObserver> observersCopy = new ArrayList<>(observers);",
    "        for (ServiceObserver observer : observersCopy) {",
    "            observer.onServiceStatusChanged(status, serviceName);",
    "        }",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Service Monitoring Beispiel\n- Observer List kopieren für Thread Safety\n- Status-Change Detection\n- Notification nur bei tatsächlicher Änderung\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Observer Pattern - Implementation", observer_implementation, notes_text)

# Slide 96: Observer Pattern - Concrete Observers
observer_concrete = [
    "```java",
    "// Dashboard Observer",
    "class MonitoringDashboard implements ServiceObserver {",
    "    private String dashboardName;",
    "    ",
    "    public void onServiceStatusChanged(ServiceStatus status, String serviceName) {",
    "        System.out.println(\"[\" + dashboardName + \"] Service \" + serviceName",
    "                         + \" status changed to \" + status);",
    "        updateDashboardDisplay(serviceName, status);",
    "    }",
    "}",
    "",
    "// Alert System Observer",
    "class AlertSystem implements ServiceObserver {",
    "    public void onServiceStatusChanged(ServiceStatus status, String serviceName) {",
    "        if (status == ServiceStatus.DOWN || status == ServiceStatus.ERROR) {",
    "            sendAlert(\"CRITICAL: Service \" + serviceName + \" is \" + status);",
    "            escalateToSupport(serviceName, status);",
    "        }",
    "    }",
    "}",
    "",
    "// Logging Observer",
    "class ServiceLogger implements ServiceObserver {",
    "    public void onServiceStatusChanged(ServiceStatus status, String serviceName) {",
    "        String logEntry = LocalDateTime.now() + \": \" + serviceName",
    "                        + \" status changed to \" + status;",
    "        writeToLog(logEntry);",
    "        if (status == ServiceStatus.DOWN) {",
    "            archiveServiceData(serviceName);",
    "        }",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Verschiedene Observer-Typen\n- Dashboard, Alerting, Logging\n- Unterschiedliche Reaktionen auf Events\n- Loose Coupling zwischen Observern\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Observer Pattern - Concrete Observers", observer_concrete, notes_text)

# Slide 97: State Pattern - Problem
state_problem = [
    "• Problem: Komplexe zustandsabhängige Logik",
    "• Beispiel: Telekom Service Lifecycle",
    "  - Pending → Active → Suspended → Terminated",
    "  - Verschiedene Operationen je nach Zustand",
    "• Herausforderungen:",
    "  - Große if-else oder switch Statements",
    "  - Schwer erweiterbar um neue Zustände",
    "  - Verstreute State-Logic im Code",
    "  - Komplexe Zustandsübergänge",
    "• State Machine Komplexität",
    "  - Viele Zustände × Viele Events = Komplexität²"
]
notes_text = "Trainer-Notizen:\n- Service Lifecycle Management\n- Problem von proceduraler State Logic\n- Finite State Machine Theorie\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "State Pattern - Problem", state_problem, notes_text)

# Slide 98: State Pattern - Structure
state_structure = [
    "• Context:",
    "  - Hält Referenz auf aktuellen State",
    "  - Delegiert State-spezifische Operationen",
    "• State Interface:",
    "  - Definiert State-spezifische Methoden",
    "• Concrete States:",
    "  - Implementieren zustandsabhängiges Verhalten",
    "  - Können State Transitions auslösen",
    "• State Transitions:",
    "  - Context.setState(newState)",
    "  - States können sich selbst ändern",
    "• Vorteile:",
    "  - OCP: Neue States ohne Context-Änderung",
    "  - State-spezifische Logik gekapselt"
]
notes_text = "Trainer-Notizen:\n- Context als State Machine\n- Delegation statt Conditionals\n- States kennen mögliche Transitions\n- Open-Closed Principle anwendbar\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "State Pattern - Struktur", state_structure, notes_text)

# Slide 99: State Pattern - Implementation
state_implementation = [
    "```java",
    "// State Interface",
    "interface ServiceState {",
    "    void activate(ServiceContext context);",
    "    void suspend(ServiceContext context);",
    "    void terminate(ServiceContext context);",
    "    void restart(ServiceContext context);",
    "    String getStatusName();",
    "}",
    "",
    "// Context",
    "class ServiceContext {",
    "    private ServiceState currentState;",
    "    private String serviceName;",
    "    ",
    "    public ServiceContext(String serviceName) {",
    "        this.serviceName = serviceName;",
    "        this.currentState = new PendingState();",
    "    }",
    "    ",
    "    public void setState(ServiceState state) {",
    "        System.out.println(\"Service \" + serviceName + \" transitioning to \" + state.getStatusName());",
    "        this.currentState = state;",
    "    }",
    "    ",
    "    public void activate() { currentState.activate(this); }",
    "    public void suspend() { currentState.suspend(this); }",
    "    public void terminate() { currentState.terminate(this); }",
    "    public void restart() { currentState.restart(this); }",
    "    ",
    "    public String getCurrentStatus() {",
    "        return currentState.getStatusName();",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Context als State Machine\n- Delegation an Current State\n- State Transition Logging\n- Einheitliche Interface für alle Operationen\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "State Pattern - Implementation", state_implementation, notes_text)

# Slide 100: State Pattern - Concrete States
state_concrete = [
    "```java",
    "class PendingState implements ServiceState {",
    "    public void activate(ServiceContext context) {",
    "        System.out.println(\"Activating service from pending state\");",
    "        context.setState(new ActiveState());",
    "    }",
    "    public void suspend(ServiceContext context) {",
    "        throw new IllegalStateException(\"Cannot suspend pending service\");",
    "    }",
    "    public void terminate(ServiceContext context) {",
    "        System.out.println(\"Terminating pending service\");",
    "        context.setState(new TerminatedState());",
    "    }",
    "    public void restart(ServiceContext context) {",
    "        System.out.println(\"Service already in pending state\");",
    "    }",
    "    public String getStatusName() { return \"PENDING\"; }",
    "}",
    "",
    "class ActiveState implements ServiceState {",
    "    public void activate(ServiceContext context) {",
    "        System.out.println(\"Service already active\");",
    "    }",
    "    public void suspend(ServiceContext context) {",
    "        System.out.println(\"Suspending active service\");",
    "        context.setState(new SuspendedState());",
    "    }",
    "    public void terminate(ServiceContext context) {",
    "        System.out.println(\"Terminating active service\");",
    "        context.setState(new TerminatedState());",
    "    }",
    "    public void restart(ServiceContext context) {",
    "        System.out.println(\"Restarting service\");",
    "        context.setState(new PendingState());",
    "        context.activate(); // Chain to activation",
    "    }",
    "    public String getStatusName() { return \"ACTIVE\"; }",
    "}"
]
notes_text = "Trainer-Notizen:\n- State-spezifisches Verhalten\n- Illegale Transitions als Exceptions\n- States können weitere Transitions auslösen\n- Klar definierte State Semantics\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "State Pattern - Concrete States", state_concrete, notes_text)

# Slide 101: Strategy Pattern - Problem
strategy_problem = [
    "• Problem: Verschiedene Algorithmen für gleiche Aufgabe",
    "• Beispiele:",
    "  - Routing-Algorithmen (Shortest Path, Load Balancing)",
    "  - Compression-Algorithmen (ZIP, GZIP, LZ4)",
    "  - Payment Processing (Credit Card, PayPal, Invoice)",
    "• Herausforderungen:",
    "  - Algorithm-Selection in Client-Code",
    "  - Schwer testbar und wartbar",
    "  - Violation of Open-Closed Principle",
    "• Telekom Example:",
    "  - Network Traffic Routing",
    "  - Different strategies based on conditions"
]
notes_text = "Trainer-Notizen:\n- Familie von Algorithmen\n- Runtime Algorithm Selection\n- Telekom: Dynamic Routing Strategies\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Strategy Pattern - Problem", strategy_problem, notes_text)

# Slide 102: Strategy Pattern - Structure
strategy_structure = [
    "• Strategy Interface:",
    "  - execute(data): Result",
    "• Concrete Strategies:",
    "  - Implementieren verschiedene Algorithmen",
    "  - Gleiche Schnittstelle, unterschiedliche Implementierung",
    "• Context:",
    "  - Hält Referenz auf Strategy",
    "  - Delegiert Algorithmus-Ausführung",
    "  - setStrategy(strategy): void",
    "• Client:",
    "  - Wählt passende Strategy",
    "  - Konfiguriert Context",
    "• Vorteile:",
    "  - Runtime Strategy Switching",
    "  - Algorithmus-Familie erweiterbar"
]
notes_text = "Trainer-Notizen:\n- Composition over Inheritance\n- Strategy als pluggable Algorithm\n- Context wird parametrisiert\n- Open-Closed Principle erfüllt\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Strategy Pattern - Struktur", strategy_structure, notes_text)

# Slide 103: Strategy Pattern - Implementation
strategy_implementation = [
    "```java",
    "// Strategy Interface",
    "interface RoutingStrategy {",
    "    Route calculateRoute(NetworkTopology topology, Node source, Node destination);",
    "    String getStrategyName();",
    "}",
    "",
    "// Concrete Strategies",
    "class ShortestPathStrategy implements RoutingStrategy {",
    "    public Route calculateRoute(NetworkTopology topology, Node source, Node destination) {",
    "        // Dijkstra's algorithm implementation",
    "        System.out.println(\"Calculating shortest path route\");",
    "        return dijkstra(topology, source, destination);",
    "    }",
    "    public String getStrategyName() { return \"SHORTEST_PATH\"; }",
    "}",
    "",
    "class LoadBalancingStrategy implements RoutingStrategy {",
    "    public Route calculateRoute(NetworkTopology topology, Node source, Node destination) {",
    "        System.out.println(\"Calculating load-balanced route\");",
    "        return findLeastCongestedPath(topology, source, destination);",
    "    }",
    "    public String getStrategyName() { return \"LOAD_BALANCING\"; }",
    "}",
    "",
    "class HighAvailabilityStrategy implements RoutingStrategy {",
    "    public Route calculateRoute(NetworkTopology topology, Node source, Node destination) {",
    "        System.out.println(\"Calculating high-availability route with redundancy\");",
    "        return findRedundantPath(topology, source, destination);",
    "    }",
    "    public String getStrategyName() { return \"HIGH_AVAILABILITY\"; }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Network Routing Beispiel\n- Verschiedene Routing-Algorithmen\n- Gleiche Schnittstelle, andere Implementierung\n- Telekom: Adaptives Routing je nach Netzlast\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Strategy Pattern - Implementation", strategy_implementation, notes_text)

# Slide 104: Strategy Pattern - Context & Usage
strategy_usage = [
    "```java",
    "// Context",
    "class NetworkRouter {",
    "    private RoutingStrategy strategy;",
    "    private NetworkTopology topology;",
    "    ",
    "    public NetworkRouter(NetworkTopology topology) {",
    "        this.topology = topology;",
    "        this.strategy = new ShortestPathStrategy(); // default",
    "    }",
    "    ",
    "    public void setRoutingStrategy(RoutingStrategy strategy) {",
    "        this.strategy = strategy;",
    "        System.out.println(\"Routing strategy changed to: \" + strategy.getStrategyName());",
    "    }",
    "    ",
    "    public Route routeTraffic(Node source, Node destination) {",
    "        return strategy.calculateRoute(topology, source, destination);",
    "    }",
    "}",
    "",
    "// Usage Example",
    "NetworkRouter router = new NetworkRouter(networkTopology);",
    "",
    "// Normal conditions - use shortest path",
    "Route route1 = router.routeTraffic(nodeA, nodeB);",
    "",
    "// High traffic - switch to load balancing",
    "if (networkLoad > 0.8) {",
    "    router.setRoutingStrategy(new LoadBalancingStrategy());",
    "}",
    "",
    "// Critical service - ensure high availability",
    "if (service.isCritical()) {",
    "    router.setRoutingStrategy(new HighAvailabilityStrategy());",
    "}"
]
notes_text = "Trainer-Notizen:\n- Context mit Strategy-Delegation\n- Runtime Strategy Switching\n- Condition-based Strategy Selection\n- Dynamic Adaptation an Netzwerkbedingungen\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Strategy Pattern - Context & Usage", strategy_usage, notes_text)

# Slide 105: Template Method Pattern - Problem
template_problem = [
    "• Problem: Gemeinsamer Algorithmus mit variablen Schritten",
    "• Beispiele:",
    "  - Data Processing Pipeline",
    "  - Service Deployment Workflow",
    "  - Quality Assurance Process",
    "• Herausforderungen:",
    "  - Code Duplication in ähnlichen Algorithmen",
    "  - Schwer konsistente Ausführungsreihenfolge",
    "  - Varianten schwer zu erweitern",
    "• Telekom Use Case:",
    "  - Service Provisioning Process",
    "  - Gemeinsame Schritte, unterschiedliche Services"
]
notes_text = "Trainer-Notizen:\n- Algorithm Skeleton Pattern\n- Template definiert Struktur\n- Subclasses implementieren Variable\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Template Method Pattern - Problem", template_problem, notes_text)

# Slide 106: Template Method Pattern - Structure
template_structure = [
    "• Abstract Class:",
    "  - templateMethod(): final - Algorithmus-Skelett",
    "  - primitiveOperations(): abstract - Variable Schritte",
    "  - hook(): void - Optionale Erweiterungspunkte",
    "• Concrete Class:",
    "  - Implementiert primitive Operations",
    "  - Überschreibt Hooks bei Bedarf",
    "• Template Method:",
    "  - Definiert Ablaufstruktur",
    "  - Ruft primitive Operations in Reihenfolge auf",
    "  - Meist final (nicht überschreibbar)",
    "• Hollywood Principle:",
    "  - \"Don't call us, we'll call you\"",
    "  - Framework ruft Application Code auf"
]
notes_text = "Trainer-Notizen:\n- Inversion of Control Pattern\n- Template Method ist final\n- Hook Methods für optionale Extensions\n- Hollywood Principle erklären\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Template Method Pattern - Struktur", template_structure, notes_text)

# Slide 107: Template Method Pattern - Implementation
template_implementation = [
    "```java",
    "// Abstract Template",
    "abstract class ServiceProvisioningTemplate {",
    "    ",
    "    // Template Method - defines the algorithm skeleton",
    "    public final ProvisioningResult provisionService(ServiceRequest request) {",
    "        System.out.println(\"Starting service provisioning process...\");",
    "        ",
    "        // Step 1: Validate request",
    "        if (!validateRequest(request)) {",
    "            return ProvisioningResult.failed(\"Request validation failed\");",
    "        }",
    "        ",
    "        // Step 2: Allocate resources (variable implementation)",
    "        ResourceAllocation allocation = allocateResources(request);",
    "        if (allocation == null) {",
    "            return ProvisioningResult.failed(\"Resource allocation failed\");",
    "        }",
    "        ",
    "        // Step 3: Configure service (variable implementation)",
    "        ServiceConfiguration config = configureService(request, allocation);",
    "        ",
    "        // Step 4: Deploy service (variable implementation)",
    "        DeploymentResult deployment = deployService(config);",
    "        ",
    "        // Step 5: Post-deployment hook (optional)",
    "        postDeploymentHook(deployment);",
    "        ",
    "        // Step 6: Finalize provisioning",
    "        return finalizeProvisioning(request, deployment);",
    "    }",
    "    ",
    "    // Common implementation",
    "    private boolean validateRequest(ServiceRequest request) {",
    "        return request != null && request.isValid();",
    "    }",
    "    ",
    "    // Abstract methods - must be implemented by subclasses",
    "    protected abstract ResourceAllocation allocateResources(ServiceRequest request);",
    "    protected abstract ServiceConfiguration configureService(ServiceRequest request, ResourceAllocation allocation);",
    "    protected abstract DeploymentResult deployService(ServiceConfiguration config);",
    "    ",
    "    // Hook method - optional override",
    "    protected void postDeploymentHook(DeploymentResult result) {",
    "        // Default: do nothing, subclasses can override",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Service Provisioning Template\n- Template Method ist final\n- Abstract Methods für variable Schritte\n- Hook Methods für optionale Erweiterungen\n- Zeitrahmen: 7 Minuten"
add_content_slide(prs, "Template Method Pattern - Implementation", template_implementation, notes_text)

# Slide 108: Template Method Pattern - Concrete Implementations
template_concrete = [
    "```java",
    "// Concrete Implementation for VoIP Service",
    "class VoipServiceProvisioning extends ServiceProvisioningTemplate {",
    "    ",
    "    protected ResourceAllocation allocateResources(ServiceRequest request) {",
    "        System.out.println(\"Allocating VoIP resources: SIP trunk, bandwidth\");",
    "        return new ResourceAllocation()",
    "                .addResource(\"sip_trunk\", request.getRequiredTrunks())",
    "                .addResource(\"bandwidth\", request.getBandwidth());",
    "    }",
    "    ",
    "    protected ServiceConfiguration configureService(ServiceRequest request, ResourceAllocation allocation) {",
    "        System.out.println(\"Configuring VoIP service: PBX, routing rules\");",
    "        return new VoipConfiguration()",
    "                .setPbxSettings(request.getPbxConfig())",
    "                .setRoutingRules(request.getRoutingRules())",
    "                .setCodecSettings(request.getPreferredCodecs());",
    "    }",
    "    ",
    "    protected DeploymentResult deployService(ServiceConfiguration config) {",
    "        System.out.println(\"Deploying VoIP service to telecom infrastructure\");",
    "        return pbxManager.deployVoipService((VoipConfiguration) config);",
    "    }",
    "    ",
    "    protected void postDeploymentHook(DeploymentResult result) {",
    "        // VoIP-specific post-deployment tasks",
    "        System.out.println(\"Running VoIP-specific tests and quality checks\");",
    "        qualityAssurance.runVoipTests(result.getServiceEndpoint());",
    "        monitoringSystem.addVoipService(result.getServiceId());",
    "    }",
    "}",
    "",
    "// Concrete Implementation for Internet Service",
    "class InternetServiceProvisioning extends ServiceProvisioningTemplate {",
    "    ",
    "    protected ResourceAllocation allocateResources(ServiceRequest request) {",
    "        System.out.println(\"Allocating Internet resources: bandwidth, IP range\");",
    "        return new ResourceAllocation()",
    "                .addResource(\"bandwidth\", request.getBandwidth())",
    "                .addResource(\"ip_range\", ipPoolManager.allocateRange(request.getRequiredIPs()));",
    "    }",
    "}"
]
notes_text = "Trainer-Notizen:\n- VoIP und Internet Service Beispiele\n- Unterschiedliche Implementierung gleicher Schritte\n- Hook Method für service-spezifische Tasks\n- Template garantiert konsistenten Ablauf\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Template Method Pattern - Concrete Implementations", template_concrete, notes_text)

# Slide 109: Visitor Pattern - Problem
visitor_problem = [
    "• Problem: Neue Operationen auf bestehende Objekthierarchie",
    "• Beispiel: Telekom Network Analysis",
    "  - Network Elements: Router, Switch, Firewall, LoadBalancer",
    "  - Operations: ConfigValidation, SecurityAudit, PerformanceReport",
    "• Herausforderungen:",
    "  - Neue Operation = Änderung aller Klassen",
    "  - Violation of Open-Closed Principle",
    "  - Code-Scattering across Hierarchy",
    "• Double Dispatch Problem:",
    "  - Operation depends on 2 types: Visitor + Element",
    "  - Single dispatch only considers receiver type"
]
notes_text = "Trainer-Notizen:\n- Expression Problem in OOP\n- Adding Operations vs Adding Types\n- Network Analysis als praktisches Beispiel\n- Zeitrahmen: 4 Minuten"
add_content_slide(prs, "Visitor Pattern - Problem", visitor_problem, notes_text)

# Slide 110: Visitor Pattern - Structure
visitor_structure = [
    "• Visitor Interface:",
    "  - visit(ConcreteElementA): void",
    "  - visit(ConcreteElementB): void",
    "  - ...(für jeden Element-Typ)",
    "• Concrete Visitors:",
    "  - Implementieren Operations für alle Element-Typen",
    "• Element Interface:",
    "  - accept(visitor): void",
    "• Concrete Elements:",
    "  - accept(visitor) { visitor.visit(this); }",
    "• Object Structure:",
    "  - Container für Elements",
    "  - Iteration über Elements",
    "• Double Dispatch:",
    "  - element.accept(visitor) → visitor.visit(element)"
]
notes_text = "Trainer-Notizen:\n- Double Dispatch Mechanismus\n- Visitor für jeden Element-Typ\n- Accept Method in Elements\n- Object Structure als Container\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Visitor Pattern - Struktur", visitor_structure, notes_text)

# Slide 111: Visitor Pattern - Implementation
visitor_implementation = [
    "```java",
    "// Element Interface",
    "interface NetworkElement {",
    "    void accept(NetworkElementVisitor visitor);",
    "    String getName();",
    "}",
    "",
    "// Concrete Elements",
    "class Router implements NetworkElement {",
    "    private String name;",
    "    private List<RoutingEntry> routingTable;",
    "    private int portCount;",
    "    ",
    "    public void accept(NetworkElementVisitor visitor) {",
    "        visitor.visit(this);",
    "    }",
    "    // getters...",
    "}",
    "",
    "class Switch implements NetworkElement {",
    "    private String name;",
    "    private List<VlanConfiguration> vlans;",
    "    private int portCount;",
    "    ",
    "    public void accept(NetworkElementVisitor visitor) {",
    "        visitor.visit(this);",
    "    }",
    "    // getters...",
    "}",
    "",
    "class Firewall implements NetworkElement {",
    "    private String name;",
    "    private List<SecurityRule> securityRules;",
    "    private String firewallPolicy;",
    "    ",
    "    public void accept(NetworkElementVisitor visitor) {",
    "        visitor.visit(this);",
    "    }",
    "    // getters...",
    "}"
]
notes_text = "Trainer-Notizen:\n- Network Elements als Concrete Elements\n- Accept Method delegiert an Visitor\n- this-Parameter für Double Dispatch\n- Element-spezifische Properties\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Visitor Pattern - Implementation", visitor_implementation, notes_text)

# Slide 112: Visitor Pattern - Concrete Visitors
visitor_concrete = [
    "```java",
    "// Visitor Interface",
    "interface NetworkElementVisitor {",
    "    void visit(Router router);",
    "    void visit(Switch networkSwitch);",
    "    void visit(Firewall firewall);",
    "}",
    "",
    "// Security Audit Visitor",
    "class SecurityAuditVisitor implements NetworkElementVisitor {",
    "    private List<SecurityIssue> issues = new ArrayList<>();",
    "    ",
    "    public void visit(Router router) {",
    "        System.out.println(\"Auditing router security: \" + router.getName());",
    "        // Check router-specific security",
    "        if (router.getRoutingTable().stream().anyMatch(entry -> entry.isDefaultRoute())) {",
    "            issues.add(new SecurityIssue(\"Router has default route\", router.getName()));",
    "        }",
    "    }",
    "    ",
    "    public void visit(Switch networkSwitch) {",
    "        System.out.println(\"Auditing switch security: \" + networkSwitch.getName());",
    "        // Check switch-specific security",
    "        if (networkSwitch.getVlans().stream().anyMatch(vlan -> vlan.isDefaultVlan())) {",
    "            issues.add(new SecurityIssue(\"Switch uses default VLAN\", networkSwitch.getName()));",
    "        }",
    "    }",
    "    ",
    "    public void visit(Firewall firewall) {",
    "        System.out.println(\"Auditing firewall security: \" + firewall.getName());",
    "        // Check firewall-specific security",
    "        if (firewall.getSecurityRules().stream().anyMatch(rule -> rule.isAllowAll())) {",
    "            issues.add(new SecurityIssue(\"Firewall has permissive rules\", firewall.getName()));",
    "        }",
    "    }",
    "    ",
    "    public List<SecurityIssue> getSecurityIssues() { return issues; }",
    "}"
]
notes_text = "Trainer-Notizen:\n- Security Audit als Beispiel-Visitor\n- Element-spezifische Audit-Logik\n- Accumulation von Audit-Ergebnissen\n- Typsichere Operation-Dispatch\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Visitor Pattern - Concrete Visitors", visitor_concrete, notes_text)

# Slide 113: Visitor Pattern - Usage & Benefits
visitor_usage = [
    "```java",
    "// Performance Analysis Visitor",
    "class PerformanceAnalysisVisitor implements NetworkElementVisitor {",
    "    private PerformanceReport report = new PerformanceReport();",
    "    ",
    "    public void visit(Router router) {",
    "        double cpuUsage = router.getCpuUsage();",
    "        double memoryUsage = router.getMemoryUsage();",
    "        report.addRouterMetrics(router.getName(), cpuUsage, memoryUsage);",
    "    }",
    "    ",
    "    public void visit(Switch networkSwitch) {",
    "        int activePortCount = networkSwitch.getActivePortCount();",
    "        double throughput = networkSwitch.getThroughput();",
    "        report.addSwitchMetrics(networkSwitch.getName(), activePortCount, throughput);",
    "    }",
    "    // ... firewall implementation",
    "}",
    "",
    "// Usage Example",
    "List<NetworkElement> network = Arrays.asList(",
    "    new Router(\"Core-Router-01\"),",
    "    new Switch(\"Access-Switch-01\"),",
    "    new Firewall(\"Edge-Firewall-01\")",
    ");",
    "",
    "// Run security audit",
    "SecurityAuditVisitor securityAudit = new SecurityAuditVisitor();",
    "network.forEach(element -> element.accept(securityAudit));",
    "List<SecurityIssue> securityIssues = securityAudit.getSecurityIssues();",
    "",
    "// Run performance analysis",
    "PerformanceAnalysisVisitor perfAnalysis = new PerformanceAnalysisVisitor();",
    "network.forEach(element -> element.accept(perfAnalysis));",
    "PerformanceReport perfReport = perfAnalysis.getReport();",
    "",
    "Benefits:",
    "• Easy to add new operations (new Visitors)",
    "• Operation logic centralized in Visitor",
    "• Type-safe dispatch to correct method",
    "• Separation of data structure and operations"
]
notes_text = "Trainer-Notizen:\n- Performance Analysis als zweiter Visitor\n- Einfaches Hinzufügen neuer Operations\n- Element Structure bleibt unverändert\n- Type Safety durch Double Dispatch\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Visitor Pattern - Usage & Benefits", visitor_usage, notes_text)

# Slide 114: Behavioral Patterns - Comparison Matrix
behavioral_comparison = [
    "```",
    "Pattern          | Use When                    | Telekom Example",
    "================|=============================|=========================",
    "Chain of Resp.  | Handler chain needed        | Support escalation",
    "Command         | Undo/Redo, Queuing needed  | Network config changes",
    "Iterator        | Complex data traversal      | Service hierarchy nav.",
    "Mediator        | Complex inter-communication | Network operations center",
    "Memento         | State backup/restore needed | Config snapshots",
    "Observer        | Event notifications needed  | Service status updates",
    "State           | State-dependent behavior    | Service lifecycle",
    "Strategy        | Runtime algorithm selection | Dynamic routing",
    "Template Method | Common algorithm structure  | Service provisioning",
    "Visitor         | Operations on object types  | Network analysis",
    "```",
    "",
    "**Auswahlkriterien:**",
    "• **Kommunikation**: Mediator, Observer, Chain of Responsibility",
    "• **Algorithm Flexibility**: Strategy, Template Method",
    "• **State Management**: State, Memento",
    "• **Operation Extension**: Visitor, Command"
]
notes_text = "Trainer-Notizen:\n- Übersicht aller Behavioral Patterns\n- Telekom-spezifische Anwendungsfälle\n- Auswahlkriterien für Pattern-Entscheidung\n- Diskussion: Welche Patterns kennen Teilnehmer?\n- Zeitrahmen: 8 Minuten"
add_content_slide(prs, "Behavioral Patterns - Comparison Matrix", behavioral_comparison, notes_text)

# Slide 115: Behavioral Patterns - Telekom Use Cases (1/2)
behavioral_telekom1 = [
    "**Network Operations & Management:**",
    "• **Chain of Responsibility**: Support-Ticket-Eskalation",
    "  - L1 → L2 → L3 → Expert Support",
    "• **Command**: Network-Konfiguration mit Rollback",
    "  - Router/Switch Config Commands mit Undo",
    "• **Mediator**: Network Operations Center (NOC)",
    "  - Zentrale Koordination aller Network-Komponenten",
    "• **Observer**: Service Health Monitoring",
    "  - Dashboard, Alerting, Logging bei Status-Änderungen",
    "• **State**: Service Lifecycle Management",
    "  - Pending → Active → Suspended → Terminated"
]
notes_text = "Trainer-Notizen:\n- Konkrete Telekom Anwendungsfälle\n- Network Operations Focus\n- Reale Implementierungen in der Praxis\n- Diskussion: Ähnliche Erfahrungen?\n- Zeitrahmen: 5 Minuten"
add_content_slide(prs, "Behavioral Patterns - Telekom Use Cases (1/2)", behavioral_telekom1, notes_text)

# Slide 116: Behavioral Patterns - Telekom Use Cases (2/2)  
behavioral_telekom2 = [
    "**Service Delivery & Analysis:**",
    "• **Strategy**: Dynamic Network Routing",
    "  - Shortest Path, Load Balancing, High Availability",
    "• **Template Method**: Service Provisioning",
    "  - VoIP, Internet, MPLS - gemeinsamer Workflow",
    "• **Iterator**: Organisationsstruktur-Navigation",
    "  - Services → SubServices → Components",
    "• **Visitor**: Network Infrastructure Analysis",
    "  - Security Audit, Performance Analysis, Compliance Check",
    "• **Memento**: Configuration Backup/Restore",
    "  - Snapshots vor kritischen Änderungen",
    "",
    "**Pattern-Kombinationen:**",
    "• Observer + Command: Event-driven Config Changes",
    "• State + Strategy: Zustandsabhängige Routing-Strategien",
    "• Template Method + Strategy: Algorithmus-Familie in Workflows"
]
notes_text = "Trainer-Notizen:\n- Service Delivery Focus\n- Pattern-Kombinationen in realen Systemen\n- Diskussion: Implementierungserfahrungen\n- Überleitung zu praktischen Übungen\n- Zeitrahmen: 6 Minuten"
add_content_slide(prs, "Behavioral Patterns - Telekom Use Cases (2/2)", behavioral_telekom2, notes_text)

# ============================================================================
# END OF BEHAVIORAL PATTERNS SECTION
# ============================================================================

# Save the presentation
prs.save('telekom-architecture-training-complete.pptx')
print("Complete presentation created successfully!")
print("- All 116 slides generated ✓")
print("- Used Cloud Fundamentals template ✓")
print("- Open Sans (Regular) for titles (NO BOLD) ✓")
print("- Open Sans Light for content ✓")
print("- Source Code Pro for code examples ✓")
print("- Content aligned to TOP ✓")
print("- Trainer notes added to key slides ✓")
print("- Professional layout from template ✓")
print("\nSlides created:")
print("1-12: Introduction Section")
print("1. Title: Software-Architektur Training")
print("2. Trainer Introduction")
print("3. Workshop Overview")
print("4. Workshop-Agenda")
print("5. Lernziele")
print("6. Methodik")
print("7. Software-Architektur Definition")
print("8. Architektur vs Design")
print("9. Clean Architecture")
print("10. SOLID Prinzipien")
print("11. DRY, KISS, YAGNI")
print("12. Code-Smells")
print("\n13-34: Creational Patterns Section")
print("13. Design Patterns Section Divider")
print("14. Creational Patterns Overview")
print("15. Factory Method - Problem")
print("16. Factory Method - Problematic Code")
print("17. Factory Method - Solution Structure")
print("18. Factory Method - Customer Interface")
print("19. Factory Method - Factory Implementation")
print("20. Factory Method - Telekom Example")
print("21. Abstract Factory - Problem")
print("22. Abstract Factory - Structure")
print("23. Abstract Factory - Telekom Use Case")
print("24. Builder Pattern - Problem")
print("25. Builder Pattern - Implementation")
print("26. Builder Pattern - Usage")
print("27. Prototype Pattern - Problem")
print("28. Prototype Pattern - Implementation")
print("29. Prototype - Deep vs Shallow Copy")
print("30. Singleton Pattern - Problem")
print("31. Singleton - Thread-Safe Implementation")
print("32. Singleton - Problems and Alternatives")
print("33. Creational Patterns - Comparison")
print("34. Pattern Selection Guide")
print("\n35-70: Structural Patterns Section")
print("35. Structural Patterns Section Divider")
print("36. Structural Patterns Overview")
print("37. Adapter Pattern - Problem")
print("38. Adapter Pattern - Structure")
print("39. Adapter Pattern - Legacy System Code")
print("40. Adapter Pattern - Implementation")
print("41. Bridge Pattern - Problem")
print("42. Bridge Pattern - Structure")
print("43. Bridge Pattern - Implementation")
print("44. Bridge Pattern - Refined Abstractions")
print("45. Composite Pattern - Problem")
print("46. Composite Pattern - Structure")
print("47. Composite Pattern - Component Interface")
print("48. Composite Pattern - Composite Implementation")
print("49. Composite Pattern - Usage Example")
print("50. Decorator Pattern - Problem")
print("51. Decorator Pattern - Structure")
print("52. Decorator Pattern - Component & Base")
print("53. Decorator Pattern - Base Decorator")
print("54. Decorator Pattern - Multiple Decorators")
print("55. Facade Pattern - Problem")
print("56. Facade Pattern - Complex Subsystems")
print("57. Facade Pattern - Facade Implementation")
print("58. Facade Pattern - Usage & Benefits")
print("59. Flyweight Pattern - Problem")
print("60. Flyweight Pattern - Structure")
print("61. Flyweight Pattern - Implementation")
print("62. Flyweight Pattern - Factory & Usage")
print("63. Proxy Pattern - Problem")
print("64. Proxy Pattern - Types")
print("65. Proxy Pattern - Virtual Proxy")
print("66. Proxy Pattern - Virtual Proxy Implementation")
print("67. Proxy Pattern - Protection Proxy")
print("68. Structural Patterns - Comparison Matrix")
print("69. Structural Patterns - Telekom Use Cases (1/2)")
print("70. Structural Patterns - Telekom Use Cases (2/2)")
print("\n🎯 STRUCTURAL PATTERNS SECTION COMPLETE!")
print("- All 7 structural patterns covered")
print("- Adapter pattern for legacy integration")
print("- Bridge pattern for notification systems")
print("- Composite pattern for hierarchical structures")
print("- Decorator pattern for dynamic service extensions")
print("- Facade pattern for complex subsystem orchestration")
print("- Flyweight pattern for memory optimization")
print("- Proxy pattern for lazy loading and access control")
print("- Comprehensive code examples with Telekom use cases")
print("- Pattern comparison matrix and real-world applications")
print("- Professional trainer notes throughout")

print("\n71-116: Behavioral Patterns Section")
print("71. Behavioral Patterns Section Divider")
print("72. Behavioral Patterns Overview")
print("73. Chain of Responsibility - Problem")
print("74. Chain of Responsibility - Structure") 
print("75. Chain of Responsibility - Implementation")
print("76. Chain of Responsibility - Concrete Handlers")
print("77. Command Pattern - Problem")
print("78. Command Pattern - Structure")
print("79. Command Pattern - Implementation")
print("80. Command Pattern - Invoker & Macro")
print("81. Iterator Pattern - Problem")
print("82. Iterator Pattern - Structure")
print("83. Iterator Pattern - Implementation")
print("84. Iterator Pattern - Usage & Benefits")
print("85. Mediator Pattern - Problem")
print("86. Mediator Pattern - Structure")
print("87. Mediator Pattern - Implementation")
print("88. Mediator Pattern - Concrete Mediator")
print("89. Memento Pattern - Problem")
print("90. Memento Pattern - Structure")
print("91. Memento Pattern - Implementation")
print("92. Memento Pattern - Caretaker")
print("93. Observer Pattern - Problem")
print("94. Observer Pattern - Structure")
print("95. Observer Pattern - Implementation")
print("96. Observer Pattern - Concrete Observers")
print("97. State Pattern - Problem")
print("98. State Pattern - Structure")
print("99. State Pattern - Implementation")
print("100. State Pattern - Concrete States")
print("101. Strategy Pattern - Problem")
print("102. Strategy Pattern - Structure")
print("103. Strategy Pattern - Implementation")
print("104. Strategy Pattern - Context & Usage")
print("105. Template Method Pattern - Problem")
print("106. Template Method Pattern - Structure")
print("107. Template Method Pattern - Implementation")
print("108. Template Method Pattern - Concrete Implementations")
print("109. Visitor Pattern - Problem")
print("110. Visitor Pattern - Structure")
print("111. Visitor Pattern - Implementation")
print("112. Visitor Pattern - Concrete Visitors")
print("113. Visitor Pattern - Usage & Benefits")
print("114. Behavioral Patterns - Comparison Matrix")
print("115. Behavioral Patterns - Telekom Use Cases (1/2)")
print("116. Behavioral Patterns - Telekom Use Cases (2/2)")
print("\n🎯 BEHAVIORAL PATTERNS SECTION COMPLETE!")
print("- All 10 behavioral patterns covered")
print("- Chain of Responsibility for support escalation")
print("- Command pattern with undo/redo functionality")
print("- Iterator for service hierarchy traversal")
print("- Mediator for network operations center coordination")
print("- Memento for configuration backup/restore")
print("- Observer for event-driven service monitoring")
print("- State pattern for service lifecycle management")
print("- Strategy for dynamic routing algorithms")
print("- Template Method for service provisioning workflows")
print("- Visitor for network infrastructure analysis")
print("- Comprehensive Telekom-specific examples")
print("- Pattern comparison matrix and selection guide")
print("- Professional trainer notes throughout")