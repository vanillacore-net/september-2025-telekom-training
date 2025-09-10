#!/usr/bin/env python3
"""
Block 4 Advanced Patterns Presentation Generator
Fixes:
1. Uses Python script with python-pptx (NO MCP server)
2. NO DOUBLE BULLETS - Layout 2 has built-in bullets
3. Adds SPEAKER NOTES from all "Note:" sections
4. Fixes CODE LINE SPACING - single-spaced using Consolas
5. Uses JAVA for all code examples
6. Uses VanillaCore.pptx template properly
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def setup_code_formatting(paragraph):
    """Apply proper code formatting with single line spacing"""
    paragraph.font.name = 'Consolas'
    paragraph.font.size = Pt(10)
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.0  # Single line spacing

def add_slide_with_layout(prs, layout_index, title_text, content_lines=None, speaker_notes=None, is_code=False):
    """
    Add slide with proper layout handling
    Layout 0: Title slide
    Layout 1: Section headers (NO content)
    Layout 2: Has built-in bullets - NO bullet characters needed
    Layout 3: Plain text without bullets
    Layout 9: Code blocks with Consolas font and single line spacing
    """
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    
    # Set title
    if hasattr(slide.shapes, 'title') and slide.shapes.title:
        slide.shapes.title.text = title_text
    
    # Add content based on layout
    if content_lines and layout_index in [2, 3, 9]:
        if layout_index == 9 or is_code:
            # Code layout - use text box with Consolas
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            # Add code content with proper formatting
            for i, line in enumerate(content_lines):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = line
                p.font.name = 'Consolas'
                p.font.size = Pt(10)
                p.line_spacing = 1.0
                p.space_after = Pt(0)
                p.space_before = Pt(0)
                
        else:
            # Regular content - use content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for i, line in enumerate(content_lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # For Layout 2: NO bullet characters - layout has built-in bullets
                    if layout_index == 2:
                        # Remove bullet characters if present
                        clean_line = line.lstrip('•- ').strip()
                        p.text = clean_line
                    else:
                        p.text = line
    
    # Add speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = speaker_notes

def main():
    """Generate Block 4 Advanced Patterns presentation"""
    
    # Paths
    template_path = "/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/templates/VanillaCore.pptx"
    output_path = "/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/block4-presentation.pptx"
    
    # Check template exists
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return False
    
    try:
        # Load template
        prs = Presentation(template_path)
        print(f"✅ Loaded template: {template_path}")
        
        # Slide 1: Title slide
        add_slide_with_layout(
            prs, 0, "Block 4: Advanced Patterns & Integration",
            speaker_notes="Block 4 behandelt die fortgeschrittenen Design Patterns, die für komplexe Enterprise-Architektur unerlässlich sind. Focus liegt auf praktischer Integration aller Patterns in produktiven Telekom-Systemen."
        )
        
        # Slide 2: Block 4 Overview
        overview_content = [
            "Mediator Pattern: Zentrale Orchestrierung statt Communication Chaos",
            "Iterator & Visitor: Intelligente Datenverarbeitung ohne Type-Casting", 
            "Memento & Interpreter: State-Recovery und Configuration-DSLs",
            "Pattern Integration: Alles zusammenfügen in produktiver Architektur",
            "Anti-Patterns vermeiden: Pattern-Obsession und Over-Engineering",
            "Team-Adoption: Graduelle Einführung und nachhaltige Nutzung",
            "Production Readiness: Monitoring, Performance und Fehlerbehandlung"
        ]
        add_slide_with_layout(
            prs, 2, "Block 4 Überblick", overview_content,
            "Block 4 behandelt die fortgeschrittenen Design Patterns, die für komplexe Enterprise-Architektur unerlässlich sind. Focus liegt auf praktischer Integration aller Patterns in produktiven Telekom-Systemen."
        )
        
        # Slide 3: Mediator Pattern Section
        add_slide_with_layout(prs, 1, "Pattern 1: Mediator Pattern")
        
        # Slide 4: Mediator - Bad Example
        bad_code = [
            "// Communication Explosion: Jeder redet mit jedem",
            "class NetworkDevice {",
            "    private List<Router> routers;           // 50+ Router",
            "    private List<Switch> switches;          // 200+ Switches", 
            "    private List<FirewallDevice> firewalls; // 30+ Firewalls",
            "    private List<MonitoringSystem> monitors; // 10+ Monitoring",
            "    ",
            "    public void statusChanged() {",
            "        // O(n²) = 50 × 200 × 30 × 10 = 3 MILLIONEN Updates!",
            "        for (Router r : routers) {",
            "            r.updateTopology(this);",
            "            for (Switch s : switches) {",
            "                s.recalculateRoutes(r, this);",
            "                for (FirewallDevice f : firewalls) {",
            "                    f.updateRules(r, s, this);", 
            "                    // HORROR: 4-fach verschachtelte Loops!",
            "                }",
            "            }",
            "        }",
            "    }",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Was ist hier schlecht?", bad_code,
            "Point-to-Point Communication führt zu exponentieller Komplexität. Bei 300+ Netzwerkgeräten entstehen über 3 Millionen Update-Operationen bei jeder Änderung.",
            is_code=True
        )
        
        # Slide 5: Code Smells
        smells_content = [
            "Performance-Kollaps: Ein Device-Change triggert 1000+ Notifications",
            "Network Storms: Broadcast-Messages überlasten Management-Netz",
            "Deadlocks: Circular Dependencies zwischen Devices",
            "Maintenance-Horror: Neue Device-Art muss mit ALLEN Types integriert werden",
            "Testing-Explosion: Jeder Change gegen 100+ Device-Kombinationen testen",
            "Root-Cause-Analysis: \"Warum ist Router-47 langsam?\" → 72h Investigation",
            "Change-Impact: Unbekannt, zu komplex zu analysieren"
        ]
        add_slide_with_layout(
            prs, 2, "Code Smells identifiziert", smells_content,
            "Die direkte Kommunikation zwischen allen Netzwerkgeräten führt zu unbeherrschbarer Komplexität und kritischen Performance-Problemen."
        )
        
        # Slide 6: Mediator Solution
        solution_content = [
            "Zentrale Koordination: Ein Mediator statt 50.000 Point-to-Point Verbindungen",
            "O(n) statt O(n²): Alle reden mit einem Mediator, nicht miteinander",
            "Handler-Architektur: Spezialisierte Handler für verschiedene Device-Types",
            "Error Isolation: Ein Handler-Fehler stoppt nicht andere Handler",
            "Priority-based Processing: Routing vor Monitoring, Security nach Routing",
            "Event-driven Design: Async Processing für Performance",
            "Extensibility: Neue Handler ohne Änderung bestehender Code"
        ]
        add_slide_with_layout(
            prs, 2, "Lösung: Mediator Pattern", solution_content,
            "Das Mediator Pattern reduziert die Komplexität von exponentiell auf linear und ermöglicht zentrale Orchestrierung aller Netzwerkoperationen."
        )
        
        # Slide 7: Mediator Implementation
        mediator_code = [
            "@Component",
            "public class TelekomNetworkOrchestrator implements NetworkMediator {",
            "    ",
            "    private final Map<DeviceType, List<DeviceHandler>> handlers = new ConcurrentHashMap<>();",
            "    ",
            "    @Override",
            "    public void deviceStatusChanged(NetworkDevice device, DeviceStatus status) {",
            "        ",
            "        log.info(\"📡 Device {} status: {} → {}\", ",
            "                device.getId(), device.getPreviousStatus(), status);",
            "        ",
            "        DeviceChangeEvent event = new DeviceChangeEvent(device, status);",
            "        List<DeviceHandler> orderedHandlers = getOrderedHandlers(device.getType());",
            "        ",
            "        // Parallel Processing für Performance",
            "        orderedHandlers.parallelStream().forEach(handler -> {",
            "            try {",
            "                handler.handle(event);",
            "            } catch (Exception e) {",
            "                log.error(\"❌ Handler {} failed\", handler.getClass().getSimpleName(), e);",
            "                // Error Isolation: Ein Handler-Fehler stoppt nicht die anderen",
            "            }",
            "        });",
            "    }",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Implementierung", mediator_code,
            "Production-ready Implementation mit Parallel Processing, Error Isolation und Metrics. Jeder Handler kümmert sich um einen spezifischen Concern.",
            is_code=True
        )
        
        # Slide 8: Iterator Pattern Section
        add_slide_with_layout(prs, 1, "Pattern 2: Iterator Pattern")
        
        # Slide 9: Iterator - Bad Example  
        iterator_bad = [
            "// Navigation und Type-Casting Horror",
            "class NetworkTopologyReport {",
            "    public String generateReport(NetworkTopology topology, ReportType reportType) {",
            "        StringBuilder report = new StringBuilder();",
            "        ",
            "        for (NetworkNode node : topology.getNodes()) {",
            "            if (node instanceof Router) {",
            "                Router r = (Router) node;",
            "                // 100 Zeilen Router-spezifischer XML Logic",
            "            } else if (node instanceof Switch) {",
            "                Switch s = (Switch) node;",
            "                // 150 Zeilen Switch-spezifischer XML Logic",
            "            } else if (node instanceof FirewallDevice) {",
            "                // 200 Zeilen Firewall XML Logic",
            "            }",
            "            // Was passiert mit neuen Device-Types???",
            "        }",
            "        ",
            "        // KOMPLETT ANDERE 500 Zeilen für JSON...",
            "        // NOCHMAL 600 Zeilen für PDF...",
            "    }",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Was ist hier schlecht?", iterator_bad,
            "instanceof-Horror führt zu N×M Komplexität: N Device-Types × M Report-Formats. Jeder neue Type erfordert Änderungen in allen Formaten.",
            is_code=True
        )
        
        # Slide 10: Iterator Code Smells
        iterator_smells = [
            "N × M Complexity: N Device-Types × M Report-Formats = exponentielle Code-Pfade",
            "Instanceof-Horror: Type-Casting überall, fehleranfällig", 
            "Code Duplication: Device-Logic wird pro Format wiederholt",
            "Maintenance-Nightmare: Bug in Router-Logic = Fix in XML, JSON UND PDF",
            "Testing-Explosion: 5 Device-Types × 4 Formate = 20 Test-Kombinationen",
            "Knowledge Requirements: Developer braucht XML, JSON UND PDF Expertise",
            "ConcurrentModificationException: Unsafe Collection Modification während Iteration"
        ]
        add_slide_with_layout(
            prs, 2, "Code Smells identifiziert", iterator_smells,
            "Die Vermischung von Navigation und Datenverarbeitung führt zu unlesbarem und wartungsfeindlichem Code."
        )
        
        # Slide 11: Iterator Solution
        iterator_solution = [
            "Separation of Concerns: Iterator (WIE navigieren) + Visitor (WAS machen)",
            "Type-Safe Operations: Keine instanceof-Checks zur Runtime",
            "Method Overloading: Jeder Device-Type hat optimierte Behandlung",
            "Safe Navigation: ConcurrentModificationException vermeiden", 
            "Extensibility: Neue Operations ohne Datenstruktur-Änderungen",
            "Parallel Processing: Stream-Integration für Performance",
            "Cycle Detection: Verhindert Infinite Loops in Mesh-Topologien"
        ]
        add_slide_with_layout(
            prs, 2, "Lösung: Iterator + Visitor Pattern", iterator_solution,
            "Die Kombination von Iterator und Visitor Pattern trennt Navigation von Verarbeitung und ermöglicht type-safe Operations ohne Casting."
        )
        
        # Slide 12: Iterator Implementation
        iterator_impl = [
            "// Iterator für sichere Navigation",
            "public class TopologyBreadthFirstIterator implements NetworkIterator<NetworkNode> {",
            "    private final Queue<NetworkNode> queue = new LinkedList<>();",
            "    private final Set<NetworkNode> visited = new HashSet<>();",
            "    ",
            "    @Override",
            "    public NetworkNode next() {",
            "        NetworkNode current = queue.poll();",
            "        if (!visited.contains(current)) {",
            "            visited.add(current);",
            "            // Add connected nodes for BFS",
            "            topology.getConnectedNodes(current).stream()",
            "                .filter(node -> !visited.contains(node))",
            "                .forEach(queue::offer);",
            "        }",
            "        return current;",
            "    }",
            "}",
            "",
            "// Visitor für type-safe Processing",
            "public class XmlReportVisitor implements NetworkNodeVisitor<String> {",
            "    @Override",
            "    public String visitRouter(Router router) {",
            "        // Router-spezifische XML-Generierung",
            "        return generateRouterXml(router);",
            "    }",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Implementierung", iterator_impl,
            "Safe Navigation mit Cycle Detection und type-safe Operations durch Method Overloading. Keine Casting-Fehler zur Runtime.",
            is_code=True
        )
        
        # Slide 13: Memento Pattern Section
        add_slide_with_layout(prs, 1, "Pattern 3: Memento Pattern")
        
        # Slide 14: Memento - Bad Example
        memento_bad = [
            "// Kein Rollback, kein Recovery",
            "class NetworkDeviceConfigurator {",
            "    public void applyConfiguration(NetworkDevice device, Configuration newConfig) {",
            "        ",
            "        // HORROR: Direkte State-Mutations ohne Backup",
            "        device.setRoutingTable(newConfig.getRoutes());        // Was war vorher?",
            "        device.setVlanConfiguration(newConfig.getVlans());    // Unbekannt!",
            "        device.setSecurityPolicies(newConfig.getSecurityPolicies()); // Lost!",
            "        ",
            "        try {",
            "            device.commitConfiguration(); // 🎲 All-or-Nothing Gamble",
            "        } catch (ConfigurationException e) {",
            "            // We're screwed - no way back!",
            "            log.error(\"🚨 Device in UNKNOWN state!\", e);",
            "            // Manual Recovery Required: 2+ hours debugging",
            "            throw new RuntimeException(\"Device unrecoverable\", e);",
            "        }",
            "    }",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Was ist hier schlecht?", memento_bad,
            "Fehlgeschlagene Configuration-Changes ohne Rollback-Möglichkeit führen zu stundenlangen manuellen Recovery-Prozessen und kritischen Ausfällen.",
            is_code=True
        )
        
        # Slide 15: Memento Code Smells
        memento_smells = [
            "No Backup: Vorheriger Zustand unwiderruflich verloren",
            "Partial Updates: Device in inkonsistentem Zustand bei Fehlern",
            "No Atomicity: Kein sauberer Rollback bei partial Failures",
            "Manual Recovery: Stunden von Expert-Zeit für Wiederherstellung",
            "Production Outages: Configuration-Fehler → 30 Min Komplettausfall",
            "Audit & Compliance: Keine Historie für Regulatory Requirements", 
            "Expert-Dependency: Nur 3-4 Personen können kritische Configs wiederherstellen"
        ]
        add_slide_with_layout(
            prs, 2, "Code Smells identifiziert", memento_smells,
            "In Telekom-Netzen kann ein Configuration-Fehler 50.000+ Kunden betreffen. Ohne Rollback-Mechanismen sind die Recovery-Zeiten inakzeptabel."
        )
        
        # Slide 16: Memento Solution
        memento_solution = [
            "Production-Safe Changes: Garantierte Rollback-Möglichkeit",
            "Atomic Operations: Entweder alles oder nichts, aber sauber",
            "Audit-Compliance: Vollständige Historie aller Changes",
            "Disaster Recovery: State-Snapshots für schnelle Wiederherstellung", 
            "Immutable Snapshots: Deep Copy verhindert versehentliche Mutations",
            "Integrity Validation: Hash-basierte Corruption Detection",
            "Multi-Device Coordination: Orchestrierte Changes mit Rollback"
        ]
        add_slide_with_layout(
            prs, 2, "Lösung: Memento Pattern", memento_solution,
            "Das Memento Pattern ermöglicht atomare Configuration-Changes mit garantiertem Rollback und vollständiger Audit-Historie."
        )
        
        # Slide 17: Memento Implementation
        memento_impl = [
            "// Immutable Memento für Network Device State",
            "public class NetworkDeviceMemento {",
            "    private final String deviceId;",
            "    private final LocalDateTime timestamp;",
            "    private final Map<String, Object> configurationSnapshot;",
            "    private final String configurationHash;",
            "    ",
            "    // Package-private Constructor - nur NetworkDevice kann Mementos erstellen",
            "    NetworkDeviceMemento(String deviceId, Map<String, Object> configuration) {",
            "        this.deviceId = deviceId;",
            "        this.timestamp = LocalDateTime.now();",
            "        this.configurationSnapshot = deepCopyConfiguration(configuration);",
            "        this.configurationHash = calculateConfigHash(configurationSnapshot);",
            "    }",
            "    ",
            "    // Integrity Validation",
            "    public boolean validateIntegrity() {",
            "        return configurationHash.equals(calculateConfigHash(configurationSnapshot));",
            "    }",
            "}",
            "",
            "// Multi-Device Emergency Rollback",
            "public EmergencyRollbackResult performEmergencyRollback(",
            "        Map<String, NetworkDeviceMemento> preChangeSnapshots) {",
            "    ",
            "    preChangeSnapshots.entrySet().parallelStream().forEach(entry -> {",
            "        NetworkDevice device = deviceService.getDevice(entry.getKey());",
            "        device.restoreFromMemento(entry.getValue());",
            "    });",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Implementierung", memento_impl,
            "Production-ready Implementation mit Integrity-Validation und parallelem Emergency-Rollback für kritische Situationen.",
            is_code=True
        )
        
        # Slide 18: Interpreter Pattern Section
        add_slide_with_layout(prs, 1, "Pattern 4: Interpreter Pattern")
        
        # Slide 19: Interpreter - Bad Example
        interpreter_bad = [
            "// String-Parsing Nightmare",
            "public void parseNetworkConfig(String configText) {",
            "    String[] lines = configText.split(\"\\n\");",
            "    ",
            "    for (String line : lines) {",
            "        if (line.startsWith(\"route\")) {",
            "            String[] parts = line.split(\" \"); // Was bei Tabs? Extra-Spaces?",
            "            if (parts.length >= 3) { // Was wenn parts.length == 2?",
            "                String destination = parts[1]; // Was wenn leer?",
            "                String gateway = parts[2];     // Was wenn invalid IP?",
            "                // Keine Syntax Validation!",
            "                addRoute(destination, gateway); // 💣 Potential Bomb",
            "            }",
            "        } else if (line.startsWith(\"vlan\")) {",
            "            int vlanId = Integer.parseInt(parts[1]); // NumberFormatException?",
            "            // ... 100+ weitere Zeilen String-Horror",
            "        }",
            "    }",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Was ist hier schlecht?", interpreter_bad,
            "String-basierte Configuration führt zu fragiler Syntax, fehlendem IDE-Support und schwer debugbaren Fehlern.",
            is_code=True
        )
        
        # Slide 20: Interpreter Code Smells
        interpreter_smells = [
            "Fragile Syntax: Whitespace-sensitive, error-prone parsing",
            "No Validation: Invalid IPs, VLAN IDs werden nicht caught",
            "Error Handling: Silent failures oder Exception-Chaos",
            "Debugging Horror: \"Zeile 247 von 2000 ist falsch\" - good luck",
            "No IDE Support: Keine Syntax-Highlighting, Auto-Completion",
            "Expert Dependency: Java-Code für jede neue Configuration-Rule",
            "Change Velocity: 2 Tage Development für einfache Firewall-Rule"
        ]
        add_slide_with_layout(
            prs, 2, "Code Smells identifiziert", interpreter_smells,
            "String-Parsing macht Network Engineers von Entwicklern abhängig und verlangsamt kritische Configuration-Changes erheblich."
        )
        
        # Slide 21: Interpreter Solution
        interpreter_solution = [
            "Configuration-as-Code: DSL für Network Engineers ohne Java-Knowledge",
            "Type-Safe Parsing: Abstract Syntax Tree statt String-Manipulation",
            "IDE Integration: Syntax-Highlighting und Auto-Completion möglich",
            "Expert Empowerment: Fachexperten werden von Entwicklern unabhängig",
            "Comprehensive Validation: Syntax- und Semantic-Checks",
            "Grammar Evolution: DSL kann schrittweise erweitert werden",
            "Tool Support: Code-Generator für Common Patterns"
        ]
        add_slide_with_layout(
            prs, 2, "Lösung: Interpreter Pattern", interpreter_solution,
            "Das Interpreter Pattern ermöglicht domänen-spezifische Sprachen, die Experten ohne Programmierkenntnisse verwenden können."
        )
        
        # Slide 22: Interpreter Implementation
        interpreter_impl = [
            "// Abstract Syntax Tree für Network Configuration",
            "public abstract class ConfigurationExpression {",
            "    public abstract void interpret(NetworkConfigurationContext context);",
            "    public abstract void validate(ValidationContext validationContext);",
            "}",
            "",
            "// Terminal Expression für Route Configuration",
            "public class RouteExpression extends ConfigurationExpression {",
            "    private final String destinationNetwork;",
            "    private final String gatewayAddress;",
            "    private final int metric;",
            "    ",
            "    @Override",
            "    public void interpret(NetworkConfigurationContext context) {",
            "        Route route = Route.builder()",
            "            .destination(destinationNetwork)",
            "            .gateway(gatewayAddress)",
            "            .metric(metric)",
            "            .build();",
            "        context.addRoute(route);",
            "    }",
            "    ",
            "    @Override",
            "    public void validate(ValidationContext validationContext) {",
            "        if (!NetworkUtils.isValidNetworkAddress(destinationNetwork)) {",
            "            validationContext.addError(\"Invalid destination: \" + destinationNetwork);",
            "        }",
            "    }",
            "}",
            "",
            "// DSL Example:",
            "// route 192.168.1.0/24 via 10.0.0.1 metric 100",
            "// vlan 100 name \"Production Network\"",
            "//   ip address 192.168.100.1/24",
            "// end"
        ]
        add_slide_with_layout(
            prs, 9, "Implementierung", interpreter_impl,
            "Strukturierter AST ermöglicht type-safe Configuration mit umfassender Validation und klarer Syntax für Network Engineers.",
            is_code=True
        )
        
        # Slide 23: Pattern Integration Section  
        add_slide_with_layout(prs, 1, "Pattern Integration")
        
        # Slide 24: Integration - Bad Example
        integration_bad = [
            "// Pattern-Obsession Monster",
            "public class OverEngineeredStringProcessor {",
            "    // OVERKILL: Factory für simple String Operations!",
            "    private final AbstractStringOperationFactory operationFactory;",
            "    // OVERKILL: Strategy für jeden String-Vorgang!",
            "    private final Map<OperationType, StringProcessingStrategy> strategies;",
            "    // OVERKILL: Observer für String Changes!",
            "    private final List<StringChangeObserver> observers;",
            "    ",
            "    // WAHNSINN: 150 Zeilen Code für \"Hello World\".toUpperCase()",
            "    public String convertToUpperCase(String input) {",
            "        // 100+ Zeilen Pattern-Overkill für eine einzige Operation",
            "        return result; // Finally! After unnecessary complexity",
            "    }",
            "    ",
            "    // Method length: 150+ lines, Pattern count: 7 (!!!)",
            "    // Time to understand: 30+ minutes",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Was ist hier schlecht?", integration_bad,
            "Pattern-Obsession führt zu Over-Engineering. Nicht jedes Problem braucht ein Pattern - manchmal ist die einfache Lösung die beste.",
            is_code=True
        )
        
        # Slide 25: Integration Code Smells
        integration_smells = [
            "Pattern für alles: Entwickler wollen Patterns überall verwenden",
            "Over-Engineering: 150 Zeilen für eine einzige String-Operation",
            "God Mediator: 10.000-Zeilen Monster-Klassen",
            "Pattern Explosion: 50+ micro-patterns für simple Use Cases",
            "Complexity Explosion: 7 Patterns für eine triviale Operation",
            "Maintenance Horror: 30 Minuten zum Verstehen, 2h zum Debuggen",
            "Team Overwhelm: Junior Developers können Code nicht mehr lesen"
        ]
        add_slide_with_layout(
            prs, 2, "Code Smells identifiziert", integration_smells,
            "Das größte Anti-Pattern ist \"Pattern für alles\". Patterns sollen Probleme lösen, nicht schaffen."
        )
        
        # Slide 26: Integration Solution
        integration_solution = [
            "KISS Principle: Keep It Simple, Stupid - einfache Probleme, einfache Lösungen",
            "Problem-First: Pattern nur wenn echtes Problem gelöst wird",
            "Team Reality Check: Pattern-Complexity muss Team-Expertise entsprechen",
            "Performance Impact: Pattern-Overhead messen, nicht annehmen",
            "Graduelle Evolution: Foundation Patterns zuerst, Advanced später",
            "Business Value: Patterns dienen Business, nicht umgekehrt",
            "Architecture Review: Pattern-bewusste Code Reviews"
        ]
        add_slide_with_layout(
            prs, 2, "Lösung: Intelligente Pattern-Integration", integration_solution,
            "Erfolgreiche Pattern-Integration erfordert pragmatische Abwägung zwischen Nutzen und Komplexität."
        )
        
        # Slide 27: Integration Implementation
        integration_impl = [
            "// Layer 1: Foundation Patterns",
            "@Component",
            "public class TelekomConfigurationManager {",
            "    private static volatile TelekomConfigurationManager instance; // Singleton",
            "}",
            "",
            "@Component", 
            "public class TelekomDeviceFactory {",
            "    public NetworkDevice createDevice(DeviceType type, DeviceSpecification spec) {",
            "        return deviceBuilders.get(type).build(); // Factory + Builder",
            "    }",
            "}",
            "",
            "// Layer 2: Behavioral Integration",
            "@Component",
            "public class NetworkEventProcessor {",
            "    private final List<NetworkEventObserver> observers; // Observer",
            "    private final Map<EventType, EventProcessingStrategy> strategies; // Strategy",
            "}",
            "",
            "// Layer 3: Advanced Coordination",
            "@Component",
            "public class TelekomNetworkOrchestrationHub implements NetworkMediator {",
            "    // Mediator + Observer + Command + Iterator integration",
            "    ",
            "    @Override",
            "    public void handleNetworkChange(NetworkChangeEvent event) {",
            "        NetworkImpactAnalysis impact = analyzeNetworkImpact(event); // Iterator",
            "        NetworkOperation operation = createResponseOperation(impact); // Command",
            "        OperationResult result = operationOrchestrator.executeOperation(operation);",
            "        eventProcessor.processNetworkEvent(result); // Observer",
            "    }",
            "}"
        ]
        add_slide_with_layout(
            prs, 9, "Implementierung - Layer Architecture", integration_impl,
            "Layered Architecture mit gradueller Pattern-Einführung: Foundation → Behavioral → Advanced Integration für produktive Enterprise-Systeme.",
            is_code=True
        )
        
        # Slide 28: Team-Adoption & Production Readiness Section
        add_slide_with_layout(prs, 1, "Team-Adoption & Production Readiness")
        
        # Slide 29: Team-Adoption Strategy
        adoption_content = [
            "Phase 1 (Wochen 1-4): Foundation Building - Survival Patterns",
            "Phase 2 (Wochen 5-8): Behavioral Integration - Observer, Strategy", 
            "Phase 3 (Wochen 9-12): Advanced Coordination - Mediator, Command",
            "Graduelle Code-Integration: Legacy-Wrapper für sanfte Migration",
            "Pattern-bewusste Code Reviews: Checklist für Pattern Usage",
            "Training & Mentoring: Internal Pattern Workshops",
            "Metrics & ROI: Pattern-Impact messen und demonstrieren"
        ]
        add_slide_with_layout(
            prs, 2, "Team-Adoption Strategie", adoption_content,
            "Nachhaltige Pattern-Adoption braucht strukturiertes Vorgehen mit Team-Training und gradueller Migration."
        )
        
        # Slide 30: Production Readiness
        production_content = [
            "Performance Monitoring: Pattern-Execution-Time tracking",
            "Error Handling: Circuit Breaker für Pattern-Failures", 
            "Resilience: Fallback-Strategien für Pattern-Ausfälle",
            "Observability: Metrics für Pattern-Usage und Impact",
            "Documentation: Pattern-Guidelines und Best Practices",
            "Testing: Pattern-Integration Testing",
            "Rollback: Feature Flags für schrittweise Einführung"
        ]
        add_slide_with_layout(
            prs, 2, "Production Readiness", production_content,
            "Production-Einsatz erfordert umfassendes Monitoring, Error Handling und Rollback-Strategien für Pattern-basierte Systeme."
        )
        
        # Slide 31: Decision Framework
        decision_content = [
            "✅ VERWENDE Patterns wenn:",
            "Komplexität rechtfertigt Pattern-Overhead",
            "Team versteht und kann Patterns maintainen", 
            "Future Changes werden durch Pattern erleichtert",
            "Performance ist akzeptabel",
            "Pattern löst tatsächlich ein Problem",
            "",
            "❌ VERMEIDE Patterns wenn:",
            "Simple Code wird komplizierter",
            "Einmaliger Use Case ohne Erweiterung",
            "Performance-kritischer Code",
            "Team nicht bereit für Pattern-Complexity", 
            "\"Cool Factor\" ist einziger Grund"
        ]
        add_slide_with_layout(
            prs, 3, "Decision Framework", decision_content,
            "Das Pattern Selection Framework hilft bei der Entscheidung, wann Patterns Mehrwert bringen und wann sie Over-Engineering sind."
        )
        
        # Slide 32: Summary Section
        add_slide_with_layout(prs, 1, "Zusammenfassung")
        
        # Slide 33: Key Takeaways
        takeaways_content = [
            "Mediator Pattern: Löst Communication-Explosion, reduziert O(n²) auf O(n)",
            "Iterator + Visitor: Trennt Navigation von Processing, eliminiert instanceof-Horror",
            "Memento + Interpreter: Production-safe State-Management + Domain-specific Languages",
            "Pattern Integration: Layer-basierte Architektur mit gradueller Evolution",
            "Anti-Pattern Awareness: Pattern-Obsession vermeiden, KISS Principle beachten",
            "Team Success: Training, graduelle Adoption, Metrics-basierte ROI-Demonstration",
            "Production Reality: Monitoring, Error Handling, Performance-Impact beachten"
        ]
        add_slide_with_layout(
            prs, 2, "Block 4 Key Takeaways", takeaways_content,
            "Block 4 hat gezeigt, wie Advanced Patterns in produktiven Enterprise-Systemen integriert werden. Der Schlüssel ist pragmatische Abwägung zwischen Nutzen und Komplexität."
        )
        
        # Slide 34: Action Plan
        action_content = [
            "Week 1-2: Pattern Assessment - Aktuelle Code-Schmerz-Punkte analysieren",
            "Month 2-3: Team Enablement - Pilot Project und interne Workshops",
            "Month 4-6: Production Integration - Patterns in kritischen Systemen", 
            "Continuous: Pattern Community - Telekom-weite Knowledge Sharing",
            "",
            "Success Factors:",
            "Business Value First - Patterns lösen Business-Probleme",
            "Team Readiness - Skills müssen Pattern-Complexity entsprechen",
            "Gradual Evolution - Schrittweise Adoption, kein \"Big Bang\"",
            "Measure Impact - ROI dokumentieren für Stakeholder Buy-in"
        ]
        add_slide_with_layout(
            prs, 3, "Your Action Plan", action_content,
            "Erfolgreiche Pattern-Mastery ist ein Journey. Beginnt mit Foundation Patterns, entwickelt euch graduell zu Advanced Integration und teilt euer Wissen im Team."
        )
        
        # Save presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        
        print(f"✅ Block 4 presentation generated successfully!")
        print(f"📁 Output: {output_path}")
        print(f"📊 Total slides: {len(prs.slides)}")
        
        # Verify the file was created
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"📏 File size: {file_size:,} bytes")
            return True
        else:
            print("❌ File was not created successfully")
            return False
            
    except Exception as e:
        print(f"❌ Error generating presentation: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)