#!/usr/bin/env python3
"""
Fix Block 2 Font Issues - Remove Calibri overrides and let template control fonts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import re

def clean_markdown(content):
    """Remove all HedgeDoc-specific markup and clean the content"""
    
    # Remove YAML frontmatter
    content = re.sub(r'^---.*?---\s*', '', content, flags=re.DOTALL)
    
    # Remove all HTML tags and divs
    content = re.sub(r'<div[^>]*>.*?</div>', '', content, flags=re.DOTALL)
    content = re.sub(r'<[^>]+>', '', content)
    
    # Remove fragment annotations
    content = re.sub(r'<!--.*?-->', '', content)
    
    # Clean up excessive blank lines
    content = re.sub(r'\n{3,}', '\n\n', content)
    
    return content.strip()

def parse_block2_markdown(content):
    """Parse Block 2 markdown into slide data"""
    
    slides = []
    raw_slides = content.split('\n---\n')
    
    for slide_text in raw_slides:
        if not slide_text.strip():
            continue
            
        # Extract speaker notes
        notes = ""
        if '**Note:**' in slide_text:
            parts = slide_text.split('**Note:**', 1)
            slide_text = parts[0].strip()
            notes = parts[1].strip()
        
        slide = {
            'title': '',
            'subtitle': '',
            'content': [],
            'notes': notes,
            'layout': 'content'  # default
        }
        
        lines = slide_text.strip().split('\n')
        in_code = False
        code_lines = []
        
        for line in lines:
            if not line.strip():
                continue
                
            # Handle code blocks
            if line.strip().startswith('```'):
                in_code = not in_code
                if not in_code and code_lines:
                    slide['content'] = code_lines
                    slide['layout'] = 'code'
                continue
            
            if in_code:
                code_lines.append(line)
                continue
            
            # Parse headers
            if line.startswith('# '):
                title = line[2:].strip()
                slide['title'] = title
                
                # Determine layout based on title
                if 'Block 2: Structural Patterns' in title:
                    slide['layout'] = 'section'
                elif any(pattern in title.lower() for pattern in ['pattern', 'lösung', 'code smells']):
                    slide['layout'] = 'content'
                    
            elif line.startswith('## '):
                if not slide['title']:
                    slide['title'] = line[3:].strip()
                else:
                    slide['subtitle'] = line[3:].strip()
                    
            elif line.startswith('### '):
                # Section within content
                slide['content'].append(line[4:].strip() + ':')
                slide['content'].append('')  # Empty line after section header
                
            elif line.startswith(('• ', '* ', '- ')):
                # Bullet point - remove the bullet character
                bullet = line.lstrip('•*- ').strip()
                slide['content'].append(bullet)
                
            elif line.strip() and not line.startswith('#'):
                # Regular text
                slide['content'].append(line.strip())
        
        if slide['title'] or slide['content'] or slide['notes']:
            slides.append(slide)
    
    return slides

def apply_text_formatting_template_fonts(text_frame, content_lines):
    """Apply text formatting WITHOUT overriding fonts - let template control fonts"""
    
    text_frame.clear()
    
    # First paragraph
    p = text_frame.paragraphs[0]
    p.clear()
    
    for i, line in enumerate(content_lines):
        if i > 0:
            p = text_frame.add_paragraph()
        
        # Check for bold text marked with **
        if '**' in line:
            parts = re.split(r'\*\*(.*?)\*\*', line)
            for j, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    # Every odd index is bold (between **)
                    if j % 2 == 1:
                        run.font.bold = True
                    # DO NOT SET font.name - let template control it
        else:
            p.text = line
            # DO NOT SET font.name - let template control it

def create_block2_presentation_fixed(template_path, slides_data, output_path):
    """Create Block 2 presentation with template font inheritance (NO font overrides)"""
    
    prs = Presentation(template_path)
    
    # Layout indices for VanillaCore template
    LAYOUTS = {
        'title': 0,          # Title Slide
        'section': 1,        # Section Header  
        'content': 3,        # Title and Content (plain)
        'two_column': 5,     # Two Columns (plain)
        'code': 9,           # Code-Block
    }
    
    for slide_data in slides_data:
        # Select layout
        layout_name = slide_data['layout']
        if layout_name not in LAYOUTS:
            layout_name = 'content'
        
        layout_idx = LAYOUTS[layout_name]
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
            # DO NOT set font.name for title - let template control it
        
        # Handle different layouts
        if layout_name == 'section':
            # Section header
            if len(slide.placeholders) > 1:
                if slide_data['subtitle']:
                    slide.placeholders[1].text = slide_data['subtitle']
                elif slide_data['content']:
                    slide.placeholders[1].text = '\n'.join(slide_data['content'])
                # DO NOT set font.name - let template control it
                    
        elif layout_name == 'code':
            # Code block - ONLY here we set font to Consolas
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = '\n'.join(slide_data['content'])
                # Apply monospace font ONLY for code blocks
                text_frame = slide.placeholders[1].text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Consolas'  # ONLY for code blocks
                        run.font.size = Pt(14)
                        
        else:
            # Regular content slides - DO NOT override fonts
            if len(slide.placeholders) > 1 and slide_data['content']:
                # Use template font inheritance
                apply_text_formatting_template_fonts(slide.placeholders[1].text_frame, slide_data['content'])
        
        # Add speaker notes
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']
    
    prs.save(output_path)
    return len(slides_data)

def main():
    """Main function to fix Block 2 fonts"""
    
    # File paths
    markdown_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/block2-content-extracted.md")
    template_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/templates/VanillaCore.pptx")
    output_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/block2-presentation.pptx")
    
    print("=== FIX-BLOCK2-FONT-PRB-019 ===")
    print("Fixing Block 2 font issues - removing Calibri overrides")
    print(f"Reading: {markdown_file}")
    print(f"Template: {template_file}")
    print(f"Output: {output_file}")
    
    # Step 1: Read and clean markdown
    print("\nStep 1: Reading and cleaning markdown...")
    if not markdown_file.exists():
        print(f"ERROR: Markdown file not found: {markdown_file}")
        return
        
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    cleaned_content = clean_markdown(content)
    print(f"✅ Content cleaned, {len(cleaned_content)} characters")
    
    # Step 2: Parse into slide data
    print("\nStep 2: Parsing markdown into slides...")
    slides_data = parse_block2_markdown(cleaned_content)
    print(f"✅ Parsed {len(slides_data)} slides")
    
    # Step 3: Create presentation with template fonts (NO overrides)
    print("\nStep 3: Creating presentation with template font inheritance...")
    if not template_file.exists():
        print(f"ERROR: Template file not found: {template_file}")
        return
        
    slide_count = create_block2_presentation_fixed(template_file, slides_data, output_file)
    print(f"✅ Created presentation with {slide_count} slides")
    
    print("\n=== FONT RULES APPLIED ===")
    print("✅ NO font.name set for regular text (Layouts 0-8)")
    print("✅ Template controls fonts through layout inheritance")  
    print("✅ ONLY Consolas font set for code blocks (Layout 9)")
    print("✅ NO explicit Calibri overrides")
    
    print(f"\n✅ Block 2 presentation fixed and saved to: {output_file}")
    print("✅ Font issues resolved - template now controls fonts properly")

if __name__ == "__main__":
    main()