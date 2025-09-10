#!/usr/bin/env python3
"""
Clean markdown and create professional PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import re

def clean_markdown(content):
    """Remove all HedgeDoc-specific markup and clean the content"""
    
    # Remove YAML frontmatter
    content = re.sub(r'^---.*?---\s*', '', content, flags=re.DOTALL)
    
    # Remove all HTML tags and divs
    content = re.sub(r'<div[^>]*>.*?</div>', '', content, flags=re.DOTALL)
    content = re.sub(r'<[^>]+>', '', content)
    
    # Remove fragment annotations
    content = re.sub(r'<!--.*?-->', '', content)
    
    # Clean up excessive blank lines
    content = re.sub(r'\n{3,}', '\n\n', content)
    
    return content.strip()

def parse_clean_markdown(content):
    """Parse cleaned markdown into slide data"""
    
    slides = []
    raw_slides = content.split('\n---\n')
    
    for slide_text in raw_slides:
        if not slide_text.strip():
            continue
            
        # Extract speaker notes
        notes = ""
        if '\nNote:' in slide_text:
            parts = slide_text.split('\nNote:', 1)
            slide_text = parts[0].strip()
            notes = parts[1].strip()
        
        slide = {
            'title': '',
            'subtitle': '',
            'content': [],
            'notes': notes,
            'layout': 'content'  # default
        }
        
        lines = slide_text.strip().split('\n')
        in_code = False
        code_lines = []
        
        for line in lines:
            if not line.strip():
                continue
                
            # Handle code blocks
            if line.strip().startswith('```'):
                in_code = not in_code
                if not in_code and code_lines:
                    slide['content'] = code_lines
                    slide['layout'] = 'code'
                continue
            
            if in_code:
                code_lines.append(line)
                continue
            
            # Parse headers
            if line.startswith('# '):
                title = line[2:].strip()
                slide['title'] = title
                
                # Determine layout
                if 'Software-Architektur' in title and not slides:
                    slide['layout'] = 'title'
                elif 'Teil' in title:
                    slide['layout'] = 'section'
                    
            elif line.startswith('## '):
                if not slide['title']:
                    slide['title'] = line[3:].strip()
                else:
                    slide['subtitle'] = line[3:].strip()
                    
            elif line.startswith('### '):
                # Section within content
                slide['content'].append(line[4:].strip() + ':')
                slide['content'].append('')  # Empty line after section header
                
            elif line.startswith(('* ', '- ')):
                # Bullet point
                bullet = line[2:].strip()
                slide['content'].append(f'• {bullet}')
                
            else:
                # Regular content
                slide['content'].append(line.strip())
        
        if slide['title']:
            slides.append(slide)
    
    return slides

def apply_text_formatting(text_frame, content_lines):
    """Apply formatting to text including bold"""
    
    # Clear existing paragraphs
    for i in reversed(range(len(text_frame.paragraphs))):
        if i > 0:
            elem = text_frame.paragraphs[i]._element
            elem.getparent().remove(elem)
    
    # First paragraph
    p = text_frame.paragraphs[0]
    p.clear()
    
    for i, line in enumerate(content_lines):
        if i > 0:
            p = text_frame.add_paragraph()
        
        # Check for bold text marked with **
        if '**' in line:
            parts = re.split(r'\*\*(.*?)\*\*', line)
            for j, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    # Every odd index is bold (between **)
                    if j % 2 == 1:
                        run.font.bold = True
        else:
            p.text = line

def create_clean_presentation(template_path, slides_data, output_path):
    """Create PowerPoint presentation with proper formatting"""
    
    prs = Presentation(template_path)
    
    # Layout indices for VanillaCore template
    LAYOUTS = {
        'title': 0,          # Title Slide
        'section': 1,         # Section Header  
        'content': 3,         # Title and Content (plain)
        'two_column': 5,      # Two Columns (plain)
        'code': 9,           # Code-Block
    }
    
    for slide_data in slides_data:
        # Select layout
        layout_name = slide_data['layout']
        if layout_name not in LAYOUTS:
            layout_name = 'content'
        
        layout_idx = LAYOUTS[layout_name]
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
        
        # Handle different layouts
        if layout_name == 'title':
            # Title slide - set subtitle
            if slide_data['subtitle'] and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data['subtitle']
            elif slide_data['content'] and len(slide.placeholders) > 1:
                slide.placeholders[1].text = '\n'.join(slide_data['content'])
                
        elif layout_name == 'section':
            # Section header - subtitle area for learning objectives
            if len(slide.placeholders) > 1:
                if slide_data['subtitle']:
                    slide.placeholders[1].text = slide_data['subtitle']
                elif slide_data['content']:
                    # Join content WITHOUT bullets for section headers
                    clean_content = []
                    for line in slide_data['content']:
                        # Remove bullet characters for section headers
                        if line.startswith('• '):
                            clean_content.append(line[2:])
                        else:
                            clean_content.append(line)
                    slide.placeholders[1].text = '\n'.join(clean_content)
                    
        elif layout_name == 'code':
            # Code block
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = '\n'.join(slide_data['content'])
                # Apply monospace font if possible
                text_frame = slide.placeholders[1].text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Courier New'
                        run.font.size = Pt(14)
                        
        else:
            # Regular content slides
            if len(slide.placeholders) > 1:
                content = slide_data['content']
                
                # Check if we need two columns (Agenda slide with Themenblöcke and Lernziele)
                if 'Themenblöcke:' in content and 'Lernziele:' in content:
                    # Use two column layout
                    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS['two_column']])
                    slide.shapes.title.text = slide_data['title']
                    
                    # Split content at Lernziele
                    idx = content.index('Lernziele:')
                    left_content = content[:idx-1]  # -1 to remove empty line
                    right_content = content[idx:]
                    
                    # Left column
                    if len(slide.placeholders) > 1:
                        apply_text_formatting(slide.placeholders[1].text_frame, left_content)
                    
                    # Right column
                    if len(slide.placeholders) > 2:
                        apply_text_formatting(slide.placeholders[2].text_frame, right_content)
                else:
                    # Single column
                    apply_text_formatting(slide.placeholders[1].text_frame, content)
        
        # Add speaker notes
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']
    
    prs.save(output_path)
    return len(slides_data)

def main():
    """Main function"""
    markdown_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/hedgedoc-intro.md")
    template_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/templates/VanillaCore.pptx")
    output_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/intro-clean.pptx")
    
    print("Step 1: Reading and cleaning markdown...")
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    clean_content = clean_markdown(content)
    
    print("Step 2: Parsing clean markdown...")
    slides = parse_clean_markdown(clean_content)
    print(f"  Found {len(slides)} slides")
    
    print("\nFirst 5 slides preview:")
    for i, slide in enumerate(slides[:5], 1):
        print(f"  Slide {i}: {slide['title'][:40]}... [{slide['layout']}]")
        if slide['notes']:
            print(f"    → Has speaker notes")
    
    print("\nStep 3: Creating PowerPoint presentation...")
    num_slides = create_clean_presentation(template_file, slides, output_file)
    
    print(f"\n✅ Success! Created {output_file}")
    print(f"   {num_slides} slides with proper layouts and formatting")
    print(f"   Speaker notes preserved where present")
    print(f"   Bold text formatting applied")

if __name__ == "__main__":
    main()