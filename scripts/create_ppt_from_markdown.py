#!/usr/bin/env python3
"""
Create PowerPoint presentation from markdown with speaker notes and animations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import re

def parse_markdown_slides(markdown_file):
    """Parse markdown and extract slides with content and speaker notes"""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separators
    slides_raw = content.split('---')
    slides = []
    
    for slide_text in slides_raw:
        slide_text = slide_text.strip()
        if not slide_text or slide_text.startswith('<!--'):
            continue
            
        # Extract speaker notes
        notes = ""
        if 'Note:' in slide_text:
            parts = slide_text.split('Note:')
            slide_text = parts[0].strip()
            notes = parts[1].strip()
        
        # Parse slide content
        lines = slide_text.split('\n')
        slide_data = {
            'title': '',
            'subtitle': '',
            'content': [],
            'notes': notes,
            'type': 'content',
            'bullets': [],
            'is_code': False
        }
        
        in_code_block = False
        code_lines = []
        
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue
                
            # Remove fragment annotations
            line_clean = re.sub(r'<!--.*?-->', '', line_stripped).strip()
            
            if line_clean.startswith('```'):
                # Toggle code block
                in_code_block = not in_code_block
                if in_code_block:
                    slide_data['is_code'] = True
                continue
            
            if in_code_block:
                code_lines.append(line)  # Keep original formatting for code
                continue
            
            if line_clean.startswith('# '):
                # Main title slide (Title page or section divider)
                slide_data['title'] = line_clean[2:]
                # Check if it starts with "Teil" for section headers
                if line_clean.startswith('# Teil'):
                    slide_data['type'] = 'section'
                else:
                    slide_data['type'] = 'title'
            elif line_clean.startswith('## '):
                # Regular slide title
                slide_data['title'] = line_clean[3:]
                slide_data['type'] = 'content'
            elif line_clean.startswith('### '):
                # Subtitle or section
                if slide_data['type'] == 'content' and not slide_data['subtitle']:
                    slide_data['subtitle'] = line_clean[4:]
                else:
                    slide_data['bullets'].append(line_clean[4:])
            elif line_clean.startswith('* ') or line_clean.startswith('- '):
                # Bullet point
                slide_data['bullets'].append(line_clean[2:])
            elif line_clean:
                # Regular content
                slide_data['content'].append(line_clean)
        
        # Add code lines if any
        if code_lines:
            slide_data['content'] = code_lines
            slide_data['type'] = 'code'
        
        if slide_data['title'] or slide_data['content']:
            slides.append(slide_data)
    
    return slides

def create_presentation(template_path, slides_data, output_path):
    """Create PowerPoint from template with slides data"""
    
    # Load template
    prs = Presentation(template_path)
    
    # Get available layouts
    slide_layouts = prs.slide_layouts
    
    # Layout mapping based on VanillaCore template
    TITLE_SLIDE = 0
    SECTION_HEADER = 1
    TITLE_CONTENT_BULLETS = 2
    TITLE_CONTENT = 3
    TWO_COLUMNS_BULLETS = 4
    TWO_COLUMNS = 5
    TITLE_ONLY = 8
    CODE_BLOCK = 9
    
    for slide_data in slides_data:
        slide = None
        
        if slide_data['type'] == 'title':
            # Title slide
            slide = prs.slides.add_slide(slide_layouts[TITLE_SLIDE])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data['title']
            if slide_data['subtitle']:
                subtitle.text = slide_data['subtitle']
            elif slide_data['content']:
                subtitle.text = '\n'.join(slide_data['content'])
                
        elif slide_data['type'] == 'section':
            # Section header for major sections (Teil X)
            slide = prs.slides.add_slide(slide_layouts[SECTION_HEADER])
            title = slide.shapes.title
            
            title.text = slide_data['title']
            # Section headers might have learning objectives in subtitle area
            if slide_data['subtitle']:
                subtitle = slide.placeholders[1]
                subtitle.text = slide_data['subtitle']
            elif slide_data['bullets']:
                # Show as plain text, no bullets for section headers
                subtitle = slide.placeholders[1]
                subtitle.text = '\n'.join(slide_data['bullets'])
                    
        elif slide_data['type'] == 'code':
            # Code block slide
            slide = prs.slides.add_slide(slide_layouts[CODE_BLOCK])
            title = slide.shapes.title
            title.text = slide_data['title'] if slide_data['title'] else 'Code Example'
            
            # Add code to content placeholder
            if len(slide.placeholders) > 1:
                code_placeholder = slide.placeholders[1]
                code_content = '\n'.join(slide_data['content'])
                code_placeholder.text = code_content
                
        else:
            # Regular content slide
            if slide_data['bullets'] and len(slide_data['bullets']) > 7:
                # Use two columns for many bullets
                slide = prs.slides.add_slide(slide_layouts[TWO_COLUMNS])
                title = slide.shapes.title
                title.text = slide_data['title']
                
                # Split bullets between columns
                mid = len(slide_data['bullets']) // 2
                left_bullets = slide_data['bullets'][:mid]
                right_bullets = slide_data['bullets'][mid:]
                
                # Left column
                left_content = slide.placeholders[1]
                left_content.text = '\n'.join(['• ' + b for b in left_bullets])
                
                # Right column  
                right_content = slide.placeholders[2]
                right_content.text = '\n'.join(['• ' + b for b in right_bullets])
                
            elif slide_data['bullets']:
                # Regular bullets slide
                slide = prs.slides.add_slide(slide_layouts[TITLE_CONTENT])
                title = slide.shapes.title
                title.text = slide_data['title']
                
                content = slide.placeholders[1]
                content.text = '\n'.join(['• ' + b for b in slide_data['bullets']])
                
            else:
                # Plain content
                slide = prs.slides.add_slide(slide_layouts[TITLE_CONTENT])
                title = slide.shapes.title
                title.text = slide_data['title']
                
                if slide_data['content']:
                    content = slide.placeholders[1]
                    content.text = '\n'.join(slide_data['content'])
        
        # Add speaker notes if available
        if slide and slide_data['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']
    
    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

def main():
    """Main function"""
    markdown_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/hedgedoc-intro.md")
    template_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/templates/VanillaCore.pptx")
    output_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/intro-design-patterns-python.pptx")
    
    print(f"Parsing markdown from {markdown_file}")
    slides_data = parse_markdown_slides(markdown_file)
    print(f"Found {len(slides_data)} slides")
    
    print(f"Creating presentation with template {template_file}")
    create_presentation(template_file, slides_data, output_file)
    
    print("\nSlide summary:")
    for i, slide in enumerate(slides_data, 1):
        print(f"  Slide {i}: {slide['title'][:50]}... ({slide['type']})")
        if slide['notes']:
            print(f"    → Has speaker notes")

if __name__ == "__main__":
    main()