#!/usr/bin/env python3
"""
Properly create PowerPoint presentation from markdown with correct layout selection
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import re

def parse_markdown_properly(markdown_file):
    """Parse markdown with proper slide type detection"""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separators
    slides_raw = content.split('\n---\n')
    slides = []
    
    for slide_text in slides_raw:
        slide_text = slide_text.strip()
        
        # Skip empty or frontmatter
        if not slide_text or slide_text.startswith('<!--') or slide_text.startswith('type:'):
            continue
        
        # Extract speaker notes
        notes = ""
        if '\nNote:' in slide_text:
            parts = slide_text.split('\nNote:')
            slide_text = parts[0].strip()
            notes = parts[1].strip()
        
        # Remove HTML divs and images
        slide_text = re.sub(r'<div[^>]*>.*?</div>', '', slide_text, flags=re.DOTALL)
        slide_text = re.sub(r'<img[^>]*>', '', slide_text)
        slide_text = re.sub(r'<strong>(.*?)</strong>', r'\1', slide_text)
        slide_text = re.sub(r'<br>', '\n', slide_text)
        
        # Remove fragment annotations
        slide_text = re.sub(r'<!--.*?-->', '', slide_text)
        
        slide_data = {
            'title': '',
            'subtitle': '',
            'content': [],
            'bullets': [],
            'notes': notes,
            'type': 'content',
            'code': [],
            'sections': {}
        }
        
        lines = slide_text.strip().split('\n')
        in_code = False
        current_section = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Code blocks
            if line.startswith('```'):
                in_code = not in_code
                if not in_code and slide_data['code']:
                    slide_data['type'] = 'code'
                continue
            
            if in_code:
                slide_data['code'].append(line)
                continue
            
            # Headers
            if line.startswith('# '):
                title = line[2:].strip()
                slide_data['title'] = title
                # Determine if this is a section header or title slide
                if 'Teil' in title and ':' in title:
                    slide_data['type'] = 'section'
                elif not slides:  # First slide
                    slide_data['type'] = 'title'
                else:
                    slide_data['type'] = 'section'
                    
            elif line.startswith('## '):
                if not slide_data['title']:
                    slide_data['title'] = line[3:].strip()
                else:
                    slide_data['subtitle'] = line[3:].strip()
                    
            elif line.startswith('### '):
                # This is a section header within content
                current_section = line[4:].strip().rstrip(':')
                slide_data['sections'][current_section] = []
                
            elif line.startswith('* ') or line.startswith('- '):
                bullet = line[2:].strip()
                # Remove bold markers
                bullet = re.sub(r'\*\*(.*?)\*\*', r'\1', bullet)
                
                if current_section:
                    slide_data['sections'][current_section].append(bullet)
                else:
                    slide_data['bullets'].append(bullet)
                    
            elif line.startswith(('1. ', '2. ', '3. ', '4. ', '5. ')):
                # Numbered list
                bullet = line[3:].strip()
                slide_data['bullets'].append(bullet)
                
            else:
                # Regular content
                slide_data['content'].append(line)
        
        if slide_data['title'] or slide_data['content']:
            slides.append(slide_data)
    
    return slides

def select_layout(slide_data, has_bullets_layout=True):
    """Select appropriate layout based on content"""
    
    # Map layout indices for VanillaCore template
    LAYOUTS = {
        'title': 0,  # Title Slide
        'section': 1,  # Section Header
        'content_bullets': 2,  # Title and Content with Bullets
        'content_plain': 3,  # Title and Content (no bullets)
        'two_col_bullets': 4,  # Two Columns Bulletpoint
        'two_col_plain': 5,  # Two Columns (no bullets)
        'code': 9,  # Code-Block
    }
    
    if slide_data['type'] == 'title':
        return LAYOUTS['title']
    elif slide_data['type'] == 'section':
        return LAYOUTS['section']
    elif slide_data['type'] == 'code':
        return LAYOUTS['code']
    else:
        # Content slides - choose based on content structure
        total_bullets = len(slide_data['bullets'])
        
        # Add section bullets
        for section_bullets in slide_data['sections'].values():
            total_bullets += len(section_bullets)
        
        # If we have sections (like Themenblöcke and Lernziele)
        if len(slide_data['sections']) == 2:
            # Two column layout
            if total_bullets > 0:
                return LAYOUTS['two_col_plain']  # We'll add bullet chars manually
            else:
                return LAYOUTS['two_col_plain']
        elif total_bullets > 7:
            # Too many bullets for single column
            return LAYOUTS['two_col_plain']
        elif total_bullets > 0:
            # Single column with bullets
            return LAYOUTS['content_plain']  # We'll add bullet chars manually
        else:
            # Plain content
            return LAYOUTS['content_plain']

def create_presentation_proper(template_path, slides_data, output_path):
    """Create PowerPoint with proper layout selection"""
    
    prs = Presentation(template_path)
    slide_layouts = prs.slide_layouts
    
    for slide_data in slides_data:
        layout_idx = select_layout(slide_data)
        slide = prs.slides.add_slide(slide_layouts[layout_idx])
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
        
        # Handle different slide types
        if slide_data['type'] == 'title':
            # Title slide - set subtitle if exists
            if slide_data['subtitle'] and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data['subtitle']
            elif slide_data['content'] and len(slide.placeholders) > 1:
                slide.placeholders[1].text = '\n'.join(slide_data['content'])
                
        elif slide_data['type'] == 'section':
            # Section header - plain text in subtitle area
            if len(slide.placeholders) > 1:
                content = []
                if slide_data['subtitle']:
                    content.append(slide_data['subtitle'])
                if slide_data['bullets']:
                    content.extend(slide_data['bullets'])
                if content:
                    slide.placeholders[1].text = '\n'.join(content)
                    
        elif slide_data['type'] == 'code':
            # Code slide
            if len(slide.placeholders) > 1:
                code_text = '\n'.join(slide_data['code'])
                slide.placeholders[1].text = code_text
                
        else:
            # Content slides
            if len(slide_data['sections']) == 2:
                # Two column layout for sections
                sections = list(slide_data['sections'].items())
                
                # Left column
                if len(slide.placeholders) > 1:
                    left_text = sections[0][0] + ':\n\n'
                    for bullet in sections[0][1]:
                        left_text += f'• {bullet}\n'
                    slide.placeholders[1].text = left_text.strip()
                
                # Right column  
                if len(slide.placeholders) > 2:
                    right_text = sections[1][0] + ':\n\n'
                    for bullet in sections[1][1]:
                        right_text += f'• {bullet}\n'
                    slide.placeholders[2].text = right_text.strip()
                    
            elif slide_data['bullets']:
                # Single column with bullets
                if len(slide.placeholders) > 1:
                    # Check if we need to split due to too many bullets
                    if len(slide_data['bullets']) > 7 and layout_idx in [5, 4]:  # Two column layout
                        mid = len(slide_data['bullets']) // 2
                        left_bullets = slide_data['bullets'][:mid]
                        right_bullets = slide_data['bullets'][mid:]
                        
                        # Left column
                        left_text = '\n'.join([f'• {b}' for b in left_bullets])
                        slide.placeholders[1].text = left_text
                        
                        # Right column
                        if len(slide.placeholders) > 2:
                            right_text = '\n'.join([f'• {b}' for b in right_bullets])
                            slide.placeholders[2].text = right_text
                    else:
                        # Single column
                        content = []
                        
                        # Add any content first
                        if slide_data['content']:
                            content.extend(slide_data['content'])
                            content.append('')  # Empty line
                        
                        # Add subtitle if exists
                        if slide_data['subtitle']:
                            content.append(slide_data['subtitle'])
                            content.append('')
                        
                        # Add sections with bullets
                        for section_name, section_bullets in slide_data['sections'].items():
                            if section_name:
                                content.append(section_name + ':')
                            for bullet in section_bullets:
                                content.append(f'• {bullet}')
                            if section_bullets:
                                content.append('')
                        
                        # Add regular bullets
                        for bullet in slide_data['bullets']:
                            content.append(f'• {bullet}')
                        
                        slide.placeholders[1].text = '\n'.join(content).strip()
            else:
                # Plain content
                if len(slide.placeholders) > 1 and slide_data['content']:
                    slide.placeholders[1].text = '\n'.join(slide_data['content'])
        
        # Add speaker notes
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']
    
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

def main():
    """Main function"""
    markdown_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/hedgedoc-intro.md")
    template_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/templates/VanillaCore.pptx")
    output_file = Path("/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/intro-design-patterns-final.pptx")
    
    print(f"Parsing markdown from {markdown_file}")
    slides_data = parse_markdown_properly(markdown_file)
    print(f"Found {len(slides_data)} slides")
    
    # Debug first few slides
    for i, slide in enumerate(slides_data[:5], 1):
        print(f"\nSlide {i}:")
        print(f"  Type: {slide['type']}")
        print(f"  Title: {slide['title']}")
        if slide['sections']:
            print(f"  Sections: {list(slide['sections'].keys())}")
        if slide['notes']:
            print(f"  Has speaker notes: Yes")
    
    print(f"\nCreating presentation with template {template_file}")
    create_presentation_proper(template_file, slides_data, output_file)

if __name__ == "__main__":
    main()