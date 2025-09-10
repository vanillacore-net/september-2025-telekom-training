#!/usr/bin/env python3
"""
PROPERLY Fix Block 2 Using VanillaCore Template and Python
PRB: PROPERLY-USE-TEMPLATE-BLOCK2-PRB-021

CRITICAL REQUIREMENTS:
- MUST actually use templates/VanillaCore.pptx as base template
- MUST use template's slide layouts properly
- MUST show REAL Java code (not comments)
- MUST remove ** formatting artifacts 
- MUST remove "Slide X:" from titles
- MUST add speaker notes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

def create_block2_presentation():
    """Create Block 2 presentation using the actual VanillaCore template"""
    
    print("Loading VanillaCore.pptx template...")
    # CRITICAL: LOAD THE ACTUAL TEMPLATE - NOT CREATE NEW!
    prs = Presentation('templates/VanillaCore.pptx')
    
    # Use the template's slide_layouts
    slide_layout_0 = prs.slide_layouts[0]  # Title
    slide_layout_1 = prs.slide_layouts[1]  # Section Header  
    slide_layout_2 = prs.slide_layouts[2]  # Bullets
    slide_layout_3 = prs.slide_layouts[3]  # Plain
    slide_layout_9 = prs.slide_layouts[9]  # Code
    
    print("Creating slides using template layouts...")
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(slide_layout_0)
    slide.shapes.title.text = "Block 2: Structural Patterns"
    slide.shapes.placeholders[1].text = "Entwurfsmuster für Objekt-Zusammensetzung und System-Architektur"
    slide.notes_slide.notes_text_frame.text = "Block 2 behandelt Structural Patterns - Entwurfsmuster, die sich mit der Komposition von Objekten und Klassen beschäftigen. Diese Patterns lösen Probleme der Objekt-Zusammensetzung und System-Architektur."

    # DECORATOR PATTERN SECTION
    
    # Slide 2: Decorator - Section Header
    slide = prs.slides.add_slide(slide_layout_1)
    slide.shapes.title.text = "Decorator Pattern"
    slide.notes_slide.notes_text_frame.text = "Decorator Pattern löst das Problem der exponentiellen Klassen-Explosion durch Komposition statt Vererbung."

    # Slide 3: Decorator - Problem (Code)
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Was ist hier schlecht?"
    
    # Add code content to the code placeholder
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    # Set font to Consolas for code
    p = code_frame.paragraphs[0]
    p.text = """// Exponentielles Wachstum-Problem
BasicTariff
BasicTariffWithData
BasicTariffWithDataAndRoaming  
// Bei n Features = 2^n Klassen!

// Vererbungs-Explosion
public class BasicTariffWithDataAndRoamingAndHotspot 
    extends BasicTariffWithDataAndRoaming {
    // Bei nur 10 Tarif-Optionen entstehen 1.024 Klassen!
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(14)
    
    slide.notes_slide.notes_text_frame.text = "Die Deutsche Telekom bietet verschiedenste Tarif-Optionen: Datenvolumen-Upgrades, Roaming-Pakete, Streaming-Services, Hotspot-Zugang. Der naive OOP-Ansatz führt zu exponentieller Klassenzahl."

    # Slide 4: Decorator - Code Smells
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Code Smells identifiziert"
    
    bullet_points = [
        "Klassen-Explosion: Bei n Features entstehen 2^n Klassen",
        "Code-Duplikation: Bugfixes müssen in hunderten Klassen repliziert werden",
        "Open/Closed Verletzung: Jede neue Option erfordert Dutzende neuer Klassen",
        "Wartungs-Alptraum: Änderungen propagieren durch die gesamte Hierarchie",
        "Kombinatorisches Problem: Unmöglich alle Feature-Kombinationen zu modellieren",
        "Tight Coupling: Features sind fest mit Basis-Implementierung gekoppelt",
        "Inflexibilität: Keine zur-Laufzeit-Konfiguration von Features möglich"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Das mathematische Problem: Bei nur 10 Tarif-Optionen entstehen 1.024 Klassen! Dies verletzt fundamental das Open/Closed Principle."

    # Slide 5: Decorator - Solution
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Lösung: Decorator Pattern"
    
    solution_points = [
        "Komposition über Vererbung: Objekte zur Laufzeit \"umhüllen\"",
        "Gemeinsame Schnittstelle: Component-Interface für einheitliche Behandlung",
        "Transparente Erweiterung: Decorators sind für Clients unsichtbar",
        "Beliebige Kombinierbarkeit: Features können flexibel kombiniert werden",
        "Einzelne Verantwortung: Jeder Decorator fügt genau ein Feature hinzu",
        "Zur-Laufzeit-Konfiguration: Dynamische Objekt-Zusammensetzung",
        "Template Method Integration: Abstraktes Decorator-Rückgrat für Standards"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(solution_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Decorator Pattern löst das Problem durch Komposition statt Vererbung. Ein Tarif-Objekt wird von \"Decorator\"-Objekten umschlossen, die jeweils eine zusätzliche Funktion bereitstellen."

    # Slide 6: Decorator - Implementation
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Implementierung"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// Gemeinsame Schnittstelle
public interface TariffService {
    BigDecimal calculatePrice();
    Set<String> getFeatures(); 
}

// Abstrakter Decorator
public abstract class TariffDecorator implements TariffService {
    protected final TariffService decoratedTariff;
    
    public BigDecimal calculatePrice() {
        return decoratedTariff.calculatePrice(); // Delegation
    }
}

// Konkreter Decorator  
public class DataOptionDecorator extends TariffDecorator {
    private final BigDecimal dataPrice;
    
    @Override
    public BigDecimal calculatePrice() {
        return super.calculatePrice().add(dataPrice); // Erweiterung
    }
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
    
    slide.notes_slide.notes_text_frame.text = "Jeder Decorator implementiert dieselbe Schnittstelle wie das umhüllte Objekt - so können Decorators beliebig geschachtelt werden. Factory Pattern übernimmt die Objektkomposition und versteckt die Komplexität."

    # COMPOSITE PATTERN SECTION
    
    # Slide 7: Composite - Section Header
    slide = prs.slides.add_slide(slide_layout_1)
    slide.shapes.title.text = "Composite Pattern"
    slide.notes_slide.notes_text_frame.text = "Composite Pattern behandelt Einzelobjekte und Objektsammlungen einheitlich."

    # Slide 8: Composite - Problem
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Was ist hier schlecht?"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// Separate Behandlung für Einzeltarife vs. Gruppen
public class TariffService {
    public void processOrder(Object tariff) {
        if (tariff instanceof IndividualTariff) {
            IndividualTariff individual = (IndividualTariff) tariff;
            // Individuelle Verarbeitung
        } else if (tariff instanceof FamilyTariff) {
            FamilyTariff family = (FamilyTariff) tariff;
            // Familiengruppen-spezifische Logik
            for (IndividualTariff member : family.getMembers()) {
                // Rekursive Verarbeitung...
            }
        }
        // Komplexe If-Verschachtelung für alle Tarif-Typen
    }
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
    
    slide.notes_slide.notes_text_frame.text = "Telekom verkauft komplexe Familien-Strukturen mit verschachtelten Abhängigkeiten. Separate Klassen für jede Konstellation führen zu kombinatorischer Explosion."

    # Slide 9: Composite - Code Smells
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Code Smells identifiziert"
    
    composite_smells = [
        "Type-Checking Nightmare: Instanceof-Checks überall im Code",
        "Rekursive Logik-Duplikation: Gleiche Traversierung in jeder Methode",
        "Kombinatorische Explosion: Neue Tarif-Typen = neue If-Zweige überall",
        "Client-Komplexität: Clients müssen alle Tarif-Typen kennen",
        "Tight Coupling: Client-Code gekoppelt an konkrete Tarif-Typen",
        "Inkonsistente Behandlung: Verschiedene Logik für ähnliche Operationen",
        "Schwer erweiterbar: Neue Hierarchie-Ebenen = massive Code-Änderungen"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(composite_smells):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Das Komplexitäts-Dilemma: Wie behandeln wir Einzeltarife und Gruppen-Tarife einheitlich? Naive Ansätze scheitern an kombinatorischer Explosion."

    # Slide 10: Composite - Solution
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Lösung: Composite Pattern"
    
    composite_solution = [
        "Einheitliche Behandlung: Einzelobjekte und Sammlungen identisch behandeln",
        "Rekursive Struktur: Composites können andere Composites enthalten",
        "Transparenz für Client: Client arbeitet nur mit Component-Interface",
        "Natürliche Hierarchie-Abbildung: Teil-Ganzes-Beziehungen elegant modelliert",
        "Liskov Substitution: Blätter und Äste vollständig austauschbar",
        "Funktionale Aggregation: Stream-API macht Traversierung elegant",
        "Strategy Pattern Integration: Gruppen-Rabatte bleiben austauschbar"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(composite_solution):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Composite Pattern behandelt Einzelobjekte und Objektsammlungen einheitlich. Ein Familien-Tarif wird genauso behandelt wie ein Einzel-Tarif."

    # Slide 11: Composite - Implementation
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Implementierung"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// Component-Interface
public interface TariffComponent {
    BigDecimal calculateTotalPrice();
    List<String> getAllFeatures();
    int getContractCount();
}

// Composite-Implementierung
public class TariffGroup implements TariffComponent {
    private final List<TariffComponent> children = new ArrayList<>();
    
    public BigDecimal calculateTotalPrice() {
        BigDecimal total = children.stream()
            .map(TariffComponent::calculateTotalPrice)
            .reduce(BigDecimal.ZERO, BigDecimal::add);
        return discountStrategy.applyGroupDiscount(total);
    }
    
    public int getContractCount() {
        return children.stream()
            .mapToInt(TariffComponent::getContractCount)
            .sum(); // Rekursive Aggregation
    }
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(11)
    
    slide.notes_slide.notes_text_frame.text = "Transparente Rekursion: TariffGroup \"weiß\" nicht, ob Kinder Leafs oder Composites sind. Visitor Pattern kann für Performance-Optimierung bei tiefen Hierarchien eingesetzt werden."

    # FACADE PATTERN SECTION
    
    # Slide 12: Facade - Section Header
    slide = prs.slides.add_slide(slide_layout_1)
    slide.shapes.title.text = "Facade Pattern"
    slide.notes_slide.notes_text_frame.text = "Facade Pattern erstellt eine \"einfache\" Schnittstelle vor einem \"komplexen\" Subsystem."

    # Slide 13: Facade - Problem
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Was ist hier schlecht?"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// Client-Albtraum: 15+ Service-Abhängigkeiten
public class TariffOrderClient {
    private TariffValidationService validation;
    private CustomerCreditCheckService creditCheck;
    private PricingService pricing;
    private InventoryService inventory;
    private NotificationService notification;
    // ... 10 weitere Services
    
    public void createOrder(TariffRequest request) {
        // 50+ Zeilen Orchestrierung
        if (validation.check() && creditCheck.approve() && 
            pricing.calculate() && inventory.reserve()) {
            // Komplexe If-Verschachtelung
        }
    }
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
    
    slide.notes_slide.notes_text_frame.text = "Eine einfache Tarif-Bestellung erfordert Koordination von 15+ Microservices. Client-Code wird unlesesbar und untestbar."

    # Slide 14: Facade - Code Smells
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Code Smells identifiziert"
    
    facade_smells = [
        "Abhängigkeits-Explosion: Client kennt 15+ Services intimement",
        "Orchestrierungs-Komplexität: 50+ Zeilen if/else-Logik pro Use Case",
        "Fehlerbehandlungs-Chaos: Exponentiell wachsende Fehlerpfade",
        "Testing-Alptraum: Mocking von 15+ Services pro Unit Test",
        "Deployment-Fragilität: Client bricht bei jeder Service-Änderung",
        "Code-Duplikation: Gleiche Orchestrierung in mehreren Clients",
        "Service-Interface-Leakage: Interne Service-Details überall sichtbar"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(facade_smells):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Das Komplexitäts-Dilemma: Jede Service-Änderung bricht den Client. Deployment-Abhängigkeiten werden unbeherrschbar."

    # Slide 15: Facade - Solution
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Lösung: Facade Pattern"
    
    facade_solution = [
        "Einheitliche Schnittstelle: Eine einfache API vor komplexem Subsystem",
        "Komplexitäts-Kapselung: Interne Orchestrierung für Clients unsichtbar",
        "Use-Case-spezifische APIs: APIs sprechen die Sprache der Business-User",
        "Zentrale Fehlerbehandlung: Standardisierte Error-Responses",
        "Versionierungs-Puffer: API-Änderungen ohne Client-Brüche möglich",
        "Performance-Optimierung: Parallelisierung und Caching zentral möglich",
        "Service-Abstraktion: Services können ausgetauscht werden ohne Client-Impact"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(facade_solution):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Facade Pattern erstellt eine \"einfache\" Schnittstelle vor einem \"komplexen\" Subsystem. Client-Sicht: 1 Methode statt 15 Service-Calls."

    # Slide 16: Facade - Implementation
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Implementierung"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """@Service
public class TariffOrderFacade {
    // Alle komplexen Dependencies sind privat gekapselt
    private final TariffValidationService validation;
    private final CreditCheckService creditCheck;
    // ... weitere Services (privat!)
    
    public TariffOrderResult createTariffOrder(TariffOrderRequest request) {
        try {
            if (!validateOrder(request).isValid()) {
                return TariffOrderResult.failed("Validation failed");
            }
            return executeOrderProcess(request); // Details versteckt
        } catch (Exception e) {
            return TariffOrderResult.failed(e.getMessage());
        }
    }
    
    // Parallele Validierungen für Performance
    private CompletableFuture<ValidationResult> validateOrderAsync(request) {
        // Unabhängige Service-Calls parallel ausführen
    }
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(10)
    
    slide.notes_slide.notes_text_frame.text = "Facade vs. Microservice Gateway: Façade ist ein architektonisches Pattern, Gateway ist infrastrukturell. Circuit Breaker, Timeout Management und Parallelisierung können integriert werden."

    # PROXY PATTERN SECTION
    
    # Slide 17: Proxy - Section Header
    slide = prs.slides.add_slide(slide_layout_1)
    slide.shapes.title.text = "Proxy Pattern"
    slide.notes_slide.notes_text_frame.text = "Proxy Pattern stellt einen \"Stellvertreter\" vor ein teures Objekt."

    # Slide 18: Proxy - Problem
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Was ist hier schlecht?"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// Synchrone, ungecachte DB-Zugriffe
public class CustomerService {
    public CustomerData getCustomer(String id) {
        // Legacy-DB: 200ms+ pro Query!
        return legacyDatabase.findById(id);
    }
}

// Bei 1000 Requests/sec:
// 1000 × 200ms = 200 Sekunden DB-Zeit pro Sekunde
// System bricht zusammen!

// Memory-Verschwendung: Jeder Customer 2KB
// 50 Million Kunden × 2KB = 100GB RAM nur für Basis-Kundendaten"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
    
    slide.notes_slide.notes_text_frame.text = "Legacy-Datenbanken sind oft der Flaschenhals moderner Systeme. Hochfrequente Zugriffe führen zu Cascading Failures."

    # Slide 19: Proxy - Code Smells
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Code Smells identifiziert"
    
    proxy_smells = [
        "Latenz-Explosion: 200ms+ Response-Zeit pro DB-Query",
        "Memory-Verschwendung: Redundante Daten millionenfach geladen",
        "Cascade-Failure-Risk: Langsame Services blockieren nachgelagerte Systeme",
        "No Caching Strategy: Identische Queries werden wiederholt ausgeführt",
        "Synchronous Blocking: Alle Threads warten auf DB-Response",
        "Resource Exhaustion: Connection Pools werden erschöpft",
        "Poor User Experience: Web-Interfaces werden unbenutzbar"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(proxy_smells):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Die Latenz-Spirale: 1000+ Customer-Abfragen pro Sekunde multipliziert mit 200ms = System-Kollaps."

    # Slide 20: Proxy - Solution
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Lösung: Proxy Pattern"
    
    proxy_solution = [
        "Intelligente Zugriffssteuerung: Proxy kontrolliert Zugriff auf teure Objekte",
        "Transparente Optimierung: Client merkt nichts von Performance-Verbesserungen",
        "Lazy Loading: Objekte nur bei tatsächlichem Bedarf laden",
        "Caching-Integration: Intelligente Zwischenspeicherung mit konfigurierbaren Policies",
        "Security-Layer: Transparent implementierte Sicherheitsprüfungen",
        "Remote Proxy: Netzwerk-Komplexität verstecken mit Resilience Patterns",
        "Monitoring-Integration: Zentrale Metriken für alle Service-Calls"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(proxy_solution):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Proxy Pattern stellt einen \"Stellvertreter\" vor ein teures Objekt. Der Proxy kann Caching, Lazy Loading, Security etc. transparent hinzufügen."

    # Slide 21: Proxy - Implementation
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Implementierung"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// Virtual Proxy für Lazy Loading
public class LazyCustomerServiceProxy implements CustomerService {
    private final CustomerService realService;
    private final Map<String, CustomerData> loadedCustomers = new HashMap<>();
    
    public CustomerData getCustomer(String customerId) {
        return loadedCustomers.computeIfAbsent(customerId, 
            id -> realService.getCustomer(id)); // Lazy Loading
    }
}

// Caching Proxy für Performance
@Component
public class CachedCustomerServiceProxy implements CustomerService {
    private final Cache<String, CustomerData> customerCache;
    
    public CustomerData getCustomer(String customerId) {
        return customerCache.get(customerId, id -> {
            return realService.getCustomer(id); // Cache Miss
        });
    }
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(10)
    
    slide.notes_slide.notes_text_frame.text = "Cache-Hit-Rate von 99.5%+ ist realistisch = 2000x Performance-Verbesserung. Security Proxy kann Authorization und Audit-Logging transparent implementieren."

    # FLYWEIGHT PATTERN SECTION
    
    # Slide 22: Flyweight - Section Header
    slide = prs.slides.add_slide(slide_layout_1)
    slide.shapes.title.text = "Flyweight Pattern"
    slide.notes_slide.notes_text_frame.text = "Flyweight Pattern trennt \"geteilte\" von \"individuellen\" Daten."

    # Slide 23: Flyweight - Problem
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Was ist hier schlecht?"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// REDUNDANT: Gleiche Tarif-Daten millionenfach dupliziert
public class CustomerData {
    private String customerId;          // INDIVIDUELL
    private String name;                // INDIVIDUELL
    
    // Nur ~50 verschiedene Werte für Millionen Kunden!
    private String tariffPlanName;      // "MagentaMobil L" - 1M× dupliziert!
    private String tariffDescription;   // Beschreibung - 1M× dupliziert!
    private BigDecimal tariffPrice;     // 49.99 - 1M× dupliziert!
    private Set<String> tariffFeatures; // Feature-Liste - 1M× dupliziert!
}

// Memory-Verschwendung: 500MB für identische Strings!"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
    
    slide.notes_slide.notes_text_frame.text = "Mit 50+ Millionen Kunden entstehen massive Datenmengen. Redundante Tarif-Beschreibungen werden millionenfach dupliziert = Gigabytes verschwendeter Speicher."

    # Slide 24: Flyweight - Code Smells
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Code Smells identifiziert"
    
    flyweight_smells = [
        "Memory-Explosion: 50 verschiedene Tarife × 1M Kunden = 50M redundante Kopien",
        "Gigabyte-Verschwendung: Identische Strings millionenfach im Memory",
        "Skalierungs-Problem: Memory-Verbrauch explodiert mit Kundenzahl",
        "Cloud-Kosten-Explosion: Ineffiziente Memory-Nutzung = hohe Infrastruktur-Kosten",
        "Cache-Ineffizienz: Redundante Daten verdrängen wichtige Daten aus Caches",
        "GC-Pressure: Garbage Collector überlastet durch redundante Objekte",
        "Network-Overhead: Redundante Daten werden über Netzwerk übertragen"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(flyweight_smells):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Mathematisches Dilemma: 50 Millionen Kunden × 2KB = 100GB RAM nur für Customer-Objekte. Redundante Daten vervielfachen den Speicher-Bedarf exponentiell."

    # Slide 25: Flyweight - Solution
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Lösung: Flyweight Pattern"
    
    flyweight_solution = [
        "Intrinsic vs. Extrinsic State: Geteilte von individuellen Daten trennen",
        "Memory-Sharing: Identische Objekte nur einmal im Memory speichern",
        "Factory-Management: Zentrale Verwaltung garantiert Singleton-Eigenschaften",
        "Immutable Flyweights: Thread-Safety durch Unveränderlichkeit",
        "Context Pattern: Individuelle Daten als Parameter-Objekte",
        "70%+ Memory-Ersparnis: Dramatische Reduktion des Speicherverbrauchs",
        "Lineare Skalierung: Memory-Verbrauch wächst linear mit Kundenzahl"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(flyweight_solution):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Flyweight Pattern trennt \"geteilte\" von \"individuellen\" Daten. Objekte mit identischen Properties werden nur einmal gespeichert."

    # Slide 26: Flyweight - Implementation
    slide = prs.slides.add_slide(slide_layout_9)
    slide.shapes.title.text = "Implementierung"
    
    code_placeholder = slide.shapes.placeholders[1]
    code_frame = code_placeholder.text_frame
    code_frame.clear()
    
    p = code_frame.paragraphs[0]
    p.text = """// Flyweight: Nur intrinsic (geteilte) Properties
public class TariffPlanFlyweight {
    private final String planName;          // Intrinsic - geteilt
    private final String description;       // Intrinsic - geteilt  
    private final BigDecimal basePrice;     // Intrinsic - geteilt
    private final Set<String> features;     // Intrinsic - geteilt
    
    // Business-Operation mit extrinsic context
    public CustomerBill calculateBill(CustomerContext context) {
        BigDecimal finalPrice = basePrice;
        if (context.hasLoyaltyDiscount()) {
            finalPrice = finalPrice.multiply(BigDecimal.valueOf(0.9));
        }
        return new CustomerBill(context.getCustomerId(), planName, finalPrice);
    }
}

// Factory für Flyweight-Management
@Component
public class TariffPlanFlyweightFactory {
    private final Map<String, TariffPlanFlyweight> flyweights = new ConcurrentHashMap<>();
    
    public TariffPlanFlyweight getTariffPlan(String planName) {
        return flyweights.computeIfAbsent(planName, this::createTariffPlan);
    }
}"""
    
    for paragraph in code_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(9)
    
    slide.notes_slide.notes_text_frame.text = "Messbare Ergebnisse: 77% Memory-Ersparnis! Von 103GB auf 24GB RAM-Verbrauch. Nur 50 Flyweight-Objekte statt 50 Millionen redundante Tarif-Kopien."

    # INTEGRATION & SUMMARY
    
    # Slide 27: Pattern-Integration & Enterprise-Readiness
    slide = prs.slides.add_slide(slide_layout_1)
    slide.shapes.title.text = "Pattern-Integration & Enterprise-Readiness"
    slide.notes_slide.notes_text_frame.text = "Die Patterns ergänzen sich perfekt: Komplexe Business-Domänen werden durch Pattern-Kombination beherrschbar."

    # Slide 28: Pattern-Synergien
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Pattern-Synergien"
    
    synergies = [
        "Decorator + Strategy: Optionale, austauschbare Algorithmen kombinieren",
        "Composite + Visitor: Performance-optimierte Hierarchie-Traversierung",
        "Facade + Circuit Breaker: Resiliente Service-Orchestrierung",
        "Proxy + Flyweight: Memory-effiziente Performance-Optimierung",
        "Bridge + Factory: Dynamische Provider-Auswahl mit Abstraktion",
        "Template Method Integration: Konsistente Flows über alle Patterns",
        "Event-Driven Enhancement: Eventual Consistency für komplexe Systeme"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(synergies):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Die Patterns ergänzen sich perfekt: Komplexe Business-Domänen werden durch Pattern-Kombination beherrschbar. Factory + Decorator lösen gemeinsam das Problem der flexiblen Objektzusammensetzung."

    # Slide 29: Production-Ready Considerations
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Production-Ready Considerations"
    
    production_ready = [
        "Distributed Caching: Redis für Multi-Instance Cache-Konsistenz",
        "Circuit Breaker: Resilience bei Provider-Integration",
        "Health Monitoring: Proaktive Performance- und Verfügbarkeits-Überwachung",
        "Cache Warming: Predictive Loading für kritische Daten",
        "Graceful Degradation: Fallback-Strategien für Provider-Ausfälle",
        "Memory Monitoring: JVM-Metriken und Alerting für Memory-Patterns",
        "Event Sourcing: Audit-Trail und System-Konsistenz"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(production_ready):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Enterprise-Tauglichkeit erfordert mehr als Design Patterns: Monitoring, Resilience, Distributed Systems Patterns und observability sind kritisch für Production-Einsatz."

    # Slide 30: Summary
    slide = prs.slides.add_slide(slide_layout_2)
    slide.shapes.title.text = "Zusammenfassung: Block 2 Erkenntnisse"
    
    summary_points = [
        "Strukturelle Probleme elegant lösen: Patterns adressieren fundamentale Kompositions-Herausforderungen",
        "Performance durch Design: Memory-Optimierung und Latenz-Reduktion als architektonische Entscheidung",
        "Enterprise-Skalierung: Patterns skalieren mit bewusster Implementierung bis zu 50M+ Kunden",
        "Microservice-Integration: Loose Coupling und Provider-Abstraktion für resiliente Systeme",
        "Pattern-Kombination: Synergien zwischen Patterns multiplizieren architektonischen Wert",
        "Production-Readiness: Mit Monitoring, Caching und Resilience enterprise-tauglich"
    ]
    
    content_placeholder = slide.shapes.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(summary_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    slide.notes_slide.notes_text_frame.text = "Block 2 zeigt: Structural Patterns sind nicht nur akademische Konzepte, sondern lösen reale Enterprise-Performance und -Wartungsprobleme. Die Kombination mehrerer Patterns ist oft mächtiger als einzelne Pattern-Anwendung."
    
    print("Saving presentation...")
    prs.save('presentations/powerpoint/block2-presentation.pptx')
    print("✅ Block 2 presentation created successfully using VanillaCore template!")
    print("📁 Saved as: presentations/powerpoint/block2-presentation.pptx")
    
    return True

if __name__ == "__main__":
    try:
        success = create_block2_presentation()
        if success:
            print("\n🎯 PRB PROPERLY-USE-TEMPLATE-BLOCK2-PRB-021 completed successfully!")
            print("✅ Used actual VanillaCore.pptx template")
            print("✅ Added real Java code examples")
            print("✅ Removed formatting artifacts")
            print("✅ Added comprehensive speaker notes")
            print("✅ Professional PowerPoint output generated")
        else:
            print("❌ Failed to create Block 2 presentation")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()