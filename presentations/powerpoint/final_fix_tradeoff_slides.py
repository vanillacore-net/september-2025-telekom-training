#!/usr/bin/env python3
"""
ULTIMATE FIX: Create trade-off slides with manual textbox layout
Since Layout 4 doesn't have the expected structure, we'll create slides with custom textboxes
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def get_tradeoffs_data():
    """Complete trade-offs data for all patterns"""
    return {
        'Layered Architecture': {
            'vorteile': [
                'Klare Trennung der Verantwortlichkeiten',
                'Einfache Testbarkeit jeder Schicht',
                'Bewährtes Pattern mit hoher Entwickler-Akzeptanz',
                'Gute Performance bei einfachen CRUD-Operationen'
            ],
            'nachteile': [
                'Monolithischer Charakter erschwert Skalierung',
                'Änderungen propagieren durch alle Schichten',
                'Datenbankschema-Changes beeinflussen alle Layer',
                'Schwierig für komplexe Domain Logic'
            ]
        },
        'Microservices': {
            'vorteile': [
                'Unabhängige Skalierung pro Service',
                'Technology Stack Diversität möglich',
                'Fehler-Isolierung zwischen Services',
                'Parallele Entwicklung durch verschiedene Teams'
            ],
            'nachteile': [
                'Hohe operationale Komplexität',
                'Distributed System Challenges (Latency, Partial Failures)',
                'Data Consistency zwischen Services schwierig',
                'Service Discovery und Load Balancing erforderlich'
            ]
        },
        'Event-Driven': {
            'vorteile': [
                'Lose Kopplung zwischen Komponenten',
                'Asynchrone Verarbeitung für bessere Performance',
                'Event Replay für Debugging und Analytics möglich',
                'Einfaches Hinzufügen neuer Event Consumer'
            ],
            'nachteile': [
                'Eventually Consistent Data',
                'Komplexere Error Handling (Dead Letter Queues)',
                'Event Schema Evolution herausfordernd',
                'Debugging von Event Flows schwierig'
            ]
        },
        'Hexagonal Architecture': {
            'vorteile': [
                'Domain Logic vollständig isoliert von Infrastructure',
                'Austauschbare External Dependencies (Database, APIs)',
                'Excellente Testbarkeit durch Dependency Inversion',
                'Framework-agnostic Domain Layer'
            ],
            'nachteile': [
                'Initial Setup komplexer als Layer Architecture',
                'Mehr Abstractions und Interfaces',
                'Kann Over-Engineering für einfache CRUD Apps sein',
                'Team muss Dependency Inversion verstehen'
            ]
        },
        'CQRS': {
            'vorteile': [
                'Optimierte Read und Write Models',
                'Unabhängige Skalierung von Queries und Commands',
                'Bessere Performance durch spezialisierte Datenmodelle',
                'Event Sourcing Integration möglich'
            ],
            'nachteile': [
                'Erhöhte Komplexität der Architektur',
                'Eventual Consistency zwischen Read/Write',
                'Synchronisation der Models erforderlich',
                'Mehr Code und Infrastruktur'
            ]
        },
        'Event Sourcing': {
            'vorteile': [
                'Vollständige Audit Trail automatisch',
                'Time Travel und Event Replay möglich',
                'Debugging durch Event History',
                'Integration mit CQRS natural'
            ],
            'nachteile': [
                'Event Store Management komplex',
                'Schema Evolution schwierig',
                'Snapshot Management erforderlich',
                'GDPR Compliance herausfordernd'
            ]
        },
        'Circuit Breaker': {
            'vorteile': [
                'Schnelles Fail-Fast Verhalten',
                'System Stabilität bei Ausfällen',
                'Automatische Recovery Mechanismen',
                'Cascading Failures Prevention'
            ],
            'nachteile': [
                'Threshold Tuning erforderlich',
                'False Positives möglich',
                'Additional Monitoring Layer',
                'Complexity in Error Handling'
            ]
        },
        'Saga Pattern': {
            'vorteile': [
                'Distributed Transaction Management',
                'Compensation Logic für Rollbacks',
                'Keine 2PC Lock Requirements',
                'Bessere Performance als 2PC'
            ],
            'nachteile': [
                'Komplexe Compensation Logic',
                'Debugging schwierig',
                'Partial Failure Handling',
                'Testing der Compensation'
            ]
        },
        'API Gateway': {
            'vorteile': [
                'Single Entry Point für Clients',
                'Cross-Cutting Concerns zentral',
                'API Composition möglich',
                'Backend Service Abstraktion'
            ],
            'nachteile': [
                'Single Point of Failure',
                'Additional Latency Layer',
                'Deployment Koordination',
                'Gateway Bottleneck möglich'
            ]
        },
        'Service Mesh': {
            'vorteile': [
                'Service-to-Service Security',
                'Automatic Load Balancing',
                'Observability out-of-the-box',
                'Traffic Management Features'
            ],
            'nachteile': [
                'Operational Overhead',
                'Resource Consumption (Sidecars)',
                'Learning Curve hoch',
                'Debugging Complexity'
            ]
        },
        'Bulkhead Pattern': {
            'vorteile': [
                'Fault Isolation zwischen Components',
                'Resource Protection',
                'Predictable Performance',
                'Prevents Resource Exhaustion'
            ],
            'nachteile': [
                'Resource Underutilization möglich',
                'Configuration Complexity',
                'More Infrastructure Required',
                'Monitoring Overhead'
            ]
        },
        'Domain-Driven Design': {
            'vorteile': [
                'Klare Bounded Contexts',
                'Ubiquitous Language etabliert',
                'Business-IT Alignment',
                'Modulare Struktur'
            ],
            'nachteile': [
                'Steile Learning Curve',
                'Upfront Design Investment',
                'Domain Expert Availability',
                'Over-Engineering Gefahr'
            ]
        }
    }

def create_tradeoff_slide_with_textboxes(prs, template_prs, title, pattern_data):
    """Create a new trade-off slide using blank layout with custom textboxes"""
    # Use a blank or simple layout (usually index 5 or 6)
    blank_layout = None
    for i, layout in enumerate(template_prs.slide_layouts):
        # Find a layout with just a title placeholder
        if len([ph for ph in layout.placeholders if hasattr(ph, 'text_frame')]) <= 2:
            blank_layout = layout
            break
    
    if not blank_layout:
        blank_layout = template_prs.slide_layouts[0]  # Use first layout as fallback
    
    # Add new slide with blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Set the title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Create left textbox for Vorteile
    left_textbox = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(1.5), 
        width=Inches(4.5),
        height=Inches(5.5)
    )
    left_frame = left_textbox.text_frame
    left_frame.clear()
    
    # Add Vorteile header
    p = left_frame.paragraphs[0]
    p.text = "Vorteile ✅"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)  # Green color
    
    # Add Vorteile items
    for item in pattern_data['vorteile']:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    
    # Create right textbox for Nachteile
    right_textbox = slide.shapes.add_textbox(
        left=Inches(5.2),
        top=Inches(1.5),
        width=Inches(4.5),
        height=Inches(5.5)
    )
    right_frame = right_textbox.text_frame
    right_frame.clear()
    
    # Add Nachteile header
    p = right_frame.paragraphs[0]
    p.text = "Nachteile ❌"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)  # Red color
    
    # Add Nachteile items
    for item in pattern_data['nachteile']:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.space_after = Pt(6)
    
    return slide

def recreate_missing_slides(prs_path, template_path):
    """Recreate all missing trade-off slides"""
    print(f"\n📝 Recreating slides in: {os.path.basename(prs_path)}")
    
    # Load presentations
    prs = Presentation(prs_path)
    template_prs = Presentation(template_path)
    
    tradeoffs_data = get_tradeoffs_data()
    
    # Define which slides should exist for each presentation
    if 'part1' in prs_path:
        expected_patterns = [
            ('Layered Architecture - Trade-offs', 'Layered Architecture'),
            ('Microservices - Trade-offs', 'Microservices'),
            ('Layered Architecture (Schichtarchitektur) - Vorteile/Nachteile', 'Layered Architecture'),
            ('Microservices Architecture - Vorteile/Nachteile', 'Microservices'),
            ('Hexagonal Architecture (Ports & Adapters) - Vorteile/Nachteile', 'Hexagonal Architecture'),
            ('Event-Driven Architecture - Vorteile/Nachteile', 'Event-Driven')
        ]
    else:  # part2
        expected_patterns = [
            ('CQRS - Trade-offs', 'CQRS'),
            ('CQRS - Vorteile/Nachteile', 'CQRS'),
            ('Event Sourcing - Vorteile/Nachteile', 'Event Sourcing'),
            ('Circuit Breaker - Vorteile/Nachteile', 'Circuit Breaker'),
            ('Saga Pattern - Vorteile/Nachteile', 'Saga Pattern'),
            ('API Gateway - Vorteile/Nachteile', 'API Gateway'),
            ('Service Mesh - Vorteile/Nachteile', 'Service Mesh'),
            ('Bulkhead Pattern - Vorteile/Nachteile', 'Bulkhead Pattern'),
            ('Domain-Driven Design - Vorteile/Nachteile', 'Domain-Driven Design')
        ]
    
    # Create all missing slides
    created_count = 0
    for slide_title, pattern_key in expected_patterns:
        if pattern_key in tradeoffs_data:
            print(f"   🔧 Creating: {slide_title}")
            
            slide = create_tradeoff_slide_with_textboxes(
                prs, template_prs, slide_title, tradeoffs_data[pattern_key]
            )
            
            if slide:
                created_count += 1
                print(f"      ✅ Created successfully")
            else:
                print(f"      ❌ Failed to create")
        else:
            print(f"   ⚠️ No data found for pattern: {pattern_key}")
    
    # Save the presentation
    prs.save(prs_path)
    print(f"   ✅ Created {created_count} new trade-off slides")
    print(f"   💾 Saved: {prs_path}")

def main():
    """Main execution"""
    print("🔧 ULTIMATE FIX: Creating Trade-off Slides with Textboxes")
    print("=" * 70)
    print("This will recreate all missing trade-off slides with proper content")
    
    base_dir = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture'
    template_path = os.path.join(base_dir, 'templates/VanillaCore.pptx')
    
    # Check if template exists
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return
    
    print(f"📋 Using template: {template_path}")
    
    # Fix both presentations
    presentations_dir = os.path.join(base_dir, 'presentations/powerpoint')
    
    recreate_missing_slides(
        os.path.join(presentations_dir, 'arch-patterns-part1.pptx'),
        template_path
    )
    
    recreate_missing_slides(
        os.path.join(presentations_dir, 'arch-patterns-part2.pptx'),
        template_path
    )
    
    print("\n" + "=" * 70)
    print("✅ ULTIMATE FIX COMPLETE!")
    print("   - All trade-off slides recreated with custom textboxes")
    print("   - Proper two-column layout with Vorteile/Nachteile")
    print("   - German content preserved")
    print("   - Color-coded headers (Green for benefits, Red for drawbacks)")

if __name__ == "__main__":
    main()