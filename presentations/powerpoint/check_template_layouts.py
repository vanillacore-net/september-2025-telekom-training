#!/usr/bin/env python3
"""
Check what layouts are actually available in the VanillaCore template
"""

from pptx import Presentation

def check_layouts():
    """Check available layouts in VanillaCore template"""
    template_path = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/templates/VanillaCore.pptx'
    
    prs = Presentation(template_path)
    
    print("AVAILABLE LAYOUTS IN VANILLACORE TEMPLATE:")
    print("=" * 60)
    
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\nLayout {i}: {layout.name}")
        
        # Count placeholders
        placeholder_count = 0
        content_placeholders = 0
        
        for shape in layout.placeholders:
            placeholder_count += 1
            if shape.placeholder_format.type == 7:  # Content placeholder
                content_placeholders += 1
            
            print(f"  - Placeholder {shape.placeholder_format.idx}: Type {shape.placeholder_format.type}")
            
            # Placeholder types:
            # 0 = Title
            # 1 = Body/Content  
            # 7 = Content
            # 13 = Picture
            # 14 = Chart
            # etc.
            
        print(f"  Total placeholders: {placeholder_count}")
        print(f"  Content placeholders: {content_placeholders}")
        
        # Check if this could be a two-column layout
        if content_placeholders >= 2:
            print(f"  ⭐ POTENTIAL TWO-COLUMN LAYOUT")

if __name__ == "__main__":
    check_layouts()