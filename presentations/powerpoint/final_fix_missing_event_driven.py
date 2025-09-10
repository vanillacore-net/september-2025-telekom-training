#!/usr/bin/env python3
"""
Add Event-Driven trade-offs that was missed
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_tradeoff_slide(prs, pattern_name, vorteile, nachteile, production=None):
    """Add a trade-off slide with two columns"""
    # Use layout 6 (Blank layout) for custom content
    slide_layout = prs.slide_layouts[6]  
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = f"{pattern_name} - Vorteile/Nachteile"
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Create two columns manually
    # Left column - Vorteile
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    # Add Vorteile header
    p = left_frame.paragraphs[0]
    p.text = "Vorteile ✅"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x99, 0x00)  # Green
    
    # Add Vorteile points
    for vorteil in vorteile:
        p = left_frame.add_paragraph()
        p.text = f"• {vorteil}"
        p.font.size = Pt(16)
        p.level = 0
    
    # Right column - Nachteile
    right_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    # Add Nachteile header
    p = right_frame.paragraphs[0]
    p.text = "Nachteile ❌"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)  # Red
    
    # Add Nachteile points
    for nachteil in nachteile:
        p = right_frame.add_paragraph()
        p.text = f"• {nachteil}"
        p.font.size = Pt(16)
        p.level = 0
    
    # Add production considerations if provided
    if production:
        prod_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
        prod_frame = prod_box.text_frame
        prod_frame.word_wrap = True
        
        p = prod_frame.paragraphs[0]
        p.text = "Production Considerations ⚠️"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0x99, 0x00)  # Orange
        
        for item in production:
            p = prod_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.level = 0
    
    return slide

def main():
    # Add Event-Driven trade-offs manually
    prs_path = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/arch-patterns-part1.pptx'
    prs = Presentation(prs_path)
    
    # Check if Event-Driven trade-offs already exist
    existing_titles = []
    for slide in prs.slides:
        if slide.shapes.title:
            existing_titles.append(slide.shapes.title.text)
    
    event_driven_tradeoff_title = "Event-Driven Architecture - Vorteile/Nachteile"
    
    if event_driven_tradeoff_title not in existing_titles:
        print("Adding Event-Driven Architecture trade-offs")
        
        # From the markdown file
        vorteile = [
            "Lose Kopplung zwischen Komponenten",
            "Asynchrone Verarbeitung für bessere Performance", 
            "Event Replay für Debugging und Analytics möglich",
            "Einfaches Hinzufügen neuer Event Consumer"
        ]
        
        nachteile = [
            "Eventually Consistent Data",
            "Komplexere Error Handling (Dead Letter Queues)",
            "Event Schema Evolution herausfordernd",
            "Debugging von Event Flows schwierig"
        ]
        
        production = [
            "Event Schema Registry für Compatibility",
            "Dead Letter Queue für Failed Events", 
            "Event Compression für hohe Throughput",
            "Monitoring von Queue Depths und Lag"
        ]
        
        add_tradeoff_slide(prs, "Event-Driven Architecture", vorteile, nachteile, production)
        
        # Save updated presentation
        prs.save(prs_path)
        print(f"Added Event-Driven trade-offs to: {prs_path}")
    else:
        print("Event-Driven trade-offs already exist")

if __name__ == "__main__":
    main()