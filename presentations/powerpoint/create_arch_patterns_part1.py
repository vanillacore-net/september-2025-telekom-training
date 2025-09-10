#!/usr/bin/env python3
"""
Create Architectural Patterns Part 1 PowerPoint presentation
Using VanillaCore.pptx template with German content
"""

from pptx import Presentation
import os

def create_arch_patterns_part1():
    """Create the architectural patterns part 1 presentation"""
    
    # Load the VanillaCore template
    template_path = "/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/templates/VanillaCore.pptx"
    prs = Presentation(template_path)
    
    # Clear existing slides (keep only layouts)
    slide_count = len(prs.slides)
    for i in range(slide_count - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    # Layout assignments
    title_layout = prs.slide_layouts[0]        # Title Slide
    section_layout = prs.slide_layouts[1]      # Section Header
    bullets_layout = prs.slide_layouts[2]      # Title and Content with Bullets
    content_layout = prs.slide_layouts[3]      # Title and Content (plain)
    two_column_layout = prs.slide_layouts[4]   # Two Columns
    code_layout = prs.slide_layouts[9]         # Code-Block layout
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Architectural Patterns - Part 1"
    slide1.placeholders[1].text = "Design Patterns Workshop - Telekom Training"
    
    # Slide 2: Overview with ASCII diagram
    slide2 = prs.slides.add_slide(code_layout)
    slide2.shapes.title.text = "Overview: Layer-Based Architectures"
    overview_diagram = """Layered Architecture        N-Tier Architecture        Clean Architecture
┌─────────────────┐        ┌─────────────────┐        ┌─────────────────┐
│   Presentation  │        │   Client Tier   │        │   Frameworks    │
├─────────────────┤        ├─────────────────┤        ├─────────────────┤
│    Business     │        │ Application     │        │   Interface     │
├─────────────────┤  <-->  │      Tier       │  <-->  │   Adapters      │
│  Persistence    │        ├─────────────────┤        ├─────────────────┤
├─────────────────┤        │   Data Tier     │        │  Use Cases      │
│    Database     │        └─────────────────┘        ├─────────────────┤
└─────────────────┘                                   │   Entities      │
                                                       └─────────────────┘
Classic Layering           Enterprise N-Tier          Domain-Centric

Focus: Clear separation of concerns through horizontal layers"""
    slide2.placeholders[1].text = overview_diagram
    
    # Slide 3: Layered Architecture Section Header
    slide3 = prs.slides.add_slide(section_layout)
    slide3.shapes.title.text = "1. Layered Architecture (Schichtarchitektur)"
    
    # Slide 4: Layered Architecture Schema
    slide4 = prs.slides.add_slide(code_layout)
    slide4.shapes.title.text = "Layered Architecture - Schema"
    layered_schema = """┌─────────────────────────────────────────────────┐
│                Presentation Layer               │ <- Web UI, REST APIs
├─────────────────────────────────────────────────┤
│                 Business Layer                  │ <- Domain Logic
├─────────────────────────────────────────────────┤
│                Persistence Layer                │ <- Data Access
├─────────────────────────────────────────────────┤
│                 Database Layer                  │ <- Data Storage
└─────────────────────────────────────────────────┘

Dependency Flow: Top -> Down (Higher layers depend on lower layers)
Data Flow: Bidirectional through defined interfaces"""
    slide4.placeholders[1].text = layered_schema
    
    # Slide 5: Layered Architecture - Telekom Use Case
    slide5 = prs.slides.add_slide(bullets_layout)
    slide5.shapes.title.text = "Layered Architecture - Telekom Use Case"
    layered_use_case = """Network Configuration Management System

Layer Implementation:
• Presentation: Web Dashboard für Network Engineers
• Business: Configuration Validation, Change Management
• Persistence: Configuration Repository, Audit Logging
• Database: PostgreSQL für Config Data, MongoDB für Audit Logs

Real-world Application:
• SNMP configuration management
• Network device inventory
• Change approval workflows
• Configuration backup and restore"""
    slide5.placeholders[1].text = layered_use_case
    
    # Slide 6: Layered Architecture - Trade-offs
    slide6 = prs.slides.add_slide(two_column_layout)
    slide6.shapes.title.text = "Layered Architecture - Trade-offs"
    
    advantages = """Vorteile ✅

• Klare Trennung der Verantwortlichkeiten
• Einfache Testbarkeit jeder Schicht
• Bewährtes Pattern mit hoher Entwickler-Akzeptanz
• Gute Performance bei einfachen CRUD-Operationen

Production Considerations ⚠️
• Horizontal Skalierung nur als Ganzes möglich
• Session-Management zwischen Layers erforderlich
• Monitoring muss alle Schichten umfassen"""
    
    disadvantages = """Nachteile ❌

• Monolithischer Charakter erschwert Skalierung
• Änderungen propagieren durch alle Schichten
• Datenbankschema-Changes beeinflussen alle Layer
• Schwierig für komplexe Domain Logic"""
    
    slide6.placeholders[1].text = advantages
    slide6.placeholders[2].text = disadvantages
    
    # Slide 7: Layered Architecture - When to Use
    slide7 = prs.slides.add_slide(two_column_layout)
    slide7.shapes.title.text = "Layered Architecture - When to Use"
    
    ideal_for = """Ideal für:

• Kleine bis mittlere Teams (<5 Entwickler)
• Einfache CRUD-Operationen
• Prototypen und Proof-of-Concepts
• Kurze Time-to-Market Anforderungen
• Bewährte Entwicklungsstandards erforderlich"""
    
    avoid_when = """Vermeiden wenn:

• Hohe Skalierungsanforderungen
• Komplexe Domain Logic
• Mehrere unabhängige Teams
• Mikroservice-Architektur gewünscht"""
    
    slide7.placeholders[1].text = ideal_for
    slide7.placeholders[2].text = avoid_when
    
    # Slide 8: Microservices Architecture Section Header
    slide8 = prs.slides.add_slide(section_layout)
    slide8.shapes.title.text = "2. Microservices Architecture"
    
    # Slide 9: Microservices Architecture Schema
    slide9 = prs.slides.add_slide(code_layout)
    slide9.shapes.title.text = "Microservices Architecture - Schema"
    microservices_schema = """Service Landscape:
┌───────────────┐  ┌───────────────┐  ┌───────────────┐
│  Device Mgmt  │  │   Monitoring  │  │  Configuration │
│   Service     │  │    Service    │  │    Service    │
├───────────────┤  ├───────────────┤  ├───────────────┤
│  PostgreSQL   │  │  InfluxDB     │  │   MongoDB     │
└───────────────┘  └───────────────┘  └───────────────┘
         │                  │                  │
         └──────────────────┼──────────────────┘
                            │
                   ┌─────────────────┐
                   │   API Gateway   │
                   ├─────────────────┤
                   │ Load Balancer   │
                   ├─────────────────┤
                   │ Service Mesh    │
                   └─────────────────┘"""
    slide9.placeholders[1].text = microservices_schema
    
    # Slide 10: Microservices - Telekom Use Case
    slide10 = prs.slides.add_slide(bullets_layout)
    slide10.shapes.title.text = "Microservices - Telekom Use Case"
    microservices_use_case = """Network Operations Center (NOC) Platform

Service Breakdown:
• Device Management: CRUD für Router, Switches, Firewalls
• Monitoring Service: Metriken, Alerts, Health Checks
• Configuration Service: Config Templates, Deployment
• Analytics Service: Performance Analysis, Capacity Planning

Technology Stack:
• Container: Docker + Kubernetes
• API Gateway: Kong oder Istio
• Service Mesh: Linkerd oder Istio
• Monitoring: Prometheus + Grafana"""
    slide10.placeholders[1].text = microservices_use_case
    
    # Slide 11: Microservices - Trade-offs
    slide11 = prs.slides.add_slide(two_column_layout)
    slide11.shapes.title.text = "Microservices - Trade-offs"
    
    micro_advantages = """Vorteile ✅

• Unabhängige Skalierung pro Service
• Technology Stack Diversität möglich
• Fehler-Isolierung zwischen Services
• Parallele Entwicklung durch verschiedene Teams

Production Considerations ⚠️
• Container Orchestration (Kubernetes) notwendig
• Service Mesh für Service-to-Service Communication
• Distributed Tracing für End-to-End Monitoring
• Event Sourcing für Data Consistency"""
    
    micro_disadvantages = """Nachteile ❌

• Hohe operationale Komplexität
• Distributed System Challenges (Latency, Partial Failures)
• Data Consistency zwischen Services schwierig
• Service Discovery und Load Balancing erforderlich"""
    
    slide11.placeholders[1].text = micro_advantages
    slide11.placeholders[2].text = micro_disadvantages
    
    # Slide 12: Microservices - When to Use
    slide12 = prs.slides.add_slide(two_column_layout)
    slide12.shapes.title.text = "Microservices - When to Use"
    
    micro_ideal = """Ideal für:

• Große verteilte Teams (5+ Entwickler)
• Independent Service Scaling erforderlich
• Complex Domain Logic
• High Availability Requirements
• Technology Stack Diversität gewünscht"""
    
    micro_avoid = """Vermeiden wenn:

• Kleine Teams (<5 Entwickler)
• Einfache Monolithische Anwendungen
• Begrenzte DevOps-Expertise
• Keine Container-Infrastruktur"""
    
    slide12.placeholders[1].text = micro_ideal
    slide12.placeholders[2].text = micro_avoid
    
    # Slide 13: Event-Driven Architecture Section Header
    slide13 = prs.slides.add_slide(section_layout)
    slide13.shapes.title.text = "3. Event-Driven Architecture"
    
    # Slide 14: Event-Driven Architecture Schema
    slide14 = prs.slides.add_slide(code_layout)
    slide14.shapes.title.text = "Event-Driven Architecture - Schema"
    event_schema = """Event Flow:
┌─────────────┐    Events     ┌─────────────┐    Events     ┌─────────────┐
│  Producers  │ ─────────────> │  Event Bus  │ ─────────────> │  Consumers  │
│             │               │             │               │             │
│ Device      │               │  Apache     │               │ Alerting    │
│ Sensors     │               │  Kafka      │               │ System      │
│             │               │             │               │             │
│ User        │               │ Topics:     │               │ Analytics   │
│ Actions     │               │ - alerts    │               │ Engine      │
│             │               │ - metrics   │               │             │
│ System      │               │ - configs   │               │ Dashboard   │
│ Events      │               │ - audit     │               │ Updates     │
└─────────────┘               └─────────────┘               └─────────────┘

    │                             │                             │
    └─── Async Pub/Sub ────────────┼─── Message Queues ────────┘
                                   │
                          ┌─────────────┐
                          │ Event Store │
                          │ (History)   │
                          └─────────────┘"""
    slide14.placeholders[1].text = event_schema
    
    # Slide 15: Event-Driven - Telekom Use Case
    slide15 = prs.slides.add_slide(bullets_layout)
    slide15.shapes.title.text = "Event-Driven - Telekom Use Case"
    event_use_case = """Network Event Processing Platform

Event Flow:
• Events: Device failures, configuration changes, performance alerts
• Producers: SNMP agents, log collectors, user interfaces
• Consumers: Alert system, analytics, dashboard updates, audit logging

Implementation:
• Event Bus: Apache Kafka
• Event Topics: alerts, metrics, configs, audit
• Event Store: Persistent event history
• Consumers: Real-time and batch processing"""
    slide15.placeholders[1].text = event_use_case
    
    # Slide 16: Event-Driven - Trade-offs
    slide16 = prs.slides.add_slide(two_column_layout)
    slide16.shapes.title.text = "Event-Driven - Trade-offs"
    
    event_advantages = """Vorteile ✅

• Lose Kopplung zwischen Komponenten
• Asynchrone Verarbeitung für bessere Performance
• Event Replay für Debugging und Analytics möglich
• Einfaches Hinzufügen neuer Event Consumer

Production Considerations ⚠️
• Event Schema Registry für Compatibility
• Dead Letter Queue für Failed Events
• Event Compression für hohe Throughput
• Monitoring von Queue Depths und Lag"""
    
    event_disadvantages = """Nachteile ❌

• Eventually Consistent Data
• Komplexere Error Handling (Dead Letter Queues)
• Event Schema Evolution herausfordernd
• Debugging von Event Flows schwierig"""
    
    slide16.placeholders[1].text = event_advantages
    slide16.placeholders[2].text = event_disadvantages
    
    # Slide 17: Event-Driven - When to Use
    slide17 = prs.slides.add_slide(two_column_layout)
    slide17.shapes.title.text = "Event-Driven - When to Use"
    
    event_ideal = """Ideal für:

• Real-time Monitoring und Alerting
• High-Throughput Event Processing
• Lose gekoppelte Systeme
• Asynchrone Verarbeitung erforderlich
• Event Replay und Analytics wichtig"""
    
    event_avoid = """Vermeiden wenn:

• Strong Consistency erforderlich
• Einfache Request-Response Patterns
• Begrenzte Message Queue Expertise
• Synchrone Verarbeitung ausreichend"""
    
    slide17.placeholders[1].text = event_ideal
    slide17.placeholders[2].text = event_avoid
    
    # Slide 18: Hexagonal Architecture Section Header
    slide18 = prs.slides.add_slide(section_layout)
    slide18.shapes.title.text = "4. Hexagonal Architecture (Ports & Adapters)"
    
    # Slide 19: Hexagonal Architecture Schema
    slide19 = prs.slides.add_slide(code_layout)
    slide19.shapes.title.text = "Hexagonal Architecture - Schema"
    hexagonal_schema = """                    Hexagonal Architecture
                           (Ports & Adapters)

    ┌─────────────────┐                    ┌─────────────────┐
    │   Web UI        │                    │   Database      │
    │   Adapter       │                    │   Adapter       │
    └─────┬───────────┘                    └─────┬───────────┘
          │                                      │
    ┌─────┴───────────┐                    ┌─────┴───────────┐
    │   HTTP Port     │                    │ Persistence Port│
    └─────┬───────────┘                    └─────┬───────────┘
          │                                      │
          │     ┌─────────────────────────┐      │
          └─────│                         │──────┘
                │     Domain Core         │
                │                         │
    ┌───────────│   Business Logic        │──────────┐
    │           │                         │          │
    │           └─────────────────────────┘          │
    │                                                │
┌───┴───────────┐                            ┌─────┴───────────┐
│  API Port     │                            │   Event Port    │
└───┬───────────┘                            └─────┬───────────┘
    │                                              │
┌───┴───────────┐                            ┌─────┴───────────┐
│   REST API    │                            │   Message       │
│   Adapter     │                            │   Queue Adapter │
└───────────────┘                            └─────────────────┘"""
    slide19.placeholders[1].text = hexagonal_schema
    
    # Slide 20: Hexagonal Architecture - Telekom Use Case
    slide20 = prs.slides.add_slide(bullets_layout)
    slide20.shapes.title.text = "Hexagonal Architecture - Telekom Use Case"
    hexagonal_use_case = """Device Configuration Service

Core Domain:
• Device configuration validation
• Template processing
• Business rules enforcement

Adapters:
• REST API: Web interface
• SNMP interface: Device communication
• File-based config: Bulk operations
• Message queues: Event notifications

Ports:
• Device management: Core business operations
• Configuration storage: Persistence interface
• Event notification: Integration interface"""
    slide20.placeholders[1].text = hexagonal_use_case
    
    # Slide 21: Hexagonal Architecture - Trade-offs
    slide21 = prs.slides.add_slide(two_column_layout)
    slide21.shapes.title.text = "Hexagonal Architecture - Trade-offs"
    
    hex_advantages = """Vorteile ✅

• Domain Logic vollständig isoliert von Infrastructure
• Austauschbare External Dependencies (Database, APIs)
• Excellente Testbarkeit durch Dependency Inversion
• Framework-agnostic Domain Layer

Production Considerations ⚠️
• Adapter Health Checks für alle External Systems
• Circuit Breaker Pattern für External Dependencies
• Adapter-specific Monitoring und Alerting
• Configuration für Adapter Switching"""
    
    hex_disadvantages = """Nachteile ❌

• Initial Setup komplexer als Layer Architecture
• Mehr Abstractions und Interfaces
• Kann Over-Engineering für einfache CRUD Apps sein
• Team muss Dependency Inversion verstehen"""
    
    slide21.placeholders[1].text = hex_advantages
    slide21.placeholders[2].text = hex_disadvantages
    
    # Slide 22: Hexagonal Architecture - When to Use
    slide22 = prs.slides.add_slide(two_column_layout)
    slide22.shapes.title.text = "Hexagonal Architecture - When to Use"
    
    hex_ideal = """Ideal für:

• Long-lived Business Applications
• Frequent External System Changes
• High Testing Requirements
• Framework Independence gewünscht
• Complex Domain Logic"""
    
    hex_avoid = """Vermeiden wenn:

• Einfache CRUD Applications
• Begrenzte DI/IoC Erfahrung
• Kurze Entwicklungszyklen
• Kleine Prototyp-Projekte"""
    
    slide22.placeholders[1].text = hex_ideal
    slide22.placeholders[2].text = hex_avoid
    
    # Slide 23: Architecture Decision Matrix
    slide23 = prs.slides.add_slide(content_layout)
    slide23.shapes.title.text = "Architecture Decision Matrix - Part 1"
    matrix_content = """| Criteria | Layered | Microservices | Event-Driven | Hexagonal |
|----------|---------|---------------|---------------|-----------|
| Team Size | 1-5 | 5+ | 3-8 | 2-6 |
| Scalability | Low | High | High | Medium |
| Complexity | Low | Very High | Medium | Medium |
| Performance | Good | Variable | High | Good |
| Maintainability | High | Medium | Medium | High |
| Testing | Easy | Complex | Medium | Easy |
| Deployment | Simple | Complex | Medium | Simple |
| Learning Curve | Low | High | Medium | Medium |"""
    slide23.placeholders[1].text = matrix_content
    
    # Slide 24: Telekom-Specific Recommendations
    slide24 = prs.slides.add_slide(bullets_layout)
    slide24.shapes.title.text = "Telekom-Specific Recommendations"
    recommendations = """Für Network Operations (NOC):
• ✅ Event-Driven für Real-time Monitoring
• ✅ Microservices für Service-specific Teams
• ⚠️ Layered nur für Simple Tools

Für Infrastructure Services:
• ✅ Microservices für Independent Scaling
• ✅ Event-Driven für System Integration
• ✅ Hexagonal für Core Business Logic

Migration Strategy:
• Start: Layered für Prototypen
• Scale: Event-Driven für Integration
• Optimize: Microservices für Team Growth"""
    slide24.placeholders[1].text = recommendations
    
    # Slide 25: Summary - Part 1
    slide25 = prs.slides.add_slide(bullets_layout)
    slide25.shapes.title.text = "Summary - Part 1"
    summary_content = """Pattern Selection Guide:

Start Simple (Layered):
• Prototypen und Proof-of-Concepts
• Kleine Teams (<5 Entwickler)
• Einfache CRUD-Operationen

Scale Smart (Microservices):
• Große verteilte Teams
• Independent Service Scaling
• High Availability Requirements

Integrate Efficiently (Event-Driven):
• Real-time Processing
• System Integration
• Asynchronous Workflows

Design for Change (Hexagonal):
• Complex Domain Logic
• Testability Requirements
• External System Integration

Next: Part 2 covers CQRS, Event Sourcing, and advanced patterns"""
    slide25.placeholders[1].text = summary_content
    
    # Save the presentation
    output_path = "/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/arch-patterns-part1.pptx"
    prs.save(output_path)
    
    print(f"✅ Presentation created successfully: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\n📋 Slide overview:")
    for i, slide in enumerate(prs.slides, 1):
        print(f"  {i:2d}. {slide.shapes.title.text}")
    
    return output_path

if __name__ == "__main__":
    create_arch_patterns_part1()