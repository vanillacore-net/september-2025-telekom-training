#!/usr/bin/env python3
"""
PRB-026: Fix trade-off slides using CORRECT comparison layouts
Use Layout 6 (Comparison Bulletpoints) with proper placeholder structure
"""

import os
from pptx import Presentation

def get_tradeoffs_data():
    """Complete trade-offs data for all patterns"""
    return {
        'Layered Architecture': {
            'vorteile': [
                'Klare Trennung der Verantwortlichkeiten',
                'Einfache Testbarkeit jeder Schicht',
                'Bewährtes Pattern mit hoher Entwickler-Akzeptanz',
                'Gute Performance bei einfachen CRUD-Operationen'
            ],
            'nachteile': [
                'Monolithischer Charakter erschwert Skalierung',
                'Änderungen propagieren durch alle Schichten',
                'Datenbankschema-Changes beeinflussen alle Layer',
                'Schwierig für komplexe Domain Logic'
            ]
        },
        'Microservices': {
            'vorteile': [
                'Unabhängige Skalierung pro Service',
                'Technology Stack Diversität möglich',
                'Fehler-Isolierung zwischen Services',
                'Parallele Entwicklung durch verschiedene Teams'
            ],
            'nachteile': [
                'Hohe operationale Komplexität',
                'Distributed System Challenges (Latency, Partial Failures)',
                'Data Consistency zwischen Services schwierig',
                'Service Discovery und Load Balancing erforderlich'
            ]
        },
        'Event-Driven': {
            'vorteile': [
                'Lose Kopplung zwischen Komponenten',
                'Asynchrone Verarbeitung für bessere Performance',
                'Event Replay für Debugging und Analytics möglich',
                'Einfaches Hinzufügen neuer Event Consumer'
            ],
            'nachteile': [
                'Eventually Consistent Data',
                'Komplexere Error Handling (Dead Letter Queues)',
                'Event Schema Evolution herausfordernd',
                'Debugging von Event Flows schwierig'
            ]
        },
        'Hexagonal Architecture': {
            'vorteile': [
                'Domain Logic vollständig isoliert von Infrastructure',
                'Austauschbare External Dependencies (Database, APIs)',
                'Excellente Testbarkeit durch Dependency Inversion',
                'Framework-agnostic Domain Layer'
            ],
            'nachteile': [
                'Initial Setup komplexer als Layer Architecture',
                'Mehr Abstractions und Interfaces',
                'Kann Over-Engineering für einfache CRUD Apps sein',
                'Team muss Dependency Inversion verstehen'
            ]
        },
        'CQRS': {
            'vorteile': [
                'Optimierte Read und Write Models',
                'Unabhängige Skalierung von Queries und Commands',
                'Bessere Performance durch spezialisierte Datenmodelle',
                'Event Sourcing Integration möglich'
            ],
            'nachteile': [
                'Erhöhte Komplexität der Architektur',
                'Eventual Consistency zwischen Read/Write',
                'Synchronisation der Models erforderlich',
                'Mehr Code und Infrastruktur'
            ]
        },
        'Event Sourcing': {
            'vorteile': [
                'Vollständige Audit Trail automatisch',
                'Time Travel und Event Replay möglich',
                'Debugging durch Event History',
                'Integration mit CQRS natural'
            ],
            'nachteile': [
                'Event Store Management komplex',
                'Schema Evolution schwierig',
                'Snapshot Management erforderlich',
                'GDPR Compliance herausfordernd'
            ]
        },
        'Circuit Breaker': {
            'vorteile': [
                'Schnelles Fail-Fast Verhalten',
                'System Stabilität bei Ausfällen',
                'Automatische Recovery Mechanismen',
                'Cascading Failures Prevention'
            ],
            'nachteile': [
                'Threshold Tuning erforderlich',
                'False Positives möglich',
                'Additional Monitoring Layer',
                'Complexity in Error Handling'
            ]
        },
        'Saga Pattern': {
            'vorteile': [
                'Distributed Transaction Management',
                'Compensation Logic für Rollbacks',
                'Keine 2PC Lock Requirements',
                'Bessere Performance als 2PC'
            ],
            'nachteile': [
                'Komplexe Compensation Logic',
                'Debugging schwierig',
                'Partial Failure Handling',
                'Testing der Compensation'
            ]
        },
        'API Gateway': {
            'vorteile': [
                'Single Entry Point für Clients',
                'Cross-Cutting Concerns zentral',
                'API Composition möglich',
                'Backend Service Abstraktion'
            ],
            'nachteile': [
                'Single Point of Failure',
                'Additional Latency Layer',
                'Deployment Koordination',
                'Gateway Bottleneck möglich'
            ]
        },
        'Service Mesh': {
            'vorteile': [
                'Service-to-Service Security',
                'Automatic Load Balancing',
                'Observability out-of-the-box',
                'Traffic Management Features'
            ],
            'nachteile': [
                'Operational Overhead',
                'Resource Consumption (Sidecars)',
                'Learning Curve hoch',
                'Debugging Complexity'
            ]
        },
        'Bulkhead Pattern': {
            'vorteile': [
                'Fault Isolation zwischen Components',
                'Resource Protection',
                'Predictable Performance',
                'Prevents Resource Exhaustion'
            ],
            'nachteile': [
                'Resource Underutilization möglich',
                'Configuration Complexity',
                'More Infrastructure Required',
                'Monitoring Overhead'
            ]
        },
        'Domain-Driven Design': {
            'vorteile': [
                'Klare Bounded Contexts',
                'Ubiquitous Language etabliert',
                'Business-IT Alignment',
                'Modulare Struktur'
            ],
            'nachteile': [
                'Steile Learning Curve',
                'Upfront Design Investment',
                'Domain Expert Availability',
                'Over-Engineering Gefahr'
            ]
        }
    }

def find_and_delete_tradeoff_slides(prs):
    """Find and delete all existing trade-off slides"""
    slides_to_delete = []
    
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.lower()
            if any(keyword in title for keyword in ['trade-off', 'vorteile', 'nachteile']):
                slides_to_delete.append(i)
                print(f"   🗑️ Found trade-off slide to delete: {slide.shapes.title.text}")
    
    # Delete slides in reverse order to maintain indices
    for i in reversed(slides_to_delete):
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(xml_slides[i])
    
    print(f"   ✅ Deleted {len(slides_to_delete)} old trade-off slides")
    return len(slides_to_delete)

def create_comparison_slide(prs, template_prs, title, pattern_data, use_layout_7=False):
    """Create a trade-off slide using Layout 6 (Comparison Bulletpoints) or Layout 7 (Comparison)"""
    
    # Use Layout 6 (Comparison Bulletpoints) or Layout 7 (Comparison)
    layout_index = 7 if use_layout_7 else 6
    comparison_layout = template_prs.slide_layouts[layout_index]
    
    print(f"   🏗️ Using Layout {layout_index}: {comparison_layout.name}")
    
    # Add new slide with comparison layout
    slide = prs.slides.add_slide(comparison_layout)
    
    # Set the title (Placeholder 0)
    if slide.shapes.title:
        slide.shapes.title.text = title
        print(f"   📝 Set title: {title}")
    
    # Find placeholders by their index
    placeholders = {}
    for shape in slide.placeholders:
        placeholders[shape.placeholder_format.idx] = shape
    
    # Layout 6 & 7 structure:
    # Placeholder 0: Title
    # Placeholder 1: Left header (Body)
    # Placeholder 2: Left content (Object) 
    # Placeholder 3: Right header (Body)
    # Placeholder 4: Right content (Object)
    
    # Set left header (Placeholder 1)
    if 1 in placeholders:
        left_header = placeholders[1]
        if hasattr(left_header, 'text_frame'):
            left_header.text_frame.text = "Vorteile ✅"
            print("   ✅ Set left header: Vorteile ✅")
    
    # Set left content (Placeholder 2)
    if 2 in placeholders:
        left_content = placeholders[2]
        if hasattr(left_content, 'text_frame'):
            left_content.text_frame.clear()
            for item in pattern_data['vorteile']:
                p = left_content.text_frame.add_paragraph()
                p.text = item
            print(f"   📋 Added {len(pattern_data['vorteile'])} Vorteile items")
    
    # Set right header (Placeholder 3)
    if 3 in placeholders:
        right_header = placeholders[3]
        if hasattr(right_header, 'text_frame'):
            right_header.text_frame.text = "Nachteile ❌"
            print("   ❌ Set right header: Nachteile ❌")
    
    # Set right content (Placeholder 4)
    if 4 in placeholders:
        right_content = placeholders[4]
        if hasattr(right_content, 'text_frame'):
            right_content.text_frame.clear()
            for item in pattern_data['nachteile']:
                p = right_content.text_frame.add_paragraph()
                p.text = item
            print(f"   📋 Added {len(pattern_data['nachteile'])} Nachteile items")
    
    return slide

def fix_presentation(prs_path, template_path):
    """Fix trade-off slides in one presentation"""
    print(f"\n🔧 Fixing presentation: {os.path.basename(prs_path)}")
    
    # Load presentations
    prs = Presentation(prs_path)
    template_prs = Presentation(template_path)
    
    # Delete existing trade-off slides
    deleted_count = find_and_delete_tradeoff_slides(prs)
    
    tradeoffs_data = get_tradeoffs_data()
    
    # Define which slides should exist for each presentation
    if 'part1' in prs_path:
        expected_patterns = [
            ('Layered Architecture - Vorteile/Nachteile', 'Layered Architecture'),
            ('Microservices - Vorteile/Nachteile', 'Microservices'),
            ('Hexagonal Architecture - Vorteile/Nachteile', 'Hexagonal Architecture'),
            ('Event-Driven Architecture - Vorteile/Nachteile', 'Event-Driven')
        ]
    else:  # part2
        expected_patterns = [
            ('CQRS - Vorteile/Nachteile', 'CQRS'),
            ('Event Sourcing - Vorteile/Nachteile', 'Event Sourcing'),
            ('Circuit Breaker - Vorteile/Nachteile', 'Circuit Breaker'),
            ('Saga Pattern - Vorteile/Nachteile', 'Saga Pattern'),
            ('API Gateway - Vorteile/Nachteile', 'API Gateway'),
            ('Service Mesh - Vorteile/Nachteile', 'Service Mesh'),
            ('Bulkhead Pattern - Vorteile/Nachteile', 'Bulkhead Pattern'),
            ('Domain-Driven Design - Vorteile/Nachteile', 'Domain-Driven Design')
        ]
    
    # Create new slides using comparison layout
    created_count = 0
    for slide_title, pattern_key in expected_patterns:
        if pattern_key in tradeoffs_data:
            print(f"   🏗️ Creating: {slide_title}")
            
            slide = create_comparison_slide(
                prs, template_prs, slide_title, tradeoffs_data[pattern_key]
            )
            
            if slide:
                created_count += 1
                print(f"      ✅ Created successfully")
            else:
                print(f"      ❌ Failed to create")
        else:
            print(f"   ⚠️ No data found for pattern: {pattern_key}")
    
    # Save the presentation
    prs.save(prs_path)
    print(f"   💾 Saved: {prs_path}")
    print(f"   📊 Summary: Deleted {deleted_count}, Created {created_count} slides")

def main():
    """Main execution"""
    print("🔧 PRB-026: Fix Trade-off Slides with CORRECT Comparison Layouts")
    print("=" * 80)
    print("Using Layout 6 (Comparison Bulletpoints) with proper placeholder structure:")
    print("  - Placeholder 0: Title")
    print("  - Placeholder 1: Left header (Vorteile ✅)")
    print("  - Placeholder 2: Left content area")
    print("  - Placeholder 3: Right header (Nachteile ❌)")
    print("  - Placeholder 4: Right content area")
    print()
    
    base_dir = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture'
    template_path = os.path.join(base_dir, 'templates/VanillaCore.pptx')
    
    # Check if template exists
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return
    
    print(f"📋 Using template: {template_path}")
    
    # Fix both presentations
    presentations_dir = os.path.join(base_dir, 'presentations/powerpoint')
    
    fix_presentation(
        os.path.join(presentations_dir, 'arch-patterns-part1.pptx'),
        template_path
    )
    
    fix_presentation(
        os.path.join(presentations_dir, 'arch-patterns-part2.pptx'),
        template_path
    )
    
    print("\n" + "=" * 80)
    print("✅ PRB-026 EXECUTION COMPLETE!")
    print("   - Deleted all broken trade-off slides")
    print("   - Created new slides using Layout 6 (Comparison Bulletpoints)")
    print("   - Proper placeholder structure with headers and content areas")
    print("   - German content preserved")
    print("   - Professional layout matching template design")

if __name__ == "__main__":
    main()