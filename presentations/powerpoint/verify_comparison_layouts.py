#!/usr/bin/env python3
"""
Verify that the comparison layout slides were created correctly
"""

import os
from pptx import Presentation

def verify_presentation(prs_path):
    """Verify trade-off slides in presentation"""
    print(f"\n🔍 Verifying: {os.path.basename(prs_path)}")
    
    if not os.path.exists(prs_path):
        print(f"   ❌ File not found: {prs_path}")
        return
    
    try:
        prs = Presentation(prs_path)
        tradeoff_slides = []
        
        for i, slide in enumerate(prs.slides):
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text
                if 'vorteile' in title.lower() and 'nachteile' in title.lower():
                    tradeoff_slides.append((i+1, title))
        
        print(f"   📋 Found {len(tradeoff_slides)} trade-off slides")
        
        for slide_num, title in tradeoff_slides:
            print(f"      📄 Slide {slide_num}: {title}")
            
            # Get the slide
            slide = prs.slides[slide_num - 1]
            
            # Check placeholders
            placeholders = {}
            for shape in slide.placeholders:
                placeholders[shape.placeholder_format.idx] = shape
            
            # Check if we have the comparison layout structure
            has_left_header = 1 in placeholders and hasattr(placeholders[1], 'text_frame')
            has_left_content = 2 in placeholders and hasattr(placeholders[2], 'text_frame')
            has_right_header = 3 in placeholders and hasattr(placeholders[3], 'text_frame')
            has_right_content = 4 in placeholders and hasattr(placeholders[4], 'text_frame')
            
            if has_left_header:
                left_header_text = placeholders[1].text_frame.text.strip()
                print(f"         ✅ Left header: {left_header_text}")
            else:
                print(f"         ❌ Missing left header placeholder")
            
            if has_left_content:
                left_paragraphs = len(placeholders[2].text_frame.paragraphs)
                print(f"         📋 Left content: {left_paragraphs} paragraphs")
            else:
                print(f"         ❌ Missing left content placeholder")
            
            if has_right_header:
                right_header_text = placeholders[3].text_frame.text.strip()
                print(f"         ❌ Right header: {right_header_text}")
            else:
                print(f"         ❌ Missing right header placeholder")
            
            if has_right_content:
                right_paragraphs = len(placeholders[4].text_frame.paragraphs)
                print(f"         📋 Right content: {right_paragraphs} paragraphs")
            else:
                print(f"         ❌ Missing right content placeholder")
        
        print(f"   ✅ Verification complete")
        
    except Exception as e:
        print(f"   ❌ Error verifying presentation: {str(e)}")

def main():
    """Main verification"""
    print("🔍 PRB-026 Verification: Comparison Layout Slides")
    print("=" * 60)
    
    base_dir = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture'
    presentations_dir = os.path.join(base_dir, 'presentations/powerpoint')
    
    verify_presentation(os.path.join(presentations_dir, 'arch-patterns-part1.pptx'))
    verify_presentation(os.path.join(presentations_dir, 'arch-patterns-part2.pptx'))
    
    print("\n" + "=" * 60)
    print("✅ Verification complete!")

if __name__ == "__main__":
    main()