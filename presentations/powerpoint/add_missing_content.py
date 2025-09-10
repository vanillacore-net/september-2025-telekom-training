#!/usr/bin/env python3
"""
Add missing Vorteile/Nachteile slides and speaker notes to architectural patterns presentations
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def extract_notes_from_markdown(markdown_content):
    """Extract all Note: sections from markdown content and create speaker notes"""
    notes = {}
    lines = markdown_content.split('\n')
    current_section = None
    
    # Simple approach: create speaker notes based on content sections
    for i, line in enumerate(lines):
        # Track current section
        if line.startswith('## ') or line.startswith('# '):
            current_section = line.strip('# ').strip()
            
            # Create speaker notes for each major section
            if current_section:
                # Look for the content after this heading to create notes
                content_lines = []
                j = i + 1
                while j < len(lines) and not (lines[j].startswith('#') or lines[j].startswith('---')):
                    if lines[j].strip() and not lines[j].startswith('```'):
                        content_lines.append(lines[j].strip())
                    j += 1
                
                # Create speaker note from first few lines of content
                if content_lines:
                    note_content = []
                    for content_line in content_lines[:3]:  # Take first 3 meaningful lines
                        if content_line and len(content_line) > 20:  # Only substantial lines
                            note_content.append(content_line)
                    
                    if note_content:
                        notes[current_section] = '\n'.join(note_content)
    
    return notes

def extract_tradeoffs_from_markdown(markdown_content):
    """Extract Vorteile/Nachteile sections from markdown"""
    tradeoffs = {}
    lines = markdown_content.split('\n')
    current_pattern = None
    current_section = None
    collecting_tradeoff = False
    
    i = 0
    while i < len(lines):
        line = lines[i]
        
        # Method 1: Pattern from Part 1 - numbered patterns with trade-offs
        pattern_match = re.match(r'^## (\d+)\.\s*([^-]+?)(?:\s*-\s*Trade-offs)?$', line)
        if pattern_match:
            current_pattern = pattern_match.group(2).strip()
        
        # Method 2: Pattern from Part 2 - direct trade-offs headers
        tradeoff_match = re.match(r'^## (.+?)\s+Trade-offs?$', line)
        if tradeoff_match:
            current_pattern = tradeoff_match.group(1).strip()
            collecting_tradeoff = True
            tradeoffs[current_pattern] = {'vorteile': [], 'nachteile': [], 'production': []}
            i += 1
            continue
        
        # Look for trade-offs section (Part 1 style)
        if current_pattern and ('- Trade-offs' in line or 'Trade-offs' in line):
            collecting_tradeoff = True
            if current_pattern not in tradeoffs:
                tradeoffs[current_pattern] = {'vorteile': [], 'nachteile': [], 'production': []}
            i += 1
            continue
        
        if collecting_tradeoff:
            # Skip empty lines and separators
            if not line.strip() or line.startswith('---'):
                i += 1
                continue
                
            # Check for new sections within trade-offs
            if line.strip().startswith('### Vorteile'):
                current_section = 'vorteile'
            elif line.strip().startswith('### Nachteile'):
                current_section = 'nachteile'
            elif line.strip().startswith('### Production Considerations'):
                current_section = 'production'
            elif line.startswith('##') and 'Trade-offs' not in line:
                # End of trade-offs section
                collecting_tradeoff = False
                current_section = None
                current_pattern = None
            elif current_section and line.strip().startswith(('✅', '❌', '-')):
                # Extract bullet point - handle both emoji and dash formats
                text = line.strip()
                if text.startswith('✅'):
                    text = text[2:].strip()
                elif text.startswith('❌'):
                    text = text[2:].strip()
                elif text.startswith('-'):
                    text = text[1:].strip()
                
                # Remove markdown formatting
                text = re.sub(r'\*\*([^*]+)\*\*:', r'\1:', text)
                
                if text:
                    tradeoffs[current_pattern][current_section].append(text)
        
        i += 1
    
    return tradeoffs

def add_tradeoff_slide(prs, pattern_name, vorteile, nachteile, production=None):
    """Add a trade-off slide with two columns"""
    # Use layout 6 (Content with Caption) which has title and content
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom content
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = f"{pattern_name} - Vorteile/Nachteile"
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Create two columns manually
    # Left column - Vorteile
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    # Add Vorteile header
    p = left_frame.paragraphs[0]
    p.text = "Vorteile ✅"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x99, 0x00)  # Green
    
    # Add Vorteile points
    for vorteil in vorteile:
        p = left_frame.add_paragraph()
        p.text = f"• {vorteil}"
        p.font.size = Pt(16)
        p.level = 0
    
    # Right column - Nachteile
    right_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    # Add Nachteile header
    p = right_frame.paragraphs[0]
    p.text = "Nachteile ❌"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)  # Red
    
    # Add Nachteile points
    for nachteil in nachteile:
        p = right_frame.add_paragraph()
        p.text = f"• {nachteil}"
        p.font.size = Pt(16)
        p.level = 0
    
    # Add production considerations if provided
    if production:
        prod_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
        prod_frame = prod_box.text_frame
        prod_frame.word_wrap = True
        
        p = prod_frame.paragraphs[0]
        p.text = "Production Considerations ⚠️"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0x99, 0x00)  # Orange
        
        for item in production:
            p = prod_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.level = 0
    
    return slide

def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide"""
    if notes_text and slide.notes_slide:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

def process_part1():
    """Process architectural patterns part 1"""
    print("Processing Part 1...")
    
    # Load existing presentation
    prs_path = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/arch-patterns-part1.pptx'
    if not os.path.exists(prs_path):
        print(f"File not found: {prs_path}")
        return
    
    prs = Presentation(prs_path)
    
    # Load markdown content
    md_path = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/architectural-patterns-part1.md'
    with open(md_path, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # Extract data
    notes = extract_notes_from_markdown(md_content)
    tradeoffs = extract_tradeoffs_from_markdown(md_content)
    
    print(f"Found {len(tradeoffs)} trade-off sections")
    print(f"Found {len(notes)} note sections")
    
    # Check what trade-off slides are missing
    patterns_with_tradeoffs = ['Layered Architecture (Schichtarchitektur)', 'Microservices Architecture', 
                              'Event-Driven Architecture', 'Hexagonal Architecture (Ports & Adapters)']
    
    existing_slide_titles = []
    for slide in prs.slides:
        if slide.shapes.title:
            existing_slide_titles.append(slide.shapes.title.text)
    
    print(f"Existing slides: {len(existing_slide_titles)}")
    print(f"Available trade-offs: {list(tradeoffs.keys())}")
    
    # Find where to insert trade-off slides and add missing ones
    for pattern in patterns_with_tradeoffs:
        tradeoff_title = f"{pattern} - Vorteile/Nachteile"
        tradeoff_title_alt = f"{pattern} - Trade-offs"
        
        if tradeoff_title not in existing_slide_titles and tradeoff_title_alt not in existing_slide_titles:
            # Look for the pattern in extracted tradeoffs
            found = False
            for key, value in tradeoffs.items():
                # More flexible matching
                pattern_keywords = pattern.lower().replace('(', '').replace(')', '').split()
                key_keywords = key.lower().replace('(', '').replace(')', '').split()
                
                if any(keyword in key.lower() for keyword in pattern_keywords[:2]):  # Match first 2 words
                    print(f"Adding trade-off slide for: {pattern} (matched with {key})")
                    if 'vorteile' in value and 'nachteile' in value:
                        add_tradeoff_slide(prs, pattern, value['vorteile'], value['nachteile'], value.get('production'))
                    found = True
                    break
            if not found:
                print(f"No trade-offs found for: {pattern}")
    
    # Add speaker notes to all slides
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title:
            title = slide.shapes.title.text
            # Find matching notes
            for section, section_notes in notes.items():
                if any(keyword in title.lower() for keyword in section.lower().split()):
                    if isinstance(section_notes, list):
                        note_text = '\n\n'.join(section_notes)
                    else:
                        note_text = section_notes
                    add_speaker_notes(slide, note_text)
                    break
    
    # Save updated presentation
    prs.save(prs_path)
    print(f"Updated presentation saved: {prs_path}")

def process_part2():
    """Process architectural patterns part 2"""
    print("Processing Part 2...")
    
    # Load existing presentation
    prs_path = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/arch-patterns-part2.pptx'
    if not os.path.exists(prs_path):
        print(f"File not found: {prs_path}")
        return
    
    prs = Presentation(prs_path)
    
    # Load markdown content
    md_path = '/Users/karsten/Nextcloud/Work/Training/2025-09_Telekom_Architecture/presentations/powerpoint/architectural-patterns-part2.md'
    with open(md_path, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # Extract data
    notes = extract_notes_from_markdown(md_content)
    tradeoffs = extract_tradeoffs_from_markdown(md_content)
    
    print(f"Found {len(tradeoffs)} trade-off sections")
    print(f"Found {len(notes)} note sections")
    
    # Add manual trade-offs for patterns that don't have explicit trade-offs sections
    manual_tradeoffs = {
        'Circuit Breaker': {
            'vorteile': [
                'Verhindert Cascading Failures in verteilten Systemen',
                'Schnelles Fail-Fast Verhalten statt langsame Timeouts',
                'Automatische Recovery durch Health Checks',
                'Monitoring und Alerting für Service Dependencies'
            ],
            'nachteile': [
                'Zusätzliche Complexity in Service Integration',
                'Konfiguration von Thresholds und Timeouts erforderlich',
                'False Positives können zu unnötigen Ausfällen führen',
                'Debugging wird komplexer durch zusätzliche Abstraction'
            ]
        },
        'Saga Pattern': {
            'vorteile': [
                'Verteilte Transaktionen ohne 2-Phase-Commit',
                'Bessere Performance als traditionelle ACID Transaktionen',
                'Flexible Compensation Logic für Rollbacks',
                'Natürliche Integration mit Event-Driven Architecture'
            ],
            'nachteile': [
                'Eventual Consistency statt ACID Guarantees',
                'Komplexe Rollback Logic (Compensating Actions)',
                'Debugging von Saga Flows schwierig',
                'Zusätzlicher Overhead für State Management'
            ]
        },
        'API Gateway': {
            'vorteile': [
                'Zentraler Entry Point für alle Client Requests',
                'Cross-cutting Concerns (Auth, Logging, Rate Limiting) zentral',
                'Service Discovery Abstraktion für Clients',
                'API Versioning und Evolution vereinfacht'
            ],
            'nachteile': [
                'Single Point of Failure ohne High Availability',
                'Performance Bottleneck bei hohem Traffic',
                'Zusätzliche Latency durch Proxy Layer',
                'Gateway selbst wird zu komplexem System'
            ]
        },
        'Service Mesh': {
            'vorteile': [
                'Infrastructure Layer für Service Communication',
                'Automatisches mTLS und Service Security',
                'Distributed Tracing und Observability out-of-the-box',
                'Traffic Management ohne Code Changes'
            ],
            'nachteile': [
                'Hohe operationale Komplexität',
                'Zusätzliche Latency durch Proxy Layer',
                'Steile Lernkurve für Operations Teams',
                'Resource Overhead durch Sidecar Proxies'
            ]
        },
        'Bulkhead Pattern': {
            'vorteile': [
                'Isolation von Critical Resources',
                'Ein problematischer Client kann andere nicht beeinträchtigen',
                'Verbesserte System Resilience',
                'Priorisierung von wichtigen Clients möglich'
            ],
            'nachteile': [
                'Reduced Resource Utilization durch Isolation',
                'Complexity bei Resource Pool Management',
                'Schwierige Konfiguration von Pool Sizes',
                'Monitoring von mehreren Resource Pools erforderlich'
            ]
        },
        'Domain-Driven Design': {
            'vorteile': [
                'Ubiquitous Language zwischen Business und Tech',
                'Bounded Context für klare Domain Grenzen',
                'Complex Domain Logic durch Domain Models',
                'Bessere Alignment zwischen Code und Business'
            ],
            'nachteile': [
                'Hoher Initial Investment für Domain Modeling',
                'Erfordert enge Zusammenarbeit mit Domain Experts',
                'Kann zu Over-Engineering bei einfachen Domains führen',
                'Steile Lernkurve für traditionelle CRUD Entwickler'
            ]
        }
    }
    
    # Merge manual trade-offs with extracted ones
    for pattern, tradeoff in manual_tradeoffs.items():
        if pattern not in tradeoffs:
            tradeoffs[pattern] = tradeoff
    
    # Patterns that should have trade-offs in Part 2
    patterns_with_tradeoffs = ['CQRS', 'Event Sourcing', 'Circuit Breaker', 'Saga Pattern',
                              'API Gateway', 'Service Mesh', 'Bulkhead Pattern', 'Domain-Driven Design']
    
    existing_slide_titles = []
    for slide in prs.slides:
        if slide.shapes.title:
            existing_slide_titles.append(slide.shapes.title.text)
    
    print(f"Existing slides: {len(existing_slide_titles)}")
    
    # Add missing trade-off slides
    for pattern in patterns_with_tradeoffs:
        tradeoff_title = f"{pattern} - Vorteile/Nachteile"
        tradeoff_title_alt = f"{pattern} Trade-offs"
        
        if tradeoff_title not in existing_slide_titles and tradeoff_title_alt not in existing_slide_titles:
            # Look for the pattern in extracted tradeoffs
            found = False
            for key, value in tradeoffs.items():
                if pattern.lower() in key.lower() or key.lower() in pattern.lower():
                    print(f"Adding trade-off slide for: {pattern}")
                    # Check if we have the structure we expect
                    if 'vorteile' in value and 'nachteile' in value:
                        add_tradeoff_slide(prs, pattern, value['vorteile'], value['nachteile'], value.get('production'))
                    found = True
                    break
            if not found:
                print(f"No trade-offs found for: {pattern}")
    
    # Add speaker notes to all slides  
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title:
            title = slide.shapes.title.text
            # Find matching notes
            for section, section_notes in notes.items():
                if any(keyword in title.lower() for keyword in section.lower().split()):
                    if isinstance(section_notes, list):
                        note_text = '\n\n'.join(section_notes)
                    else:
                        note_text = section_notes
                    add_speaker_notes(slide, note_text)
                    break
    
    # Save updated presentation
    prs.save(prs_path)
    print(f"Updated presentation saved: {prs_path}")

def main():
    """Main processing function"""
    print("Adding missing Vorteile/Nachteile and speaker notes to architectural patterns presentations...")
    
    # Process both presentations
    process_part1()
    process_part2()
    
    print("Processing completed!")

if __name__ == "__main__":
    main()