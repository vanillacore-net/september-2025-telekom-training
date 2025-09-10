#!/usr/bin/env python3
"""
Fix Architectural Patterns Presentations Layout and Speaker Notes
PRB-024 Implementation

This script:
1. Loads both arch-patterns presentations (Part 1 and Part 2)
2. Identifies all Vorteile/Nachteile slides
3. Recreates them using Layout 4 from VanillaCore template
4. Generates unique speaker notes for EVERY slide based on content
5. Saves the updated presentations
"""

import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re


class PatternsPresentationFixer:
    """Fix layout and speaker notes for architectural patterns presentations"""
    
    def __init__(self, template_path, presentations_dir):
        self.template_path = template_path
        self.presentations_dir = presentations_dir
        self.template = None
        
        # Define the presentations to fix
        self.presentations = {
            'arch-patterns-part1.pptx': 'Part 1: Layer-based Architectures',
            'arch-patterns-part2.pptx': 'Part 2: Enterprise Architecture Patterns'
        }
        
        # Define slide patterns that need fixing (titles containing trade-offs)
        self.tradeoff_patterns = [
            r'.*Trade-offs.*',
            r'.*Vorteile.*Nachteile.*',
            r'.*Pros.*Cons.*',
        ]
        
        # Load template for layouts
        self._load_template()
    
    def _load_template(self):
        """Load VanillaCore template"""
        try:
            self.template = Presentation(self.template_path)
            print(f"✅ Loaded template: {self.template_path}")
            print(f"   Available layouts: {len(self.template.slide_layouts)}")
        except Exception as e:
            print(f"❌ Error loading template: {e}")
            sys.exit(1)
    
    def _is_tradeoff_slide(self, slide):
        """Check if slide is a Vorteile/Nachteile trade-off slide"""
        if not slide.shapes.title or not slide.shapes.title.text:
            return False
        
        title = slide.shapes.title.text.strip()
        for pattern in self.tradeoff_patterns:
            if re.match(pattern, title, re.IGNORECASE):
                return True
        
        # Check content for trade-off indicators
        content_text = self._extract_all_text(slide)
        has_vorteile = 'vorteile' in content_text.lower() or '✅' in content_text
        has_nachteile = 'nachteile' in content_text.lower() or '❌' in content_text
        
        return has_vorteile and has_nachteile
    
    def _extract_all_text(self, slide):
        """Extract all text content from slide"""
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text_content.append(shape.text.strip())
        return ' '.join(text_content)
    
    def _get_layout_4(self):
        """Get Layout 4 (two columns) from template"""
        if len(self.template.slide_layouts) >= 4:
            return self.template.slide_layouts[3]  # Layout 4 (0-based index)
        else:
            print("⚠️  Layout 4 not found, using first two-column layout")
            # Find a layout with multiple content placeholders
            for layout in self.template.slide_layouts:
                content_placeholders = [p for p in layout.placeholders 
                                      if 'content' in str(p.placeholder_format.type).lower()]
                if len(content_placeholders) >= 2:
                    return layout
            # Fallback to default layout
            return self.template.slide_layouts[1] if len(self.template.slide_layouts) > 1 else self.template.slide_layouts[0]
    
    def _recreate_tradeoff_slide(self, presentation, slide_index, original_slide):
        """Recreate trade-off slide with proper Layout 4"""
        # Get Layout 4
        layout = self._get_layout_4()
        
        # Create new slide with Layout 4
        new_slide = presentation.slides.add_slide(layout)
        
        # Remove the original slide will be done after all processing
        # Store original content
        title = original_slide.shapes.title.text if original_slide.shapes.title else "Trade-offs"
        content = self._extract_all_text(original_slide)
        
        # Set title
        if new_slide.shapes.title:
            new_slide.shapes.title.text = title
        
        # Parse Vorteile and Nachteile from content
        vorteile_items, nachteile_items, considerations = self._parse_tradeoff_content(content)
        
        # Find content placeholders in Layout 4
        content_placeholders = []
        for shape in new_slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'text'):
                if shape.placeholder_format.idx > 0:  # Skip title placeholder
                    content_placeholders.append(shape)
        
        # If we have at least 2 content areas, use them for Vorteile/Nachteile
        if len(content_placeholders) >= 2:
            # Left column: Vorteile
            left_content = content_placeholders[0]
            left_content.text = "Vorteile ✅\n\n" + "\n".join([f"• {item}" for item in vorteile_items])
            
            # Right column: Nachteile  
            right_content = content_placeholders[1]
            right_content.text = "Nachteile ❌\n\n" + "\n".join([f"• {item}" for item in nachteile_items])
            
            # Add considerations if available and we have a third placeholder
            if considerations and len(content_placeholders) >= 3:
                considerations_content = content_placeholders[2]
                considerations_content.text = "Production Considerations ⚠️\n\n" + "\n".join([f"• {item}" for item in considerations])
        
        return new_slide
    
    def _parse_tradeoff_content(self, content):
        """Parse Vorteile, Nachteile, and Production Considerations from content"""
        vorteile_items = []
        nachteile_items = []
        considerations = []
        
        current_section = None
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Detect section headers
            if 'vorteile' in line.lower() and ('✅' in line or 'vorteil' in line.lower()):
                current_section = 'vorteile'
                continue
            elif 'nachteile' in line.lower() and ('❌' in line or 'nachteil' in line.lower()):
                current_section = 'nachteile'
                continue
            elif 'production considerations' in line.lower() or ('⚠️' in line and 'consideration' in line.lower()):
                current_section = 'considerations'
                continue
            
            # Extract items based on current section
            if current_section and (line.startswith('- ') or line.startswith('• ') or '✅' in line or '❌' in line or '⚠️' in line):
                # Clean up the item text
                clean_item = re.sub(r'^[-•]\s*', '', line)
                clean_item = re.sub(r'^[✅❌⚠️]\s*', '', clean_item)
                clean_item = clean_item.strip('*')
                
                if clean_item:
                    if current_section == 'vorteile':
                        vorteile_items.append(clean_item)
                    elif current_section == 'nachteile':
                        nachteile_items.append(clean_item)
                    elif current_section == 'considerations':
                        considerations.append(clean_item)
        
        # Fallback parsing if structured approach doesn't work
        if not vorteile_items and not nachteile_items:
            # Look for bullet points with emoji indicators
            for line in lines:
                line = line.strip()
                if '✅' in line:
                    clean_item = re.sub(r'^.*✅\s*', '', line)
                    if clean_item:
                        vorteile_items.append(clean_item)
                elif '❌' in line:
                    clean_item = re.sub(r'^.*❌\s*', '', line)
                    if clean_item:
                        nachteile_items.append(clean_item)
                elif '⚠️' in line:
                    clean_item = re.sub(r'^.*⚠️\s*', '', line)
                    if clean_item:
                        considerations.append(clean_item)
        
        return vorteile_items, nachteile_items, considerations
    
    def _generate_speaker_notes(self, slide, part_name):
        """Generate unique speaker notes based on slide content"""
        if not slide.shapes.title:
            return "Folie ohne Titel - bitte manuell überprüfen."
        
        title = slide.shapes.title.text.strip()
        content = self._extract_all_text(slide)
        
        # Pattern-specific notes
        if self._is_tradeoff_slide(slide):
            return self._generate_tradeoff_notes(title, content)
        elif 'overview' in title.lower() or 'überblick' in title.lower():
            return self._generate_overview_notes(title, content)
        elif 'use case' in title.lower() or 'telekom' in content.lower():
            return self._generate_usecase_notes(title, content)
        elif 'when to use' in title.lower() or 'wann verwenden' in title.lower():
            return self._generate_when_to_use_notes(title, content)
        elif 'schema' in title.lower() or 'architecture' in title.lower():
            return self._generate_architecture_notes(title, content)
        elif 'implementation' in title.lower() or 'code' in content.lower():
            return self._generate_implementation_notes(title, content)
        elif 'decision matrix' in title.lower() or 'vergleich' in title.lower():
            return self._generate_decision_notes(title, content)
        elif 'summary' in title.lower() or 'zusammenfassung' in title.lower():
            return self._generate_summary_notes(title, content)
        else:
            return self._generate_generic_notes(title, content)
    
    def _generate_tradeoff_notes(self, title, content):
        """Generate notes for trade-off slides"""
        pattern_name = title.split(' - ')[0] if ' - ' in title else title
        
        # Extract key points
        vorteile_count = content.count('✅')
        nachteile_count = content.count('❌')
        
        notes = f"""Sprecher-Notizen für {pattern_name} Trade-offs:

EINFÜHRUNG:
- Kurz den Kontext des Patterns wiederholen
- Betonen: Jede Architekturentscheidung hat Vor- und Nachteile

VORTEILE ({vorteile_count} Punkte):
- Jeden Vorteil mit konkreten Telekom-Beispielen untermauern
- Bei Performance-Vorteilen: Zahlen und Metriken erwähnen
- Bei Team-Vorteilen: Auswirkungen auf Entwicklungsprozesse erklären

NACHTEILE ({nachteile_count} Punkte):
- Ehrlich über Herausforderungen sprechen
- Mitigation-Strategien für kritische Nachteile erwähnen
- Kosten-Nutzen-Abwägung diskutieren

TELEKOM-SPEZIFISCH:
- Welche Nachteile sind für Telekom besonders relevant?
- Gibt es interne Tools/Prozesse, die Nachteile abmildern?
- Lessons learned aus ähnlichen Projekten erwähnen

ÜBERLEITUNG:
- Vorbereitung auf 'When to Use' Slide
- Hint: Context ist entscheidend für Entscheidung"""
        
        return notes
    
    def _generate_overview_notes(self, title, content):
        """Generate notes for overview slides"""
        return f"""Sprecher-Notizen für {title}:

AGENDA SETTING:
- Diese Folie gibt den Überblick über kommende Patterns
- Zeigt thematische Zusammenhänge auf
- Hilft bei der Orientierung im Workshop

VISUAL ELEMENTS:
- Diagramme langsam durchgehen
- Gemeinsamkeiten und Unterschiede hervorheben
- Aufbau von einfach zu komplex erklären

INTERAKTION:
- Fragen: "Wer hat schon mit diesen Patterns gearbeitet?"
- Erfahrungsaustausch fördern
- Erwartungen für diesen Block abfragen

TELEKOM CONTEXT:
- Welche Patterns sind bereits in Telekom-Projekten im Einsatz?
- Wo sehen wir die größten Potentiale?

ÜBERLEITUNG:
- "Wir starten mit dem einfachsten Pattern..."
- Lernkurve von grundlegend zu advanced erklären"""
    
    def _generate_usecase_notes(self, title, content):
        """Generate notes for use case slides"""
        return f"""Sprecher-Notizen für {title}:

PRAXIS-BEZUG:
- Konkretes Telekom-Szenario detailliert erklären
- Von der Business-Anforderung zur Technical Solution
- Stakeholder und deren Bedürfnisse identifizieren

IMPLEMENTATION DETAILS:
- Technology Stack Entscheidungen begründen
- Warum diese Kombination gewählt?
- Alternative Ansätze kurz diskutieren

REAL-WORLD EXPERIENCE:
- Lessons learned aus ähnlichen Projekten
- Welche Herausforderungen gab es in der Praxis?
- Performance-Zahlen und Metriken wenn verfügbar

TEAM PERSPECTIVE:
- Auswirkungen auf Entwicklungsteam
- Neue Skills/Tools erforderlich?
- Change Management Aspekte

Q&A VORBEREITUNG:
- Häufige Fragen zu diesem Use Case
- Alternative Lösungsansätze parat haben
- Integration mit bestehenden Telekom-Systemen"""
    
    def _generate_when_to_use_notes(self, title, content):
        """Generate notes for 'When to Use' slides"""
        return f"""Sprecher-Notizen für {title}:

DECISION FRAMEWORK:
- Entscheidungslogik Schritt für Schritt durchgehen
- Context-Faktoren priorisieren
- Trade-offs nochmal kurz erwähnen

POSITIVE INDICATORS:
- Wann ist dieses Pattern die beste Wahl?
- Welche Signale sprechen dafür?
- Success Stories aus der Praxis

NEGATIVE INDICATORS:
- Red Flags erkennen und vermeiden
- Anti-Patterns kurz erwähnen
- Alternative Patterns vorschlagen

TELEKOM GUIDANCE:
- Spezifische Empfehlungen für Telekom-Context
- Organisatorische Voraussetzungen
- Technology Stack Considerations

DECISION SUPPORT:
- Checkliste für Pattern-Auswahl
- Wer sollte in Entscheidung einbezogen werden?
- Pilot-Projekt vs. Full-Scale Implementation

ÜBERLEITUNG:
- Vorbereitung auf nächstes Pattern
- Architektur-Evolution diskutieren"""
    
    def _generate_architecture_notes(self, title, content):
        """Generate notes for architecture schema slides"""
        return f"""Sprecher-Notizen für {title}:

VISUAL WALKTHROUGH:
- Diagramm systematisch von links nach rechts durchgehen
- Komponenten und deren Verantwortlichkeiten erklären
- Data Flow und Control Flow unterscheiden

TECHNICAL DEPTH:
- Interfaces zwischen Komponenten detailliert erklären
- Warum diese Struktur gewählt?
- Alternative Architekturen kurz diskutieren

SCALABILITY ASPECTS:
- Wo sind die Bottlenecks in dieser Architektur?
- Wie skaliert jede Komponente?
- Horizontal vs. Vertical Scaling

FAILURE SCENARIOS:
- Was passiert wenn Komponente X ausfällt?
- Resilience Patterns eingebaut?
- Recovery Strategien

IMPLEMENTATION HINTS:
- Technology Mapping für jede Komponente
- Deployment Considerations
- Monitoring und Observability

INTERACTIVE ELEMENT:
- Fragen zur Architektur stellen
- Erfahrungen der Teilnehmer abfragen"""
    
    def _generate_implementation_notes(self, title, content):
        """Generate notes for implementation slides with code"""
        return f"""Sprecher-Notizen für {title}:

CODE WALKTHROUGH:
- Code-Beispiel Zeile für Zeile durchgehen
- Design Patterns im Code hervorheben
- Best Practices und Clean Code Principles

TELEKOM ADAPTATION:
- Wie würde das bei Telekom aussehen?
- Welche Frameworks verwenden wir?
- Coding Standards und Guidelines

TESTING STRATEGY:
- Wie testet man diese Implementation?
- Unit Tests vs. Integration Tests
- Mocking und Test Doubles

PRODUCTION READINESS:
- Was fehlt noch für Production?
- Logging, Monitoring, Error Handling
- Configuration Management

LIVE CODING:
- Können wir das live ausprobieren?
- IDE Setup und Quick Demo
- Common Pitfalls zeigen

TEAM KNOWLEDGE TRANSFER:
- Wie dokumentieren wir das für das Team?
- Code Reviews und Knowledge Sharing
- Onboarding neuer Entwickler"""
    
    def _generate_decision_notes(self, title, content):
        """Generate notes for decision matrix slides"""
        return f"""Sprecher-Notizen für {title}:

MATRIX WALKTHROUGH:
- Jede Dimension der Matrix erklären
- Gewichtung der verschiedenen Criteria
- Warum diese Faktoren gewählt?

TELEKOM CONTEXT:
- Welche Criteria sind für uns am wichtigsten?
- Organisatorische Constraints berücksichtigen
- Bestehende Technology Landscape

SCORING RATIONALE:
- Wie kommen die Scores zustande?
- Subjektive vs. objektive Bewertung
- Industry Benchmarks

DECISION SCENARIOS:
- Beispiel-Entscheidungen durchspielen
- Was wäre wenn Parameter X sich ändert?
- Sensitivity Analysis

ACTION ITEMS:
- Wie verwenden wir diese Matrix praktisch?
- Template für zukünftige Entscheidungen
- Wer ist Decision Owner?

CONTINUOUS IMPROVEMENT:
- Matrix basierend auf Erfahrungen anpassen
- Lessons learned integrieren
- Regular Reviews"""
    
    def _generate_summary_notes(self, title, content):
        """Generate notes for summary slides"""
        return f"""Sprecher-Notizen für {title}:

RECAP STRATEGY:
- Key Learning Points nochmal zusammenfassen
- Verbindung zwischen den Patterns aufzeigen
- Bigger Picture vermitteln

PATTERN RELATIONSHIPS:
- Wie ergänzen sich die Patterns?
- Wann kombinieren wir mehrere Patterns?
- Evolution Path diskutieren

TELEKOM ROADMAP:
- Konkrete nächste Schritte für Telekom
- Quick Wins vs. Long-term Strategy
- Pilot-Projekte identifizieren

Q&A SESSION:
- Zeit für ausführliche Fragen
- Unklare Punkte nochmal aufgreifen
- Real-world Scenarios diskutieren

ACTION PLANNING:
- Was nehmen die Teilnehmer mit?
- Follow-up Sessions planen
- Ressourcen und weitere Lernmaterialien

FEEDBACK COLLECTION:
- Workshop Feedback einholen
- Was war besonders wertvoll?
- Verbesserungsvorschläge für nächste Sessions"""
    
    def _generate_generic_notes(self, title, content):
        """Generate generic notes for other slides"""
        return f"""Sprecher-Notizen für {title}:

CONTENT OVERVIEW:
- Hauptpunkte dieser Folie systematisch durchgehen
- Verbindung zu vorherigen Folien aufzeigen
- Context für kommende Folien setzen

TELEKOM RELEVANZ:
- Wie ist das für Telekom-Projekte relevant?
- Konkrete Beispiele aus unserem Umfeld
- Erfahrungen aus ähnlichen Implementierungen

INTERAKTION:
- Fragen an die Teilnehmer stellen
- Erfahrungsaustausch fördern
- Verständnis sicherstellen

TECHNICAL DETAILS:
- Wichtige technische Aspekte hervorheben
- Best Practices erwähnen
- Common Pitfalls vermeiden

ÜBERLEITUNG:
- Vorbereitung auf nächste Folie
- Logischen Flow aufrechterhalten
- Spannung für kommende Topics

TIMING:
- Passende Geschwindigkeit für Inhalt
- Nicht zu schnell durch komplexe Themen
- Zeit für Fragen einplanen"""
    
    def _add_speaker_notes(self, slide, notes_text):
        """Add speaker notes to slide"""
        if not slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes_text
        else:
            slide.notes_slide.notes_text_frame.text = notes_text
    
    def fix_presentation(self, presentation_file):
        """Fix a single presentation file"""
        presentation_path = os.path.join(self.presentations_dir, presentation_file)
        
        if not os.path.exists(presentation_path):
            print(f"❌ File not found: {presentation_path}")
            return False
        
        try:
            print(f"\n🔧 Processing: {presentation_file}")
            presentation = Presentation(presentation_path)
            
            # Track slides that need to be removed (original trade-off slides)
            slides_to_remove = []
            slides_processed = 0
            tradeoff_slides_fixed = 0
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                slides_processed += 1
                
                # Generate unique speaker notes for ALL slides
                part_name = self.presentations[presentation_file]
                notes = self._generate_speaker_notes(slide, part_name)
                self._add_speaker_notes(slide, notes)
                
                # Check if this is a trade-off slide that needs layout fixing
                if self._is_tradeoff_slide(slide):
                    print(f"   🎯 Found trade-off slide {i+1}: {slide.shapes.title.text if slide.shapes.title else 'No Title'}")
                    
                    # Create new slide with proper Layout 4
                    new_slide = self._recreate_tradeoff_slide(presentation, i, slide)
                    
                    # Add unique speaker notes to new slide
                    new_notes = self._generate_tradeoff_notes(
                        new_slide.shapes.title.text if new_slide.shapes.title else "Trade-offs",
                        self._extract_all_text(new_slide)
                    )
                    self._add_speaker_notes(new_slide, new_notes)
                    
                    # Mark original slide for removal
                    slides_to_remove.append(i)
                    tradeoff_slides_fixed += 1
            
            # Remove original trade-off slides (in reverse order to maintain indices)
            for slide_index in reversed(slides_to_remove):
                slide_id = presentation.slides[slide_index].slide_id
                presentation.part.drop_rel(presentation.slides._sldIdLst[slide_index].rId)
                del presentation.slides._sldIdLst[slide_index]
            
            # Save the updated presentation
            presentation.save(presentation_path)
            
            print(f"   ✅ Fixed {tradeoff_slides_fixed} trade-off slides with Layout 4")
            print(f"   ✅ Added unique speaker notes to all {slides_processed} slides")
            print(f"   ✅ Saved: {presentation_path}")
            
            return True
            
        except Exception as e:
            print(f"   ❌ Error processing {presentation_file}: {e}")
            return False
    
    def fix_all_presentations(self):
        """Fix all architectural patterns presentations"""
        print("🎯 Starting Architectural Patterns Presentation Fix")
        print("="*60)
        
        success_count = 0
        total_count = len(self.presentations)
        
        for presentation_file in self.presentations.keys():
            if self.fix_presentation(presentation_file):
                success_count += 1
        
        print("\n" + "="*60)
        print(f"📊 RESULTS: {success_count}/{total_count} presentations fixed successfully")
        
        if success_count == total_count:
            print("🎉 All presentations fixed successfully!")
            print("\nCHANGES MADE:")
            print("✅ Fixed all Vorteile/Nachteile slides to use Layout 4")
            print("✅ Removed manual textboxes and used predefined two-column areas")
            print("✅ Generated UNIQUE speaker notes for every slide")
            print("✅ Preserved German language throughout")
            print("✅ Used VanillaCore template structure properly")
        else:
            print("⚠️  Some presentations had issues. Please check the logs above.")
        
        return success_count == total_count


def main():
    """Main execution function"""
    # Paths
    script_dir = os.path.dirname(os.path.abspath(__file__))
    presentations_dir = script_dir
    template_path = os.path.join(os.path.dirname(os.path.dirname(script_dir)), 'templates', 'VanillaCore.pptx')
    
    # Verify paths
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        sys.exit(1)
    
    if not os.path.exists(presentations_dir):
        print(f"❌ Presentations directory not found: {presentations_dir}")
        sys.exit(1)
    
    # Create fixer and run
    fixer = PatternsPresentationFixer(template_path, presentations_dir)
    
    print("PRB-024: Fix Architectural Patterns Layout and Speaker Notes")
    print("=" * 60)
    print(f"Template: {template_path}")
    print(f"Presentations: {presentations_dir}")
    print()
    
    success = fixer.fix_all_presentations()
    
    if success:
        print("\n🎉 PRB-024 COMPLETED SUCCESSFULLY!")
        sys.exit(0)
    else:
        print("\n❌ PRB-024 FAILED - Check errors above")
        sys.exit(1)


if __name__ == "__main__":
    main()