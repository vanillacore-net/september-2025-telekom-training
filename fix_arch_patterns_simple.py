#!/usr/bin/env python3
"""
Fix architectural patterns presentations with proper VanillaCore template usage.
Simplified approach focusing on critical layout fixes.
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def create_arch_patterns_presentations():
    """Create both architectural patterns presentations with proper layouts."""
    
    # Load template
    template_path = "templates/VanillaCore.pptx"
    
    # Create Part 1 presentation
    create_part1_presentation(template_path)
    
    # Create Part 2 presentation  
    create_part2_presentation(template_path)


def create_part1_presentation(template_path):
    """Create Part 1: Layered, Microservices, Event-Driven, Hexagonal Architecture."""
    prs = Presentation(template_path)
    
    # Title slide (Layout 0)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Architectural Patterns - Part 1"
    slide.placeholders[1].text = "Design Patterns Workshop - Telekom Training"
    
    # Overview section header (Layout 1 - title only)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Overview: Layer-Based Architectures"
    
    # Overview ASCII diagram (Layout 9 - Code block)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Architecture Types Comparison"
    overview_ascii = """Layered Architecture        N-Tier Architecture        Clean Architecture
┌─────────────────┐        ┌─────────────────┐        ┌─────────────────┐
│   Presentation  │        │   Client Tier   │        │   Frameworks    │
├─────────────────┤        ├─────────────────┤        ├─────────────────┤
│    Business     │        │ Application     │        │   Interface     │
├─────────────────┤  <-->  │      Tier       │  <-->  │   Adapters      │
│  Persistence    │        ├─────────────────┤        ├─────────────────┤
├─────────────────┤        │   Data Tier     │        │  Use Cases      │
│    Database     │        └─────────────────┘        ├─────────────────┤
└─────────────────┘                                   │   Entities      │
                                                       └─────────────────┘
Classic Layering           Enterprise N-Tier          Domain-Centric"""
    add_code_content(slide, overview_ascii)
    
    # 1. Layered Architecture section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Layered Architecture (Schichtarchitektur)"
    
    # Layered Architecture schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Layered Architecture - Schema"
    layered_ascii = """┌─────────────────────────────────────────────────┐
│                Presentation Layer               │ <- Web UI, REST APIs
├─────────────────────────────────────────────────┤
│                 Business Layer                  │ <- Domain Logic
├─────────────────────────────────────────────────┤
│                Persistence Layer                │ <- Data Access
├─────────────────────────────────────────────────┤
│                 Database Layer                  │ <- Data Storage
└─────────────────────────────────────────────────┘

Dependency Flow: Top -> Down (Higher layers depend on lower layers)
Data Flow: Bidirectional through defined interfaces"""
    add_code_content(slide, layered_ascii)
    
    # Layered Architecture Use Case (Layout 2)
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "Layered Architecture - Telekom Use Case"
    add_bullet_content(slide, [
        "**Network Configuration Management System**",
        "**Layer Implementation:**",
        "• Presentation: Web Dashboard für Network Engineers",
        "• Business: Configuration Validation, Change Management", 
        "• Persistence: Configuration Repository, Audit Logging",
        "• Database: PostgreSQL für Config Data, MongoDB für Audit Logs",
        "",
        "**Real-world Application:**",
        "• SNMP configuration management",
        "• Network device inventory",
        "• Change approval workflows",
        "• Configuration backup and restore"
    ])
    
    # Layered Architecture Trade-offs (Layout 4 - Two columns)
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = "Layered Architecture - Trade-offs"
    add_advantages_disadvantages(slide, 
        advantages=[
            "Klare Trennung der Verantwortlichkeiten",
            "Einfache Testbarkeit jeder Schicht", 
            "Bewährtes Pattern mit hoher Entwickler-Akzeptanz",
            "Gute Performance bei einfachen CRUD-Operationen"
        ],
        disadvantages=[
            "Monolithischer Charakter erschwert Skalierung",
            "Änderungen propagieren durch alle Schichten",
            "Datenbankschema-Changes beeinflussen alle Layer",
            "Schwierig für komplexe Domain Logic"
        ]
    )
    
    # 2. Microservices section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Microservices Architecture"
    
    # Microservices schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Microservices Architecture - Schema"
    microservices_ascii = """Service Landscape:
┌───────────────┐  ┌───────────────┐  ┌───────────────┐
│  Device Mgmt  │  │   Monitoring  │  │  Configuration │
│   Service     │  │    Service    │  │    Service    │
├───────────────┤  ├───────────────┤  ├───────────────┤
│  PostgreSQL   │  │  InfluxDB     │  │   MongoDB     │
└───────────────┘  └───────────────┘  └───────────────┘
         │                  │                  │
         └──────────────────┼──────────────────┘
                            │
                   ┌─────────────────┐
                   │   API Gateway   │
                   ├─────────────────┤
                   │ Load Balancer   │
                   ├─────────────────┤
                   │ Service Mesh    │
                   └─────────────────┘"""
    add_code_content(slide, microservices_ascii)
    
    # Microservices Trade-offs (Layout 4)
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = "Microservices - Trade-offs"
    add_advantages_disadvantages(slide,
        advantages=[
            "Unabhängige Skalierung pro Service",
            "Technology Stack Diversität möglich",
            "Fehler-Isolierung zwischen Services", 
            "Parallele Entwicklung durch verschiedene Teams"
        ],
        disadvantages=[
            "Hohe operationale Komplexität",
            "Distributed System Challenges (Latency, Partial Failures)",
            "Data Consistency zwischen Services schwierig",
            "Service Discovery und Load Balancing erforderlich"
        ]
    )
    
    # 3. Event-Driven section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. Event-Driven Architecture"
    
    # Event-Driven schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Event-Driven Architecture - Schema" 
    event_ascii = """Event Flow:
┌─────────────┐    Events     ┌─────────────┐    Events     ┌─────────────┐
│  Producers  │ ─────────────> │  Event Bus  │ ─────────────> │  Consumers  │
│             │               │             │               │             │
│ Device      │               │  Apache     │               │ Alerting    │
│ Sensors     │               │  Kafka      │               │ System      │
│             │               │             │               │             │
│ User        │               │ Topics:     │               │ Analytics   │
│ Actions     │               │ - alerts    │               │ Engine      │
│             │               │ - metrics   │               │             │
│ System      │               │ - configs   │               │ Dashboard   │
│ Events      │               │ - audit     │               │ Updates     │
└─────────────┘               └─────────────┘               └─────────────┘

    │                             │                             │
    └─── Async Pub/Sub ────────────┼─── Message Queues ────────┘
                                   │
                          ┌─────────────┐
                          │ Event Store │
                          │ (History)   │
                          └─────────────┘"""
    add_code_content(slide, event_ascii)
    
    # 4. Hexagonal Architecture section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. Hexagonal Architecture (Ports & Adapters)"
    
    # Hexagonal schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Hexagonal Architecture - Schema"
    hexagonal_ascii = """                    Hexagonal Architecture
                           (Ports & Adapters)

    ┌─────────────────┐                    ┌─────────────────┐
    │   Web UI        │                    │   Database      │
    │   Adapter       │                    │   Adapter       │
    └─────┬───────────┘                    └─────┬───────────┘
          │                                      │
    ┌─────┴───────────┐                    ┌─────┴───────────┐
    │   HTTP Port     │                    │ Persistence Port│
    └─────┬───────────┘                    └─────┬───────────┘
          │                                      │
          │     ┌─────────────────────────┐      │
          └─────│                         │──────┘
                │     Domain Core         │
                │                         │
    ┌───────────│   Business Logic        │──────────┐
    │           │                         │          │
    │           └─────────────────────────┘          │
    │                                                │
┌───┴───────────┐                            ┌─────┴───────────┐
│  API Port     │                            │   Event Port    │
└───┬───────────┘                            └─────┬───────────┘
    │                                              │
┌───┴───────────┐                            ┌─────┴───────────┐
│   REST API    │                            │   Message       │
│   Adapter     │                            │   Queue Adapter │
└───────────────┘                            └─────────────────┘"""
    add_code_content(slide, hexagonal_ascii)
    
    # Summary slide (Layout 2)
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "Summary - Part 1"
    add_bullet_content(slide, [
        "**Pattern Selection Guide:**",
        "",
        "**Start Simple (Layered):**",
        "• Prototypen und Proof-of-Concepts",
        "• Kleine Teams (<5 Entwickler)",
        "• Einfache CRUD-Operationen",
        "",
        "**Scale Smart (Microservices):**", 
        "• Große verteilte Teams",
        "• Independent Service Scaling",
        "• High Availability Requirements",
        "",
        "**Integrate Efficiently (Event-Driven):**",
        "• Real-time Processing",
        "• System Integration",
        "",
        "**Design for Change (Hexagonal):**",
        "• Complex Domain Logic",
        "• Testability Requirements"
    ])
    
    # Save Part 1
    prs.save("presentations/powerpoint/arch-patterns-part1.pptx")
    print("✅ Created Part 1: arch-patterns-part1.pptx")


def create_part2_presentation(template_path):
    """Create Part 2: CQRS, Event Sourcing, Circuit Breaker, Saga, API Gateway, Service Mesh, Bulkhead, DDD."""
    prs = Presentation(template_path)
    
    # Title slide (Layout 0)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Architectural Patterns - Part 2"
    slide.placeholders[1].text = "Enterprise Architecture Patterns für Skalierbarkeit und Integration"
    
    # 1. CQRS section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "CQRS (Command Query Responsibility Segregation)"
    
    # CQRS schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "CQRS Pattern"
    cqrs_ascii = """Command Side              Event Store           Query Side
┌──────────────────────┐  ┌─────────────────┐  ┌──────────────────────┐
│ Write Model          │  │   Events       │  │ Read Model         │
│                      │  │                 │  │                    │
│ - Normalized         │─>│ Event Bus      │─>│ - Denormalized     │
│ - PostgreSQL         │  │ Command/Query  │  │ - Elasticsearch    │
│ - ACID Transactions  │  │ Sync           │  │ - Read Optimized   │
└──────────────────────┘  └─────────────────┘  └──────────────────────┘"""
    add_code_content(slide, cqrs_ascii)
    
    # CQRS Trade-offs (Layout 4)
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = "CQRS - Trade-offs"
    add_advantages_disadvantages(slide,
        advantages=[
            "Performance: Optimized Data Models für Read/Write",
            "Skalierung: Read/Write Operations unabhängig skalierbar",
            "Complex Queries: Performance durch denormalized Read Models",
            "Event Sourcing Integration: Natürliche Kombination möglich"
        ],
        disadvantages=[
            "Eventual Consistency: zwischen Read/Write Models", 
            "Maintenance: Doppelte Data Models erhöhen Aufwand",
            "Synchronization: Event Sync zwischen Models komplex",
            "Over-Engineering: für simple CRUD Applications"
        ]
    )
    
    # 2. Event Sourcing section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Event Sourcing"
    
    # Event Sourcing schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Event Sourcing Pattern"
    es_ascii = """Commands            Events               Current State
┌───────────┐    ┌─────────────┐      ┌─────────────────┐
│ Create    │───>│ DeviceCreated│───> │                 │
│ Device    │    │   Event     │     │     Current     │
└───────────┘    └─────────────┘     │      State      │
                                     │                 │
┌───────────┐    ┌─────────────┐     │   (Projection   │
│ Configure │───>│ConfigChanged│───> │   from Events)  │
│ Device    │    │   Event     │     │                 │
└───────────┘    └─────────────┘     └─────────────────┘"""
    add_code_content(slide, es_ascii)
    
    # 3. Circuit Breaker section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Circuit Breaker Pattern"
    
    # Circuit Breaker schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Circuit Breaker State Machine"
    cb_ascii = """    CLOSED                    OPEN                    HALF-OPEN
┌─────────────────┐      ┌─────────────────┐      ┌─────────────────┐
│ Normal          │      │ Service Down    │      │ Testing         │
│ Operation       │─────>│ Fast Fail       │─────>│ Recovery        │
│                 │      │                 │      │                 │
│ Success Rate    │      │ Timeout Period  │      │ Limited Calls   │
│ > Threshold     │      │ Elapsed         │      │                 │
└─────────────────┘      └─────────────────┘      └─────────────────┘
         ▲                                                  │
         └──────────────────────────────────────────────────┘
                         Success Rate > Threshold"""
    add_code_content(slide, cb_ascii)
    
    # 4. Saga Pattern section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Saga Pattern"
    
    # Saga schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Saga Pattern Types"
    saga_ascii = """Choreography Saga:           Orchestration Saga:
                            
Service A ──> Service B       ┌─────────────────┐
    │             │          │ Saga Manager   │
    ▼             ▼          │                 │
Service C ──> Service D      │ 1. Service A    │
                             │ 2. Service B    │
Event-driven,                │ 3. Service C    │
Distributed control          │ 4. Rollback?    │
                             └─────────────────┘
                             
                             Centralized control,
                             Explicit workflow"""
    add_code_content(slide, saga_ascii)
    
    # 5. API Gateway section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "API Gateway Pattern"
    
    # API Gateway schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "API Gateway Architecture"
    gateway_ascii = """          Client Requests
               │
               ▼
    ┌─────────────────────┐
    │   API Gateway       │
    │                     │
    │ - Authentication    │
    │ - Authorization     │
    │ - Rate Limiting     │
    │ - Load Balancing    │
    │ - Request Routing   │
    │ - Response Caching  │
    │ - Monitoring        │
    └─────────────────────┘
               │
               ▼
    ┌─────────┬─────────┬─────────┐
    │Service A│Service B│Service C│
    └─────────┴─────────┴─────────┘"""
    add_code_content(slide, gateway_ascii)
    
    # 6. Service Mesh section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Service Mesh"
    
    # Service Mesh schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Service Mesh Architecture (Istio)"
    mesh_ascii = """┌──────────────────────────────────────────────────────┐
│                Control Plane                         │
│ ┌─────────────┐ ┌─────────────┐ ┌─────────────┐      │
│ │   Pilot     │ │   Citadel   │ │   Mixer     │      │
│ │ (Discovery) │ │ (Security)  │ │ (Telemetry) │      │
│ └─────────────┘ └─────────────┘ └─────────────┘      │
└──────────────────────────────────────────────────────┘
                        │
                        ▼ (Configuration)
┌─────────────┐     ┌─────────────┐     ┌─────────────┐
│ Service A   │     │ Service B   │     │ Service C   │
│             │     │             │     │             │
│ ┌─────────┐ │◄────┤ ┌─────────┐ │◄────┤ ┌─────────┐ │
│ │ Proxy   │ │     │ │ Proxy   │ │     │ │ Proxy   │ │
│ │(Envoy)  │ │     │ │(Envoy)  │ │     │ │(Envoy)  │ │
│ └─────────┘ │     │ └─────────┘ │     │ └─────────┘ │
└─────────────┘     └─────────────┘     └─────────────┘

Data Plane: Envoy Proxies handle all traffic
Control Plane: Configures and monitors proxies"""
    add_code_content(slide, mesh_ascii)
    
    # 7. Bulkhead Pattern section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Bulkhead Pattern"
    
    # Bulkhead schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "Bulkhead Pattern - Resource Isolation"
    bulkhead_ascii = """Thread Pools per Client Type:
┌─────────────────────────────────────────────────┐
│                Service                          │
│                                                 │
│ ┌─────────────┐ ┌─────────────┐ ┌─────────────┐ │
│ │Premium      │ │Standard     │ │Basic        │ │
│ │Clients      │ │Clients      │ │Clients      │ │
│ │             │ │             │ │             │ │
│ │Thread Pool  │ │Thread Pool  │ │Thread Pool  │ │
│ │Size: 50     │ │Size: 30     │ │Size: 20     │ │
│ └─────────────┘ └─────────────┘ └─────────────┘ │
└─────────────────────────────────────────────────┘

Benefit: Basic clients cannot starve Premium clients"""
    add_code_content(slide, bulkhead_ascii)
    
    # 8. DDD section header (Layout 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Domain-Driven Design (DDD)"
    
    # DDD schema (Layout 9)
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    slide.shapes.title.text = "DDD Bounded Contexts - Telekom"
    ddd_ascii = """┌─────────────────┐    ┌─────────────────┐    ┌─────────────────┐
│ Customer        │    │ Billing         │    │ Network         │
│ Management      │    │ Context         │    │ Operations      │
│                 │    │                 │    │                 │
│ - Customer      │◄──►│ - Invoice       │◄──►│ - Device        │
│ - Contract      │    │ - Payment       │    │ - Configuration │
│ - Profile       │    │ - Tariff        │    │ - Monitoring    │
└─────────────────┘    └─────────────────┘    └─────────────────┘

Relationships:
◄──► Anti-Corruption Layer (ACL)
═══► Shared Kernel
───► Customer-Supplier"""
    add_code_content(slide, ddd_ascii)
    
    # Summary slide (Layout 2)
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = "Zusammenfassung - Part 2"
    add_bullet_content(slide, [
        "**Pattern-Kombinationen nutzen:**",
        "• CQRS + Event Sourcing: Für Audit und Performance",
        "• Microservices + Service Mesh: Für Infrastructure Concerns",
        "• Circuit Breaker + Bulkhead: Für System Resilience",
        "• API Gateway + DDD: Für Business-oriented APIs",
        "",
        "**Evolution über Revolution:**",
        "• Start Simple: Layered Architecture für Prototypen",
        "• Scale Systematically: Microservices für große Teams", 
        "• Optimize Specifically: CQRS/Event Sourcing für Performance",
        "",
        "**Telekom-spezifische Guidance:**",
        "✅ Network Operations: Event-Driven + CQRS",
        "✅ Customer Management: DDD + Event Sourcing",
        "✅ Billing Systems: Saga Pattern",
        "✅ Legacy Integration: Strangler Fig Pattern"
    ])
    
    # Save Part 2
    prs.save("presentations/powerpoint/arch-patterns-part2.pptx")
    print("✅ Created Part 2: arch-patterns-part2.pptx")


def add_code_content(slide, code_text):
    """Add monospace code content to slide using Layout 9."""
    if slide.placeholders[1]:
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = code_text.strip()
        
        # Apply monospace formatting for ASCII diagrams
        for run in p.runs:
            run.font.name = 'Consolas'
            run.font.size = Pt(9)
        
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = 1.0


def add_bullet_content(slide, bullet_points):
    """Add bullet point content to slide."""
    if slide.placeholders[1]:
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        first = True
        for point in bullet_points:
            if point.strip():  # Skip empty lines
                p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
                p.text = point
                p.level = 1 if point.startswith('•') else 0
                first = False


def add_advantages_disadvantages(slide, advantages, disadvantages):
    """Add advantages/disadvantages using two-column layout (Layout 4)."""
    # Left column (advantages)
    if slide.placeholders[1]:
        left_placeholder = slide.placeholders[1]
        left_frame = left_placeholder.text_frame
        left_frame.clear()
        
        # Add header
        header_p = left_frame.paragraphs[0]
        header_p.text = "Vorteile ✅"
        header_p.font.bold = True
        header_p.font.size = Pt(16)
        
        # Add advantages
        for advantage in advantages:
            p = left_frame.add_paragraph()
            p.text = f"• {advantage}"
            p.space_after = Pt(6)
    
    # Right column (disadvantages)  
    if slide.placeholders[2]:
        right_placeholder = slide.placeholders[2]
        right_frame = right_placeholder.text_frame
        right_frame.clear()
        
        # Add header
        header_p = right_frame.paragraphs[0]
        header_p.text = "Nachteile ❌"
        header_p.font.bold = True
        header_p.font.size = Pt(16)
        
        # Add disadvantages
        for disadvantage in disadvantages:
            p = right_frame.add_paragraph()
            p.text = f"• {disadvantage}"
            p.space_after = Pt(6)


if __name__ == "__main__":
    print("🚀 Fixing architectural patterns presentations...")
    print("   ✅ Using VanillaCore.pptx template correctly")
    print("   ✅ ASCII diagrams in Layout 9 (monospace preservation)")
    print("   ✅ Section headers in Layout 1 (title only)")
    print("   ✅ Vorteile/Nachteile in Layout 4 (two columns)")
    print("   ✅ German content preserved throughout")
    print()
    
    create_arch_patterns_presentations()
    
    print()
    print("✅ CRITICAL FIXES APPLIED:")
    print("   📄 Part 1: presentations/powerpoint/arch-patterns-part1.pptx")
    print("   📄 Part 2: presentations/powerpoint/arch-patterns-part2.pptx")
    print("   🎯 Template: VanillaCore.pptx used correctly")
    print("   🎯 Layouts: Proper layout assignments for all content types")
    print("   🎯 Formatting: ASCII diagrams preserved with monospace fonts")